"""Gera apresentação PowerPoint v3.2.1 — narrativa pedagógica.

Estrutura em 4 partes:
1. Onde paramos (recap da última reunião)
2. O que avancei nesta semana (resposta aos 4 pedidos do orientador)
3. Como a tese está sendo construída (v3.2 com conceitos explicados)
4. Próximos passos

Princípios:
- Linguagem narrativa, não tabular dense
- Cada termo técnico explicado ANTES do uso
- Voz coletiva em 1ª pessoa do plural ("validamos", "combinamos")
- Analogias para conceitos complexos
- Máximo 5-6 bullets por slide
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- Paleta acadêmica sóbria ---
NAVY = RGBColor(0x1F, 0x2A, 0x4E)
GRAY_DK = RGBColor(0x3D, 0x42, 0x4E)
GRAY_MD = RGBColor(0x70, 0x76, 0x82)
GRAY_LT = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(2.5), Inches(1.5), Inches(10.5), Inches(2.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Equidade Racial em Classificação Facial"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Atualização da pesquisa — o que evoluiu nesta semana"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GRAY_DK

    meta = slide.shapes.add_textbox(Inches(2.5), Inches(4.2), Inches(10.5), Inches(2.5))
    mf = meta.text_frame
    mf.word_wrap = True

    rows = [
        ("Mestrando:", "Marcello Ozzetti"),
        ("Orientador:", "Prof. Marcos Quiles"),
        ("Programa:", "Mestrado em Ciência da Computação — Unifesp / ICT"),
        ("Reunião:", "Junho de 2026 (segunda reunião do mês)"),
    ]
    for i, (k, v) in enumerate(rows):
        p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        p.text = f"{k}  {v}"
        p.font.size = Pt(17)
        p.font.color.rgb = GRAY_DK


def add_section_divider(prs: Presentation, num: str, title: str, subtitle: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tx_num = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3), Inches(3))
    pn = tx_num.text_frame.paragraphs[0]
    pn.text = num
    pn.font.size = Pt(130)
    pn.font.bold = True
    pn.font.color.rgb = WHITE

    tx_t = slide.shapes.add_textbox(Inches(4.3), Inches(2.5), Inches(8.5), Inches(2))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = GRAY_LT


def add_content_slide(prs: Presentation, title: str, bullets: list, footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    tx_b = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.3), Inches(5.4))
    bf = tx_b.text_frame
    bf.word_wrap = True

    for i, b in enumerate(bullets):
        if isinstance(b, tuple):
            text, level = b
        else:
            text, level = b, 0

        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = text
        p.level = level
        if level == 0:
            p.font.size = Pt(20)
            p.font.color.rgb = GRAY_DK
            p.font.bold = False
        else:
            p.font.size = Pt(17)
            p.font.color.rgb = GRAY_MD
        p.space_after = Pt(8)

    if footer:
        tx_f = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.4))
        pf = tx_f.text_frame.paragraphs[0]
        pf.text = footer
        pf.font.size = Pt(15)
        pf.font.color.rgb = GRAY_MD
        pf.font.italic = True


def add_explainer_slide(prs: Presentation, title: str, simple_def: str, why_matters: str, concrete: str = "", footer: str = "") -> None:
    """Slide explicativo: 'O que é X?' — definição simples + por que importa + exemplo concreto."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    p = tx_t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    # Em palavras simples
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.7))
    box1.fill.solid()
    box1.fill.fore_color.rgb = GRAY_LT
    box1.line.fill.background()
    tf1 = box1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.3); tf1.margin_right = Inches(0.3); tf1.margin_top = Inches(0.15)
    p1 = tf1.paragraphs[0]
    p1.text = "Em palavras simples"
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p1b = tf1.add_paragraph()
    p1b.text = simple_def
    p1b.font.size = Pt(18)
    p1b.font.color.rgb = GRAY_DK

    # Por que importa
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(3.5), Inches(12.1), Inches(1.7))
    box2.fill.solid()
    box2.fill.fore_color.rgb = WHITE
    box2.line.color.rgb = NAVY
    box2.line.width = Pt(1)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.3); tf2.margin_right = Inches(0.3); tf2.margin_top = Inches(0.15)
    p2 = tf2.paragraphs[0]
    p2.text = "Por que importa para a tese"
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2b = tf2.add_paragraph()
    p2b.text = why_matters
    p2b.font.size = Pt(18)
    p2b.font.color.rgb = GRAY_DK

    # Exemplo concreto
    if concrete:
        box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.6))
        box3.fill.solid()
        box3.fill.fore_color.rgb = WHITE
        box3.line.color.rgb = GREEN
        box3.line.width = Pt(1)
        tf3 = box3.text_frame
        tf3.word_wrap = True
        tf3.margin_left = Inches(0.3); tf3.margin_right = Inches(0.3); tf3.margin_top = Inches(0.15)
        p3 = tf3.paragraphs[0]
        p3.text = "Exemplo concreto"
        p3.font.size = Pt(15)
        p3.font.bold = True
        p3.font.color.rgb = GREEN
        p3b = tf3.add_paragraph()
        p3b.text = concrete
        p3b.font.size = Pt(17)
        p3b.font.color.rgb = GRAY_DK

    if footer:
        tx_f = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
        pf = tx_f.text_frame.paragraphs[0]
        pf.text = footer
        pf.font.size = Pt(14)
        pf.font.color.rgb = GRAY_MD
        pf.font.italic = True


def add_thesis_slide(prs: Presentation, title: str, statement: str, footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    p = tx_t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.4))
    box.fill.solid()
    box.fill.fore_color.rgb = GRAY_LT
    box.line.color.rgb = NAVY
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5); tf.margin_top = Inches(0.4)

    p = tf.paragraphs[0]
    p.text = statement
    p.font.size = Pt(20)
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT

    if footer:
        tx_f = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
        pf = tx_f.text_frame.paragraphs[0]
        pf.text = footer
        pf.font.size = Pt(14)
        pf.font.color.rgb = GRAY_MD
        pf.font.italic = True


def add_table_slide(prs: Presentation, title: str, headers: list, rows: list, footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    p = tx_t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    nrows = len(rows) + 1
    ncols = len(headers)
    table = slide.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.5), Inches(12.5), Inches(5.4)).table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = WHITE
                r.font.size = Pt(16)

    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(v)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else GRAY_LT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(15)
                    r.font.color.rgb = GRAY_DK

    if footer:
        tx_f = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.4))
        pf = tx_f.text_frame.paragraphs[0]
        pf.text = footer
        pf.font.size = Pt(15)
        pf.font.color.rgb = GRAY_MD
        pf.font.italic = True


def build_presentation(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ==================== TÍTULO + SUMÁRIO ====================
    add_title_slide(prs)

    add_content_slide(
        prs,
        "Principais alinhamentos e decisões",
        [
            "Parte 1 — Onde paramos na última reunião",
            ("Recap da última reunião e orientações", 1),
            "Parte 2 — O que avancei nesta semana",
            ("Resposta direta aos 4 pedidos: ampliar venues, revisar método, validar SOTA, reformular tese", 1),
            "Parte 3 — Como a tese está sendo construída agora",
            ("Linguagem clara: cada conceito explicado antes de usar", 1),
            "Parte 4 — Próximos passos",
            ("O que preciso decidir nesta reunião", 1),
        ],
    )

    # ==================== PARTE 1: ONDE PARAMOS ====================
    add_section_divider(prs, "1", "Onde paramos na última reunião", "Recap da última reunião e orientações")

    add_table_slide(
        prs,
        "Como a tese estava na última reunião — visão resumida",
        ["Dimensão", "Estado naquele momento"],
        [
            ["Postura", "Diagnóstica — explicar por que o erro existe"],
            ["Tese central", "75% F1 é teto fenotípico irredutível: raças se sobrepõem em tom de pele, especialmente Latinx"],
            ["Tarefa", "Apenas race classification (sem face recognition)"],
            ["Critério de sucesso", "Decomposição: irredutível vs redutível"],
            ["Corpus", "23 fichas em 6 tracks (A: race clf; B: FR fairness; C: skin tone; D: mitigação; E: auditoria; F: fundamentação)"],
            ["Validado na reunião", "Evolução metodológica, SOTA (FaceScanPaliGemma 75.7%), linha de pesquisa proposta"],
        ],
        footer="Esta é a base sobre a qual estamos construindo a nova versão (prescritiva) da tese",
    )

    add_content_slide(
        prs,
        "Os 4 pedidos que combinamos",
        [
            "1. Revisar o método em profundidade de cada artigo selecionado",
            ("Queremos enxergar não só o resultado, mas o rigor metodológico", 1),
            "",
            "2. Ampliar a pesquisa para venues de ML, Redes Neurais e temas relacionados",
            ("Identificamos lacuna de papers fundadores destes campos", 1),
            "",
            "3. Double-check do SOTA (75.7% F1)",
            ("Garantir que não emergiu competidor recente", 1),
            "",
            "4. Reformular a tese de DIAGNÓSTICA para PRESCRITIVA",
            ("Treinar classifier de tom → usar para condicionar classifier de raça → estender a face recognition", 1),
        ],
    )

    # ==================== PARTE 2: O QUE AVANCEI ESTA SEMANA ====================
    add_section_divider(prs, "2", "O que avancei nesta semana", "Resposta direta aos 4 pedidos")

    add_content_slide(
        prs,
        "Resposta ao pedido 2: ampliar venues ML / Redes Neurais",
        [
            "Executei Rodada 5 de triagem com foco em ML e Redes Neurais.",
            "",
            "Adicionei 6 fichas novas, todas de venues top da área:",
            ("• Hardt, Price & Srebro (NeurIPS 2016) — origem das métricas formais de fairness", 1),
            ("• Perez et al. (AAAI 2018) — FiLM, mecanismo de condicionamento neural", 1),
            ("• Zemel et al. (ICML 2013, Test-of-Time Award 2023) — Fair Representation Learning", 1),
            ("• Madras et al. (ICML 2018) — LAFTR, prova teórica de fair transferência", 1),
            ("• Zhang, Lemoine & Mitchell (AIES 2018) — adversarial debiasing", 1),
            ("• Kleinberg, Mullainathan & Raghavan (ITCS 2017) — teorema da impossibilidade", 1),
            "",
            "Cada um destes é base de algum baseline ou mecanismo do pipeline.",
        ],
    )

    add_content_slide(
        prs,
        "Resposta ao pedido 1: revisar método em profundidade",
        [
            "Criei uma nova Seção 12 normativa em cada ficha: ‘Análise crítica do método’.",
            "",
            "Avalia cada paper em 5 dimensões:",
            ("(a) Rigor formal — matemática defensável? Pressupostos declarados?", 1),
            ("(b) Reprodutibilidade — hiperparâmetros, código, seeds?", 1),
            ("(c) Aplicabilidade ao nosso pipeline — onde funciona e onde falha?", 1),
            ("(d) Design choices — o que foi justificado vs assumido?", 1),
            ("(e) Conexão com o que aprendi na Rodada 5", 1),
            "",
            "Adicionada nas 10 fichas centrais (lista de leitura prioritária).",
            "Vamos expandir para as outras 19 após aprovação conjunta da nova versão da tese.",
        ],
    )

    add_content_slide(
        prs,
        "Resposta ao pedido 3: double-check do SOTA",
        [
            "Executei a Rodada 2.6 com janela expandida fevereiro–junho de 2026.",
            "",
            "Fonte: 12 papers que citam o FaceScanPaliGemma no Semantic Scholar.",
            "",
            "Resultado: NENHUM dos 12 papers usa FairFace race 7-class como métrica principal.",
            ("Citantes trabalham em tarefas diferentes: gender, age, emotion, detection, audit, etc.", 1),
            "",
            "Conclusão: FaceScanPaliGemma 75.7% F1 segue sendo o SOTA único em junho/2026.",
            "",
            "Achados adicionais: validação cruzada confirmou também o baseline ResNet-34 = 72%",
            ("Reportado por AlDahoul (2024/26) E por Lin (FairGRAPE 2022) independentemente", 1),
        ],
    )

    add_content_slide(
        prs,
        "Resposta ao pedido 4: reformular para PRESCRITIVA",
        [
            "Mudança central: deixa de ‘decompor o erro’ e passa a ‘construir um pipeline’ que melhora fairness em race classification E em face recognition, usando tom de pele como sinal auxiliar.",
            "",
            "Pipeline construtivo em 6 etapas:",
            ("1. Classificador de tom de pele (MST) — usar SkinToneNet pré-treinado", 1),
            ("2. Auditar FairFace e publicar a matriz pública MST × raça", 1),
            ("3. Race classifier com tom de pele como contexto (mecanismo FiLM)", 1),
            ("4. Comparar fairness COM vs SEM tom de pele, contra 6 baselines", 1),
            ("5. Aplicar pipeline a face recognition (RFW ou BFW)", 1),
            ("6. Medir melhora em Black/African e decomposição Latinx", 1),
        ],
    )

    add_content_slide(
        prs,
        "Triagem do corpus de pesquisa",
        [
            "Dos 57 papers avaliados, aprovamos 29 como fichas do corpus.",
            "",
            "Aprovação seguiu os critérios de seleção acordados:",
            ("Venue científico forte (top conferences ou journals peer-reviewed)", 1),
            ("Citações acima do threshold do ano de publicação", 1),
            ("Aplicabilidade direta a uma das 14 perguntas de pesquisa", 1),
            ("Aprovação por exceção (cobertura única) com justificativa registrada", 1),
            "",
            "Cada decisão (aprovado / standby / descartado) tem data e justificativa registradas.",
        ],
    )

    # ==================== PARTE 3: COMO A TESE ESTÁ SENDO CONSTRUÍDA ====================
    add_section_divider(prs, "3", "Como a tese está sendo construída", "Agora prescritiva, com conceitos explicados")

    add_thesis_slide(
        prs,
        "A tese em uma frase clara",
        "Treinar uma rede neural para reconhecer o tom de pele de uma foto, "
        "e usar essa informação como contexto extra ao treinar uma rede de "
        "classificação racial, melhora as métricas de fairness (a paridade "
        "entre raças). Essa melhoria também se transfere para tarefas "
        "downstream de reconhecimento facial, beneficiando especialmente "
        "grupos sub-representados (Black/African).\n\n"
        "Em outras palavras: o tom de pele é informação que o classificador "
        "de raça precisa ter explicitamente — não pode descobrir sozinho.",
    )

    # Conceitos explicados (3 slides) — ANTES do pipeline
    add_explainer_slide(
        prs,
        "O que é a escala MST (Monk Skin Tone)?",
        "Uma escala de 10 tons de pele (do tom 1 mais claro ao tom 10 mais escuro), "
        "criada em 2023 pelo sociólogo Dr. Ellis Monk (Harvard) em parceria com o Google. "
        "Foi desenhada especificamente para auditar fairness em sistemas de IA — não é "
        "a escala antiga de dermatologia (Fitzpatrick), que tem viés para pele clara.",
        "Nossa tese precisa medir tom de pele de forma confiável e granular. MST é o "
        "padrão moderno aceito pela comunidade de fairness research (publicado em "
        "NeurIPS 2023). Substituir a Fitzpatrick pela MST evita 3 problemas conhecidos "
        "da escala antiga.",
        "A Fitzpatrick (1975) tem só 6 tons, sendo 3 deles para variações de 'pele "
        "considerada branca'. Um terço dos próprios dermatologistas confunde Fitzpatrick "
        "com classificação racial — erro categorial endêmico.",
        footer="Classificador MST pré-treinado validado pela Rodada 6 (SkinToneNet, Pereira 2026)",
    )

    _add_film_math_slide(prs)

    _add_laftr_theory_slide(prs)

    # Pipeline com analogia
    add_content_slide(
        prs,
        "O pipeline em 6 etapas",
        [
            "Etapa 1 — Treinar um classificador de tom de pele (MST)",
            ("Dados: MST-E (Schumann 2023) + Casual Conversations (Meta 2021)", 1),
            "Etapa 2 — Auditar o FairFace: como tom de pele se distribui entre as 7 raças?",
            ("Aplicar o classificador MST sobre 10.954 imagens; validar manualmente um subset", 1),
            "Etapa 3 — Treinar o classificador de raça USANDO tom de pele como contexto (via FiLM)",
            ("Backbone: ConvNeXt-T (rede moderna, 28 milhões de parâmetros — leve)", 1),
            "Etapa 4 — Comparar fairness: COM tom de pele vs SEM",
            ("Baselines: ConvNeXt-T puro, FSCL+, Group DRO, Adversarial debiasing", 1),
            "Etapa 5 — Aplicar pipeline idêntico em reconhecimento facial (RFW ou BFW)",
            ("Verificar transferência da propriedade fair", 1),
            "Etapa 6 — Medir se Black/African melhora especificamente",
            ("Métrica: TAR @ FAR fixo por raça", 1),
        ],
        footer="Etapas 1, 4 e 5 reforçadas pela Rodada 6 (SkinToneNet, Adversarial debiasing, fair transfer empírico)",
    )

    _add_metrics_slide(prs)

    add_table_slide(
        prs,
        "As 5 hipóteses que vamos testar (na forma de perguntas)",
        ["#", "Pergunta", "Como saberemos se SIM"],
        [
            ["H1", "O pipeline com tom de pele melhora fairness em classificação?", "F1 sobe ≥+2pp E DR cai ≥20% vs ResNet-34"],
            ["H2", "Trocar a rede (ResNet-34 → ConvNeXt-T) por si só ajuda Latinx?", "Latinx permanece em ≈60% (não muda)"],
            ["H3", "Latinx tem tom de pele espalhado em várias categorias?", "Spread cobre ≥5 das 10 categorias MST"],
            ["H4", "A maior parte dos erros do Latinx é por sobreposição de tom?", "≥50% dos erros estão em zonas de overlap"],
            ["H5", "A melhoria transfere para reconhecimento facial?", "Black/African melhora ≥+3pp em RFW ou BFW"],
        ],
        footer="H5 em revisão após Pangelinan 2023 (Rodada 6): pixel info pode ser confounder — discussão na pauta",
    )

    # ==================== PARTE 4: PRÓXIMOS PASSOS ====================
    add_section_divider(prs, "4", "Próximos passos", "Plano experimental e calendário")

    add_table_slide(
        prs,
        "Cronograma estimado",
        ["Fase", "Target", "O que será entregue"],
        [
            ["Aprovação conjunta da nova versão da tese", "Junho 2026", "Ajustes na tese se necessário"],
            ["Submissão para qualificação", "Julho 2026", "Documento de qualificação ao programa"],
            ["Setup metodológico", "Agosto 2026", "Documentos com especificações detalhadas"],
            ["Capítulo 1 — Classificador MST + matriz", "Setembro 2026", "Resultado de H3"],
            ["Capítulo 2 — Race + condicionamento", "Dezembro 2026", "Resultados de H1, H2, H4"],
            ["Capítulo 3 — Face recognition", "Fevereiro 2027", "Resultado de H5"],
            ["Síntese final", "Março 2027", "Decomposição do erro Latinx"],
            ["Escrita dos capítulos", "Em paralelo (Set 2026 → Mar 2027)", "Texto final da dissertação"],
            ["Defesa prevista", "Abril 2027", "Defesa da dissertação"],
        ],
    )

    add_content_slide(
        prs,
        "Os 3 capítulos experimentais",
        [
            "Capítulo 1 — Construir o classificador de tom",
            ("Treinar sobre MST-E + Casual Conversations", 1),
            ("Aplicar no FairFace e construir a matriz pública tom × raça", 1),
            ("Testa H3 (Latinx tem spread amplo)", 1),
            "",
            "Capítulo 2 — Classificador de raça com tom como contexto",
            ("Pipeline ConvNeXt-T + FiLM, comparado com 4 baselines independentes", 1),
            ("3 seeds para significância estatística", 1),
            ("Testa H1 (pipeline funciona), H2 (Latinx é estrutural), H4 (overlap explica erros)", 1),
            "",
            "Capítulo 3 — Aplicar a reconhecimento facial",
            ("Mesmo pipeline em RFW ou BFW (datasets que já estão catalogados)", 1),
            ("Foco em accuracy de Black/African", 1),
            ("Testa H5 (fair transferência funciona)", 1),
        ],
    )

    add_content_slide(
        prs,
        "O que precisamos decidir nesta reunião",
        [
            "1. Aprovação da nova versão da tese (prescritiva, pipeline integrado)?",
            ("Se aprovarmos, seguimos para o detalhamento metodológico.", 1),
            ("Se houver ajustes, fazemos e revisamos juntos antes de prosseguir.", 1),
            "",
            "2. Alguma das 5 hipóteses precisa ser reformulada?",
            ("São o esqueleto da dissertação — se uma estiver mal formulada, melhor descobrir agora.", 1),
            "",
            "3. Escolha definitiva entre RFW e BFW para o Capítulo 3?",
            ("Temos preferência por RFW (mais escala, mais histórico) — precisamos validar essa escolha juntos.", 1),
            "",
            "4. Alguma nova frente de literatura que ficou faltando?",
            ("Podemos fazer Rodada 6 se necessário, mas o corpus já está consistente.", 1),
        ],
    )

    # ==================== SLIDE FINAL ====================
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Discussão e feedback"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Aprovação da nova versão da tese — decidir os 4 itens acima"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GRAY_LT
    p2.alignment = PP_ALIGN.CENTER

    # ==================== SLIDES DE APOIO ====================
    add_section_divider(prs, "★", "Slides de apoio", "Material complementar para consulta durante a discussão")

    _add_perguntas_slide(prs)

    _add_paper_card_slide(
        prs,
        "Pangelinan et al. (2023) — Causas de variação demográfica em FR accuracy",
        "arXiv:2304.07175 — Notre Dame / Florida Tech / IST — preprint abril 2023",
        contexto="Buolamwini & Gebru 2018 documentaram disparities em FR sem isolar mecanismo causal. Trabalhos subsequentes assumiram skin tone como variável dominante. Esta limpeza causal era um buraco aberto na literatura.",
        pergunta="A disparity racial em accuracy de face recognition é causada por skin tone direto, face size/shape, desbalanceamento de treino ou qualidade da imagem (pixel info útil após crop/alinhamento)?",
        metodo="Auditoria de ArcFace e AdaFace (SOTA em FR) sobre RFW e MORPH. Mede 4 variáveis: (1) skin tone via ITA, (2) face pixel area (pós-detecção), (3) balanceamento de treino, (4) face geometry. Correlaciona com TAR@FAR por grupo demográfico. Análise estratificada para isolar efeitos.",
        resultados="Face pixel information explica a MAIOR PARTE da variação de accuracy entre grupos. Skin tone direto é fator secundário. Balanceamento de treino tem efeito menor que o esperado. Sugere que muito do 'race bias' em FR é mediado por qualidade da imagem.",
        critica="Trabalha com FR (verificação 1:1), NÃO com race classification multi-classe. Análise correlacional, não causal estrito. Não testa intervenção arquitetural para mitigar pixel info. RFW e MORPH têm vieses próprios de coleta.",
        conexoes="Conversa com Grother 2019 (NIST FRVT) que documenta disparity industry-wide. Refuta implicitamente parte da narrativa de Buolamwini 2018 (PPB) que enfatiza skin tone direto. Compatível com Dooley 2022 que aponta arquitetura como fator independente.",
        impacto="REFUTAÇÃO POTENCIAL DE H5. Se pixel info é causa primária em FR, a transferência de fairness do Cap 2 para Cap 3 pode ser confounded por qualidade da imagem. Motivou as 3 versões de reformulação de H5 (V1/V2/V3) na pauta de hoje.",
        accent=ACCENT,
    )

    _add_paper_card_slide(
        prs,
        "Pereira et al. (2026) — SkinToneNet + dataset STW",
        "arXiv:2603.02475 — ICMC-USP / IMPA — preprint março 2026",
        contexto="Schumann et al. 2023 (NeurIPS) propuseram MST como padrão para auditoria mas sem classificador automático em larga escala. Datasets MST-anotados eram pequenos (MST-E ~1.500). Faltava um classificador SOTA generalizável e um dataset 'in-the-wild' grande.",
        pergunta="É possível treinar um classificador de tom de pele (MST 10-classe) robusto e cross-domain, e qual é a distribuição de tons em datasets faciais largamente usados pela comunidade de fairness?",
        metodo="Constrói STW (Skin Tone in the Wild): 42.313 imagens de 3.564 indivíduos, anotadas pela escala Monk Skin Tone (10 classes), agregadas de múltiplas fontes web. Treina SkinToneNet (ViT-Small fine-tuned a partir de ImageNet). Audita FairFace, CelebA, BUPT, UTKFace, CASIA, CK+ usando o classificador treinado.",
        resultados="SkinToneNet atinge SOTA em skin-tone classification cross-domain (precisão >90% out-of-distribution). FairFace exibe distribuição AGREGADA com forte ausência de MST 6-10 (subrepresentação severa de pele escura). NÃO publica matriz cruzada MST × race — apenas distribuição global.",
        critica="STW concentra sujeitos de imagens web (viés de quem é fotografado e publicado). Dataset não tem split público fixo para benchmarks reprodutíveis. Anotação MST sem protocolo multi-anotador como Schumann 2023. Apenas um classificador (ViT-Small) avaliado.",
        conexoes="Estende Schumann et al. 2023 (MST scale + protocolo) com classificador SOTA. Conversa com Dominguez-Catena et al. 2024 (DSAP) que auditou datasets via métricas unificadas. Insumo direto para nossa C2 (matriz cruzada MST × race do FairFace).",
        impacto="DECISÃO TÉCNICA: usar SkinToneNet pré-treinado como insumo do pipeline (economiza ~3 semanas de cronograma). Nossa C2 (matriz cruzada MST × raça) segue contribuição original. STW pode ser dataset de validação externa do Cap 1.",
        accent=GREEN,
    )

    _add_paper_card_slide(
        prs,
        "Dooley et al. (2022/23) — Fairer architectures make for fairer FR",
        "arXiv:2210.09943 — Maryland / Bosch / ELLIS / Tübingen — versão dez 2023",
        contexto="A literatura de fairness mitigation focava majoritariamente em dados (re-sampling) e loss functions (FSCL, Group DRO). Arquitetura era assumida como variável neutra. Faltava investigação sobre se NAS pode encontrar arquiteturas inerentemente mais fair.",
        pergunta="Bias em FR é apenas consequência dos dados de treino ou também emerge da arquitetura? Existe arquitetura que reduza disparity sem qualquer intervenção em dados ou loss?",
        metodo="Neural Architecture Search (NAS) bi-objetivo: accuracy E fairness (medida por DR / EOD). Busca jointly arquitetura + hiperparâmetros em espaço de centenas de candidatos. Datasets de avaliação: CelebA, RFW, BFW. Métricas: TAR@FAR por grupo + Demographic Parity Difference.",
        resultados="Frente de Pareto inclui arquiteturas que DOMINAM ResNet-50/100 e MobileNet em accuracy E fairness simultaneamente. Modificações arquiteturais SOZINHAS (sem dados extra, sem mitigação algorítmica) reduzem disparity em até 30%. Arquiteturas vencedoras têm padrões distintos (mais skip connections, menos depth).",
        critica="Custo computacional altíssimo — NAS proibitivo para a maioria dos labs. Arquiteturas vencedoras não têm interpretação semântica clara (por que essa estrutura é mais fair?). Não testa em FairFace 7-class race especificamente. ConvNeXt-T não está no espaço de busca.",
        conexoes="Conversa com Manzoor & Rattani 2024 (FineFACE) que também aposta em mudança arquitetural (cross-layer attention). Implicitamente reforça preocupação de Lin et al. 2022 (FairGRAPE) com modificações estruturais. Contrasta com Park et al. 2022 (FSCL) que muda apenas a loss.",
        impacto="REFORÇA H2 (testar se trocar de ResNet-34 para ConvNeXt-T já reduz disparity). Justifica escolha de backbone moderno. Se H1 falhar mas ConvNeXt-T puro já reduzir disparity, parte do efeito é atribuível à arquitetura — resultado científico válido.",
        accent=NAVY,
    )

    _add_paper_card_slide(
        prs,
        "Aguirre & Dredze (2023) — Transferring fairness via multi-task learning",
        "arXiv:2305.12671 — Johns Hopkins (NLP) — versão revisada abr 2024",
        contexto="LAFTR (Madras 2018) provou teoricamente que fairness pode transferir entre tarefas via representação compartilhada, mas a validação empírica era limitada. Quando demographic labels existem em uma tarefa mas não na outra, esse princípio teórico tinha utilidade prática não explorada.",
        pergunta="Quando demographic labels estão disponíveis APENAS em uma tarefa relacionada (não na tarefa-alvo), é possível transferir fairness via multi-task learning? Em que magnitude e sob quais condições?",
        metodo="Adapta single-task fairness loss para framework multi-task. Treina classificador na tarefa-alvo SEM demographic labels, usando demographic labels APENAS na tarefa auxiliar. Avalia em datasets de NLP (toxicidade Jigsaw, sentiment Twitter) com atributos sensíveis de raça e gênero. Compara contra single-task baseline e contra multi-task sem fairness loss auxiliar.",
        resultados="Objetivos de fairness demográfico se transferem empiricamente para a tarefa-alvo em todos os datasets testados. Permite intersectional fairness combinando datasets com atributos demográficos diferentes (single-axis). Magnitude da transferência: ~70-80% da redução de disparity vista em single-task fair.",
        critica="Experimentos em NLP — não validados em CV/face. Magnitude da transferência depende de correlação entre tarefas auxiliar e alvo. Não testa robustez sob distribution shift. Único atributo sensível por experimento.",
        conexoes="Reforço empírico direto de Madras et al. 2018 (LAFTR). Conversa com Zemel et al. 2013 (LFR — paradigma fundador). Compatível com Hardt et al. 2016 (métricas de fairness).",
        impacto="REFORÇO EMPÍRICO do princípio teórico do LAFTR. Suporta etapa 3 (race + tom auxiliar) e etapa 5 (transferência para FR) do pipeline. Princípio é independente da modalidade, mas precisamos validar empiricamente em face — é o que faremos no Cap 3.",
        accent=NAVY,
    )

    _add_paper_card_slide(
        prs,
        "Kolla & Savadamuthu (2022/23) — Impact of racial distribution in FR training",
        "arXiv:2211.14498 — WACVW 2023 (IEEE/CVF Winter Conf. Workshops)",
        contexto="A linha de mitigação 'rebalancear o dataset' (Karkkainen & Joo 2021 com FairFace) era assumida como suficiente. Mas se balanceamento bastasse, o gap Latinx no FairFace já estaria resolvido — e não está. Faltava investigação sistemática do limite dessa abordagem.",
        pergunta="Balancear a distribuição racial no dataset de treino é SUFICIENTE para eliminar disparities raciais em face recognition? Em que cenários falha?",
        metodo="Treina ArcFace em variantes do dataset MS1MV2 com proporções variáveis de raças: uniforme (33% cada), skewed-toward-Black (60% Black), skewed-toward-White (60% White). Introduz métrica 'racial gradation' para medir correlação intra/inter-classe. Avalia em RFW e auditoria estratificada por qualidade da imagem.",
        resultados="Distribuição uniforme de raças no treino NÃO garante FR sem viés. Skewed-toward-Black supera uniforme para grupo Black em ~3pp TAR@FAR. Qualidade da imagem introduz disparity adicional que NENHUMA distribuição de treino captura. Sugere limite fundamental da estratégia 'só balancear dados'.",
        critica="Workshop paper — revisão menos profunda que main conference. Foco em ArcFace específico (outras losses podem se comportar diferente). Não compara contra intervenções arquiteturais. Skewed-toward-Black testado apenas com 60% (não outras proporções).",
        conexoes="Refina Karkkainen & Joo 2021 (FairFace — premissa do balanceamento). Compatível com Pangelinan 2023 (pixel info como confounder). Reforça argumento de Dooley 2022 (arquitetura importa).",
        impacto="JUSTIFICA INTERVENÇÃO ARQUITETURAL além de balanceamento de dados. FairFace JÁ é balanceado, e o gap Latinx persiste — coerente com nossa tese de que precisamos do FiLM-conditioning. Reforça argumento contra resolver fairness apenas com dados.",
        accent=NAVY,
    )

    _add_paper_card_slide(
        prs,
        "Liu et al. (2025) — BNMR: Bayesian Network meta-learning para fairness",
        "arXiv:2505.01699 — ACM FAccT 2025 (top venue ética AI)",
        contexto="Mitigation methods existentes (FSCL, Group DRO, AdvDebias) requerem hiperparâmetros fixos a priori que controlam o trade-off accuracy×fairness. Esses hiperparâmetros são frágeis: ótimo varia por subgrupo, por época de treino e por configuração de dados. Faltava abordagem adaptativa.",
        pergunta="É possível mitigar bias em face attribute classification de forma adaptativa, sem hiperparâmetros de fairness fixos a priori, ajustando o trade-off em tempo de treino?",
        metodo="BNMR (Bayesian Network-informed Meta Reweighting): (1) rede bayesiana modela dependências entre atributos faciais (componentes); (2) meta-learning calibra reweighting de amostras durante treino; (3) tracking dinâmico de viés do modelo em cada época. Avalia em CelebA com atributos sensíveis raça e gênero. Compara contra FSCL+, Group DRO, AdvDebias.",
        resultados="Supera baselines em CelebA em DEMOGRAPHIC PARITY (+12% vs FSCL+) e EQUAL OPPORTUNITY (+8% vs Group DRO). Permite calibração adaptativa do trade-off accuracy×fairness sem grid search — economiza ordens de magnitude em compute de tuning.",
        critica="Custo computacional adicional da rede bayesiana (precisa ser treinada conjuntamente). Testado em CelebA, NÃO em FairFace race 7-class. Pressupõe que o usuário define corretamente a estrutura causal entre atributos — pressuposto forte. Sem ablação clara do componente bayesiano.",
        conexoes="Alternativa metodológica a Park et al. 2022 (FSCL+) e Sagawa et al. 2020 (Group DRO). Estende ideia de Zhang et al. 2018 (AdvDebias) com adaptividade. Usa métricas de Hardt 2016.",
        impacto="BASELINE COMPETITIVO recente do Cap 2. Mecanismo ortogonal ao nosso (sample reweighting vs FiLM-conditioning) — comparação justa contra SOTA em fairness 2025. Se BNMR superar nosso pipeline, é informação útil para a tese (não falha — calibra expectativa).",
        accent=NAVY,
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def _add_metrics_slide(prs: Presentation) -> None:
    """Slide das 3 métricas em layout HORIZONTAL (3 colunas lado a lado) + faixa final do teorema."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7))
    p = tx_t.text_frame.paragraphs[0]
    p.text = "As métricas que vamos reportar"
    p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(12.5), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY; line.line.fill.background()

    # Subtítulo
    tx_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.5), Inches(0.4))
    psub = tx_sub.text_frame.paragraphs[0]
    psub.text = "Todas públicas, padrão da literatura — reportar as 3 simultaneamente é exigência metodológica, não escolha."
    psub.font.size = Pt(14); psub.font.color.rgb = GRAY_DK; psub.font.italic = True

    # 3 colunas lado a lado
    metricas = [
        {
            "nome": "F1 macro",
            "subtitulo": "performance média",
            "definicao": "Média harmônica de precision/recall por classe, depois média simples entre as 7 raças.",
            "ref": "van Rijsbergen 1979",
            "valor_titulo": "SOTA atual",
            "valor": "75%",
            "valor_sub": "FaceScanPaliGemma",
            "color": NAVY,
        },
        {
            "nome": "DR",
            "subtitulo": "Disparity Ratio",
            "definicao": "Razão entre o F1 da pior raça e o F1 da melhor. Mede o gap entre subgrupos.",
            "ref": "Hardt, Price & Srebro 2016 (NeurIPS)",
            "valor_titulo": "Estado atual",
            "valor": "0.67",
            "valor_sub": "60% Latinx ÷ 90% Black",
            "color": ACCENT,
        },
        {
            "nome": "Worst-class F1",
            "subtitulo": "pior subgrupo",
            "definicao": "F1 da raça em que o modelo erra mais. Garante que ninguém fique para trás.",
            "ref": "Sagawa et al. 2020 (ICLR) — Group DRO",
            "valor_titulo": "Estado atual",
            "valor": "60%",
            "valor_sub": "Latinx",
            "color": GREEN,
        },
    ]

    col_w = 4.0
    gap = 0.25
    total_w = 3 * col_w + 2 * gap
    x0 = (13.33 - total_w) / 2
    y_top = 1.6
    h_card = 4.5

    for i, m in enumerate(metricas):
        x = x0 + i * (col_w + gap)
        # Card de fundo
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y_top), Inches(col_w), Inches(h_card))
        card.fill.solid(); card.fill.fore_color.rgb = GRAY_LT; card.line.fill.background()

        # Cabeçalho colorido
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y_top), Inches(col_w), Inches(0.9))
        header.fill.solid(); header.fill.fore_color.rgb = m["color"]; header.line.fill.background()
        tfh = header.text_frame; tfh.word_wrap = True
        tfh.margin_left = Inches(0.2); tfh.margin_right = Inches(0.2); tfh.margin_top = Inches(0.1)
        ph = tfh.paragraphs[0]; ph.text = m["nome"]
        ph.font.size = Pt(20); ph.font.bold = True; ph.font.color.rgb = WHITE
        ph.alignment = PP_ALIGN.CENTER
        phb = tfh.add_paragraph(); phb.text = m["subtitulo"]
        phb.font.size = Pt(12); phb.font.color.rgb = GRAY_LT; phb.font.italic = True
        phb.alignment = PP_ALIGN.CENTER

        # Conteúdo abaixo do cabeçalho
        tx_def = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y_top + 1.0), Inches(col_w - 0.3), Inches(1.3))
        tf_def = tx_def.text_frame; tf_def.word_wrap = True
        pdef = tf_def.paragraphs[0]; pdef.text = m["definicao"]
        pdef.font.size = Pt(13); pdef.font.color.rgb = GRAY_DK

        # Referência
        tx_ref = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y_top + 2.35), Inches(col_w - 0.3), Inches(0.4))
        tf_ref = tx_ref.text_frame; tf_ref.word_wrap = True
        pref = tf_ref.paragraphs[0]; pref.text = "Referência: " + m["ref"]
        pref.font.size = Pt(11); pref.font.color.rgb = GRAY_MD; pref.font.italic = True

        # Linha separadora
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.3), Inches(y_top + 2.85), Inches(col_w - 0.6), Inches(0.03))
        sep.fill.solid(); sep.fill.fore_color.rgb = GRAY_MD; sep.line.fill.background()

        # Valor (destaque)
        tx_vt = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y_top + 2.95), Inches(col_w - 0.3), Inches(0.35))
        pvt = tx_vt.text_frame.paragraphs[0]
        pvt.text = m["valor_titulo"]
        pvt.font.size = Pt(11); pvt.font.color.rgb = GRAY_MD; pvt.font.italic = True
        pvt.alignment = PP_ALIGN.CENTER

        tx_v = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y_top + 3.25), Inches(col_w - 0.3), Inches(0.7))
        pv = tx_v.text_frame.paragraphs[0]
        pv.text = m["valor"]
        pv.font.size = Pt(34); pv.font.bold = True; pv.font.color.rgb = m["color"]
        pv.alignment = PP_ALIGN.CENTER

        tx_vs = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y_top + 3.95), Inches(col_w - 0.3), Inches(0.35))
        pvs = tx_vs.text_frame.paragraphs[0]
        pvs.text = m["valor_sub"]
        pvs.font.size = Pt(11); pvs.font.color.rgb = GRAY_DK
        pvs.alignment = PP_ALIGN.CENTER

    # Faixa final — teorema da impossibilidade
    box_imp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(12.5), Inches(0.95))
    box_imp.fill.solid(); box_imp.fill.fore_color.rgb = NAVY; box_imp.line.fill.background()
    tfi = box_imp.text_frame; tfi.word_wrap = True
    tfi.margin_left = Inches(0.3); tfi.margin_right = Inches(0.3); tfi.margin_top = Inches(0.1)
    pi0 = tfi.paragraphs[0]
    pi0.text = "Por que 3 métricas e não uma?"
    pi0.font.size = Pt(13); pi0.font.bold = True; pi0.font.color.rgb = WHITE
    pi1 = tfi.add_paragraph()
    pi1.text = "Teorema da impossibilidade (Kleinberg, Mullainathan & Raghavan 2017, ITCS): não existe métrica única de fairness que satisfaça simultaneamente calibração, equal FPR e equal FNR — reportar as 3 simultaneamente é a forma honesta de comunicar trade-offs."
    pi1.font.size = Pt(12); pi1.font.color.rgb = GRAY_LT


def _add_film_math_slide(prs: Presentation) -> None:
    """Slide do FiLM com diagrama + detalhamento científico (analogia vai para speaker notes)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7))
    p = tx_t.text_frame.paragraphs[0]
    p.text = "O que é FiLM (Feature-wise Linear Modulation)?"
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(12.5), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY; line.line.fill.background()

    # Subtítulo
    tx_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.5), Inches(0.45))
    psub = tx_sub.text_frame.paragraphs[0]
    psub.text = "Mecanismo de condicionamento neural que permite a uma rede 'consultar' um sinal de contexto antes de decidir."
    psub.font.size = Pt(14); psub.font.color.rgb = GRAY_DK; psub.font.italic = True

    # Diagrama em 3 blocos
    y_box = 1.7
    box_h = 1.2
    box_w = 3.7
    gap = 0.4
    x_start = 0.5

    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_start), Inches(y_box), Inches(box_w), Inches(box_h))
    b1.fill.solid(); b1.fill.fore_color.rgb = GRAY_LT; b1.line.color.rgb = NAVY; b1.line.width = Pt(1.5)
    tfb1 = b1.text_frame; tfb1.word_wrap = True
    tfb1.margin_left = Inches(0.15); tfb1.margin_right = Inches(0.15); tfb1.margin_top = Inches(0.1)
    pb1 = tfb1.paragraphs[0]; pb1.text = "1.  Reconhecer tom de pele"
    pb1.font.size = Pt(13); pb1.font.bold = True; pb1.font.color.rgb = NAVY
    pb1b = tfb1.add_paragraph()
    pb1b.text = "SkinToneNet recebe a foto e devolve um vetor de contexto representando o tom (MST 1 a 10)."
    pb1b.font.size = Pt(11); pb1b.font.color.rgb = GRAY_DK

    a1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x_start + box_w + 0.02), Inches(y_box + box_h/2 - 0.18), Inches(gap - 0.04), Inches(0.36))
    a1.fill.solid(); a1.fill.fore_color.rgb = NAVY; a1.line.fill.background()

    x2 = x_start + box_w + gap
    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x2), Inches(y_box), Inches(box_w), Inches(box_h))
    b2.fill.solid(); b2.fill.fore_color.rgb = NAVY; b2.line.fill.background()
    tfb2 = b2.text_frame; tfb2.word_wrap = True
    tfb2.margin_left = Inches(0.15); tfb2.margin_right = Inches(0.15); tfb2.margin_top = Inches(0.1)
    pb2 = tfb2.paragraphs[0]; pb2.text = "2.  Camada FiLM"
    pb2.font.size = Pt(13); pb2.font.bold = True; pb2.font.color.rgb = WHITE
    pb2b = tfb2.add_paragraph()
    pb2b.text = "Usa o vetor de contexto para modular as features intermediárias da rede de raça (escala e deslocamento por canal)."
    pb2b.font.size = Pt(11); pb2b.font.color.rgb = GRAY_LT

    a2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x2 + box_w + 0.02), Inches(y_box + box_h/2 - 0.18), Inches(gap - 0.04), Inches(0.36))
    a2.fill.solid(); a2.fill.fore_color.rgb = NAVY; a2.line.fill.background()

    x3 = x2 + box_w + gap
    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x3), Inches(y_box), Inches(box_w), Inches(box_h))
    b3.fill.solid(); b3.fill.fore_color.rgb = GRAY_LT; b3.line.color.rgb = NAVY; b3.line.width = Pt(1.5)
    tfb3 = b3.text_frame; tfb3.word_wrap = True
    tfb3.margin_left = Inches(0.15); tfb3.margin_right = Inches(0.15); tfb3.margin_top = Inches(0.1)
    pb3 = tfb3.paragraphs[0]; pb3.text = "3.  Decidir a raça"
    pb3.font.size = Pt(13); pb3.font.bold = True; pb3.font.color.rgb = NAVY
    pb3b = tfb3.add_paragraph()
    pb3b.text = "ConvNeXt-T decide entre as 7 raças do FairFace já com o tom como contexto."
    pb3b.font.size = Pt(11); pb3b.font.color.rgb = GRAY_DK

    # Detalhamento científico (substitui analogia)
    box_d = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.05), Inches(12.5), Inches(2.6))
    box_d.fill.solid(); box_d.fill.fore_color.rgb = GRAY_LT; box_d.line.fill.background()
    tfd = box_d.text_frame; tfd.word_wrap = True
    tfd.margin_left = Inches(0.25); tfd.margin_right = Inches(0.25); tfd.margin_top = Inches(0.12)
    pd = tfd.paragraphs[0]; pd.text = "Detalhamento científico"
    pd.font.size = Pt(13); pd.font.bold = True; pd.font.color.rgb = NAVY
    detalhes_film = [
        "Origem (Perez et al. 2018, AAAI): proposto originalmente para Visual Question Answering, generaliza Conditional Batch Normalization (de Vries et al. 2017, NeurIPS). Família mais ampla: conditional normalization.",
        "Mecanismo central: para cada canal de uma feature map, FiLM aplica um par (escala, deslocamento) aprendido a partir do contexto. É linear nas features e não linear no contexto.",
        "Onde inserir no nosso pipeline: após cada bloco residual do ConvNeXt-T (4 blocos hierárquicos). Apenas as MLPs de modulação são novas; o backbone segue treinável end-to-end.",
        "Por que FiLM e não alternativas: concatenar o vetor MST nas features perde escala posicional; cross-attention seria mais caro e menos estável para vetor de 10 dimensões; conditional BN é menos expressivo (afeta apenas batchnorm).",
        "Evidências de uso: ~1.000+ papers citantes em CV, NLP e RL desde 2018 (Google Scholar). Estabelecido como técnica padrão de conditioning em redes neurais profundas.",
    ]
    for ln in detalhes_film:
        pln = tfd.add_paragraph(); pln.text = "•  " + ln
        pln.font.size = Pt(12); pln.font.color.rgb = GRAY_DK
        pln.space_after = Pt(2)

    # Por que importa para a tese
    box_i = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.75), Inches(12.5), Inches(1.4))
    box_i.fill.solid(); box_i.fill.fore_color.rgb = WHITE
    box_i.line.color.rgb = GREEN; box_i.line.width = Pt(1.5)
    tfi = box_i.text_frame; tfi.word_wrap = True
    tfi.margin_left = Inches(0.25); tfi.margin_right = Inches(0.25); tfi.margin_top = Inches(0.12)
    pi = tfi.paragraphs[0]; pi.text = "Por que importa para a tese"
    pi.font.size = Pt(13); pi.font.bold = True; pi.font.color.rgb = GREEN
    pib = tfi.add_paragraph()
    pib.text = "FiLM transforma a ideia 'usar tom de pele como contexto' em mecanismo formal, reproduzível e auditável. O ganho marginal sobre baseline puro é a evidência direta de H1. Custo de parâmetros adicional é mínimo (apenas duas MLPs pequenas por bloco)."
    pib.font.size = Pt(12); pib.font.color.rgb = GRAY_DK

    # Speaker notes — analogia
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "ANALOGIA PARA EXPLICAR FiLM ORALMENTE:\n\n"
        "Imagine um médico fazendo um diagnóstico. Antes de decidir, ele consulta exames "
        "complementares (ECG, raio-X) que dão contexto sobre o paciente. O diagnóstico final "
        "depende tanto do que ele observa no exame físico quanto do contexto desses exames.\n\n"
        "FiLM faz exatamente isso para a rede de raça: antes de decidir a classe, ela "
        "'consulta' a saída do classificador de tom de pele e ajusta sua decisão com base "
        "nesse contexto.\n\n"
        "USE A ANALOGIA SE A AUDIÊNCIA PEDIR UMA INTUIÇÃO MENOS TÉCNICA — caso contrário, "
        "fique no detalhamento científico do slide."
    )


def _add_laftr_theory_slide(prs: Presentation) -> None:
    """Slide do LAFTR com diagrama + detalhamento científico (analogia vai para speaker notes)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7))
    p = tx_t.text_frame.paragraphs[0]
    p.text = "O que é fair transferência (LAFTR)?"
    p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.0), Inches(12.5), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY; line.line.fill.background()

    # Subtítulo
    tx_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.5), Inches(0.45))
    psub = tx_sub.text_frame.paragraphs[0]
    psub.text = "Framework de adversarial fair representation learning que demonstra fairness como propriedade transferível entre tarefas."
    psub.font.size = Pt(14); psub.font.color.rgb = GRAY_DK; psub.font.italic = True

    # Diagrama compacto
    y_box = 1.7
    box_h = 1.2

    b1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_box), Inches(4.0), Inches(box_h))
    b1.fill.solid(); b1.fill.fore_color.rgb = GRAY_LT; b1.line.color.rgb = NAVY; b1.line.width = Pt(1.5)
    tfb1 = b1.text_frame; tfb1.word_wrap = True
    tfb1.margin_left = Inches(0.15); tfb1.margin_right = Inches(0.15); tfb1.margin_top = Inches(0.1)
    pb1 = tfb1.paragraphs[0]; pb1.text = "Cap 2 — Race classification"
    pb1.font.size = Pt(13); pb1.font.bold = True; pb1.font.color.rgb = NAVY
    pb1b = tfb1.add_paragraph()
    pb1b.text = "Treino do backbone com objetivo fair sobre raça. A representação aprendida (Z) é, por construção, equilibrada entre grupos."
    pb1b.font.size = Pt(11); pb1b.font.color.rgb = GRAY_DK

    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.6), Inches(y_box + box_h/2 - 0.18), Inches(0.7), Inches(0.36))
    arr.fill.solid(); arr.fill.fore_color.rgb = NAVY; arr.line.fill.background()

    b2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(y_box), Inches(2.5), Inches(box_h))
    b2.fill.solid(); b2.fill.fore_color.rgb = NAVY; b2.line.fill.background()
    tfb2 = b2.text_frame; tfb2.word_wrap = True
    tfb2.margin_left = Inches(0.15); tfb2.margin_right = Inches(0.15); tfb2.margin_top = Inches(0.2)
    pb2 = tfb2.paragraphs[0]; pb2.text = "Representação Z fair"
    pb2.font.size = Pt(13); pb2.font.bold = True; pb2.font.color.rgb = WHITE
    pb2.alignment = PP_ALIGN.CENTER
    pb2b = tfb2.add_paragraph()
    pb2b.text = "A propriedade fair é embutida AQUI."
    pb2b.font.size = Pt(11); pb2b.font.color.rgb = GRAY_LT
    pb2b.alignment = PP_ALIGN.CENTER

    arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.0), Inches(y_box + box_h/2 - 0.18), Inches(0.7), Inches(0.36))
    arr2.fill.solid(); arr2.fill.fore_color.rgb = NAVY; arr2.line.fill.background()

    b3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(y_box), Inches(4.2), Inches(box_h))
    b3.fill.solid(); b3.fill.fore_color.rgb = GRAY_LT; b3.line.color.rgb = NAVY; b3.line.width = Pt(1.5)
    tfb3 = b3.text_frame; tfb3.word_wrap = True
    tfb3.margin_left = Inches(0.15); tfb3.margin_right = Inches(0.15); tfb3.margin_top = Inches(0.1)
    pb3 = tfb3.paragraphs[0]; pb3.text = "Cap 3 — Face recognition"
    pb3.font.size = Pt(13); pb3.font.bold = True; pb3.font.color.rgb = NAVY
    pb3b = tfb3.add_paragraph()
    pb3b.text = "Reuso do backbone como ponto de partida. A nova tarefa herda a propriedade fair sem re-treinar do zero."
    pb3b.font.size = Pt(11); pb3b.font.color.rgb = GRAY_DK

    # Detalhamento científico (substitui analogia)
    box_d = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.05), Inches(12.5), Inches(2.6))
    box_d.fill.solid(); box_d.fill.fore_color.rgb = GRAY_LT; box_d.line.fill.background()
    tfd = box_d.text_frame; tfd.word_wrap = True
    tfd.margin_left = Inches(0.25); tfd.margin_right = Inches(0.25); tfd.margin_top = Inches(0.12)
    pd = tfd.paragraphs[0]; pd.text = "Detalhamento científico"
    pd.font.size = Pt(13); pd.font.bold = True; pd.font.color.rgb = NAVY
    detalhes_laftr = [
        "Linhagem teórica: parte da família Fair Representation Learning iniciada por Zemel et al. 2013 (ICML, Test-of-Time Award 2023). LAFTR avança o paradigma com formulação adversarial.",
        "Arquitetura: 3 componentes treinados conjuntamente — encoder (gera a representação Z), classificador da tarefa principal, adversário que tenta recuperar o atributo sensível a partir de Z.",
        "Resultado teórico chave (Teorema 1): se o adversário falha em recuperar o atributo sensível, então qualquer classificador downstream treinado sobre Z herda um limite superior de violação de fairness.",
        "Métricas implementadas: demographic parity, equal opportunity, equalized odds (Hardt et al. 2016). Compatível com a triangulação de métricas que vamos usar no Cap 2.",
        "Validações empíricas: UCI Adult, Heritage Health, CelebA. Aguirre & Dredze 2023 estendeu o princípio para multi-task NLP — evidência de generalização do mecanismo.",
        "Implementação: código aberto em github.com/VectorInstitute/laftr. Reduz custo de adaptação para o nosso pipeline.",
    ]
    for ln in detalhes_laftr:
        pln = tfd.add_paragraph(); pln.text = "•  " + ln
        pln.font.size = Pt(11); pln.font.color.rgb = GRAY_DK
        pln.space_after = Pt(2)

    # Por que importa para a tese
    box_i = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.75), Inches(12.5), Inches(1.25))
    box_i.fill.solid(); box_i.fill.fore_color.rgb = WHITE
    box_i.line.color.rgb = GREEN; box_i.line.width = Pt(1.5)
    tfi = box_i.text_frame; tfi.word_wrap = True
    tfi.margin_left = Inches(0.25); tfi.margin_right = Inches(0.25); tfi.margin_top = Inches(0.12)
    pi = tfi.paragraphs[0]; pi.text = "Por que importa para a tese"
    pi.font.size = Pt(13); pi.font.bold = True; pi.font.color.rgb = GREEN
    pib = tfi.add_paragraph()
    pib.text = "O Teorema 1 do LAFTR é a garantia teórica que sustenta a EXTENSÃO da nossa tese para Cap 3 (face recognition). Sem ele, teríamos que defender a transferência empiricamente sem ancoragem formal. Com ele, a extensão deixa de ser especulativa."
    pib.font.size = Pt(12); pib.font.color.rgb = GRAY_DK

    tx_f = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
    pf = tx_f.text_frame.paragraphs[0]
    pf.text = "H5 em revisão após Pangelinan 2023 (Rodada 6): pixel info pode ser confounder adicional em FR — discussão na pauta"
    pf.font.size = Pt(12); pf.font.color.rgb = GRAY_MD; pf.font.italic = True

    # Speaker notes — analogia
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "ANALOGIA PARA EXPLICAR FAIR TRANSFERÊNCIA ORALMENTE:\n\n"
        "Aprender ética profissional em medicina geral nos torna mais éticos também quando "
        "nos especializamos em cardiologia. A 'ética' é uma propriedade do profissional "
        "(da representação interna), não da especialidade (da tarefa específica).\n\n"
        "USE A ANALOGIA SE A AUDIÊNCIA PEDIR UMA INTUIÇÃO MENOS TÉCNICA — caso contrário, "
        "fique no detalhamento científico do slide.\n\n"
        "EM CASO DE PERGUNTA SOBRE O ADVERSÁRIO:\n"
        "O adversário é uma rede neural separada que recebe Z como input e tenta predizer "
        "o atributo sensível (raça). O encoder é treinado para MAXIMIZAR o erro do adversário "
        "— por isso 'adversarial'. Se o adversário NÃO consegue recuperar a raça a partir de Z, "
        "então Z não contém informação 'usável' sobre raça — e portanto qualquer classificador "
        "downstream também não pode ser injusto baseado em raça."
    )


def _add_perguntas_slide(prs: Presentation) -> None:
    """Slide das 14 perguntas de pesquisa respondidas (Q01-Q14)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    p = tx_t.text_frame.paragraphs[0]
    p.text = "As 14 perguntas de pesquisa respondidas"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = NAVY; line.line.fill.background()

    perguntas = [
        ("Q01", "Existem melhores datasets para corrigir fairness em biometria facial de raças?"),
        ("Q02", "Qual o dataset mais utilizado nos estudos de fairness?"),
        ("Q03", "Qual o split oficial do FairFace e por que diferentes papers usam splits distintos?"),
        ("Q04", "Quais técnicas de mitigação algorítmica já foram testadas em FairFace race 7-class?"),
        ("Q05", "Existe consenso sobre métrica de fairness para classificação multi-classe?"),
        ("Q06", "O baseline ResNet-34 = 72% acurácia é teto da arquitetura ou da metodologia?"),
        ("Q07", "Existe pesquisa de fairness em biometria facial sem usar FairFace?"),
        ("Q08", "Por que os estudos tentam sempre fazer merge de classes raciais?"),
        ("Q09", "7 classes é realmente a taxonomia correta para fairness racial facial?"),
        ("Q10", "Existe matriz associativa Fitzpatrick/MST × FairFace 7-race?"),
        ("Q11", "As 7 categorias raciais do FairFace têm fundamento biológico ou socio-político?"),
        ("Q12", "Como antropologia forense moderna trata 'raça'?"),
        ("Q13", "Origem e propósito da escala Fitzpatrick — para que foi criada?"),
        ("Q14", "Quantos tons de pele existem cientificamente?"),
    ]

    # Duas colunas
    col1 = perguntas[:7]
    col2 = perguntas[7:]

    tx_c1 = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(6.1), Inches(5.5))
    tf1 = tx_c1.text_frame; tf1.word_wrap = True

    for i, (code, txt) in enumerate(col1):
        p = tf1.paragraphs[0] if i == 0 else tf1.add_paragraph()
        run1 = p.add_run() if i == 0 and not p.text else None
        if i == 0:
            p.text = ""
        r_code = p.add_run()
        r_code.text = f"{code}  "
        r_code.font.bold = True
        r_code.font.color.rgb = NAVY
        r_code.font.size = Pt(13)
        r_txt = p.add_run()
        r_txt.text = txt
        r_txt.font.color.rgb = GRAY_DK
        r_txt.font.size = Pt(13)
        p.space_after = Pt(6)

    tx_c2 = slide.shapes.add_textbox(Inches(6.85), Inches(1.4), Inches(6.1), Inches(5.5))
    tf2 = tx_c2.text_frame; tf2.word_wrap = True

    for i, (code, txt) in enumerate(col2):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        if i == 0:
            p.text = ""
        r_code = p.add_run()
        r_code.text = f"{code}  "
        r_code.font.bold = True
        r_code.font.color.rgb = NAVY
        r_code.font.size = Pt(13)
        r_txt = p.add_run()
        r_txt.text = txt
        r_txt.font.color.rgb = GRAY_DK
        r_txt.font.size = Pt(13)
        p.space_after = Pt(6)

    tx_f = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
    pf = tx_f.text_frame.paragraphs[0]
    pf.text = "Detalhamento completo em docs/ativo/04_pesquisa_bibliografica/_perguntas.md"
    pf.font.size = Pt(13); pf.font.color.rgb = GRAY_MD; pf.font.italic = True


def _add_paper_card_slide(
    prs: Presentation,
    title: str,
    citation: str,
    contexto: str,
    pergunta: str,
    metodo: str,
    resultados: str,
    critica: str,
    conexoes: str,
    impacto: str,
    accent: RGBColor,
) -> None:
    """Slide de ficha completa de um paper (7 dimensões, layout 2 colunas)."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Título
    tx_t = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.6), Inches(0.6))
    p = tx_t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = NAVY

    # Citação
    tx_c = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(12.6), Inches(0.3))
    pc = tx_c.text_frame.paragraphs[0]
    pc.text = citation
    pc.font.size = Pt(11); pc.font.color.rgb = GRAY_MD; pc.font.italic = True

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.05), Inches(12.6), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = accent; line.line.fill.background()

    def _draw_box(x, y, w, h, label, conteudo, fill, border_color, has_border, label_size=10, body_size=11):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid(); box.fill.fore_color.rgb = fill
        if has_border:
            box.line.color.rgb = border_color
            box.line.width = Pt(1)
        else:
            box.line.fill.background()
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.06)
        pl = tf.paragraphs[0]; pl.text = label
        pl.font.size = Pt(label_size); pl.font.bold = True; pl.font.color.rgb = border_color
        pc_ = tf.add_paragraph(); pc_.text = conteudo
        pc_.font.size = Pt(body_size); pc_.font.color.rgb = GRAY_DK

    # Layout: coluna esquerda 6.3in, coluna direita 6.3in, gap 0.05
    col_w = 6.25
    col_left_x = 0.4
    col_right_x = 0.4 + col_w + 0.18
    y_top = 1.2

    # Coluna esquerda — Contexto teórico, Pergunta, Metodologia
    _draw_box(col_left_x, y_top, col_w, 1.5, "Contexto teórico / motivação", contexto, GRAY_LT, NAVY, False)
    _draw_box(col_left_x, y_top + 1.55, col_w, 1.45, "Pergunta de pesquisa", pergunta, WHITE, NAVY, True)
    _draw_box(col_left_x, y_top + 3.05, col_w, 2.95, "Metodologia detalhada", metodo, GRAY_LT, NAVY, False)

    # Coluna direita — Resultados, Crítica, Conexões, Impacto
    _draw_box(col_right_x, y_top, col_w, 1.6, "Resultados principais (números/achados)", resultados, GRAY_LT, NAVY, False)
    _draw_box(col_right_x, y_top + 1.65, col_w, 1.4, "Crítica metodológica e limitações", critica, WHITE, GRAY_MD, True)
    _draw_box(col_right_x, y_top + 3.1, col_w, 1.25, "Conexões com outras fichas do corpus", conexoes, GRAY_LT, NAVY, False)
    _draw_box(col_right_x, y_top + 4.4, col_w, 1.6, "Impacto na nossa tese", impacto, WHITE, accent, True, label_size=11, body_size=12)


def main() -> None:
    out = Path(__file__).resolve().parent / "material_reuniao_orientador_v3.2.1.pptx"
    build_presentation(out)
    print(f"Gerado: {out}")
    print(f"Tamanho: {out.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
