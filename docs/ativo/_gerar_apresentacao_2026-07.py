"""Gera apresentacao PowerPoint para reuniao de julho/2026.

Escrita da qualificacao fechada + 5 slides tematicos:
  1. Cronograma proposto
  2. Resumo de toda a escrita da qualificacao
  3. Pipeline experimental (6 etapas)
  4. Contribuicoes esperadas (3 eixos, 7 contribuicoes)
  5. Estrategia backbone + FiLM

Uso:
    python _gerar_apresentacao_2026-07.py
    -> produz: docs/ativo/material_reuniao_orientador_2026-07.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Paleta academica sobria (consistente com scripts anteriores)
NAVY = RGBColor(0x1F, 0x2A, 0x4E)
BLUE_MID = RGBColor(0x3D, 0x68, 0xBB)
BLUE_LIGHT = RGBColor(0xB4, 0xC4, 0xE8)
GRAY_DK = RGBColor(0x3D, 0x42, 0x4E)
GRAY_MD = RGBColor(0x70, 0x76, 0x82)
GRAY_LT = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
YELLOW = RGBColor(0xF5, 0xB7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ============================================================
# Helpers de layout
# ============================================================
def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(2.5), Inches(1.3), Inches(10.5), Inches(2.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Qualificacao — escrita fechada"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Consolidacao da monografia + revisao critica"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GRAY_DK

    p3 = tf.add_paragraph()
    p3.text = "5 capitulos, 104 fichas bibliograficas, pipeline em 6 etapas"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRAY_MD
    p3.font.italic = True

    meta = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(10.5), Inches(2.5))
    mf = meta.text_frame
    mf.word_wrap = True
    rows = [
        ("Mestrando:", "Marcello Ozzetti"),
        ("Orientador:", "Prof. Marcos Quiles"),
        ("Programa:", "Mestrado em Ciencia da Computacao — Unifesp / ICT"),
        ("Reuniao:", "Julho de 2026"),
    ]
    for i, (k, v) in enumerate(rows):
        p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        p.text = f"{k}  {v}"
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY_DK


def add_title(slide, text: str) -> None:
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.5), Inches(0.9))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, text: str) -> None:
    if not text:
        return
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
    pf = ft.text_frame.paragraphs[0]
    pf.text = text
    pf.font.size = Pt(10)
    pf.font.color.rgb = GRAY_MD
    pf.font.italic = True


def add_bullets(prs: Presentation, title: str, bullets: list, footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(5.4))
    tf = tx.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            head, body = item
            run1 = p.add_run()
            run1.text = head + "  "
            run1.font.size = Pt(15)
            run1.font.bold = True
            run1.font.color.rgb = NAVY
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(15)
            run2.font.color.rgb = GRAY_DK
        else:
            p.text = "— " + item
            p.font.size = Pt(15)
            p.font.color.rgb = GRAY_DK
        p.space_after = Pt(6)

    add_footer(slide, footer)


def add_table_slide(prs: Presentation, title: str, headers: list, rows: list, footer: str = "",
                    col_widths: list | None = None, highlight_rows: list | None = None,
                    font_size: int = 11) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.5)
    height = Inches(5.2)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = WHITE

    highlight_rows = highlight_rows or []
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = GRAY_DK

    add_footer(slide, footer)


# ============================================================
# SLIDE 1 — Cronograma proposto
# ============================================================
def slide_cronograma(prs: Presentation) -> None:
    add_table_slide(
        prs,
        "Cronograma proposto — qualificacao ate defesa",
        ["Periodo", "Etapa do pipeline", "Marco"],
        [
            ["Jul/2026", "—", "Solicitacao da qualificacao ao PPG-CC"],
            ["Set-Out/2026", "—", "EXAME DE QUALIFICACAO"],
            ["Nov/2026", "Etapa 1 — Classificador MST", "Infraestrutura + SkinToneNet + datasets"],
            ["Dez/2026", "Etapa 2 — Auditoria fenotipica FairFace", "Matriz MST x raca + validacao humana"],
            ["Jan-Mar/2027", "Etapa 3 — Classificador racial + FiLM", "ConvNeXt-T + 4 configs de ablation"],
            ["Abr/2027", "Etapa 4 — Comparacao com baselines", "6 baselines + triangulacao de metricas"],
            ["Mai/2027", "Etapa 5 — Transferencia fair", "RFW/BFW + pixel information"],
            ["Jun/2027", "Etapa 6 — Sintese decompositiva", "Decomposicao do erro Latinx"],
            ["Nov/2026 - Jul/2027", "—", "Redacao em paralelo"],
            ["2o sem 2027", "—", "DEFESA DA DISSERTACAO"],
        ],
        col_widths=[2.4, 4.5, 5.6],
        highlight_rows=[1, 9],
        footer="Mapeamento 1:1 com as 6 Etapas do Cap 4. Abril/2027 (Etapa 4) preenche gap anterior do cronograma.",
        font_size=10,
    )


# ============================================================
# SLIDE 2 — Resumo da escrita da qualificacao
# ============================================================
def slide_resumo_escrita(prs: Presentation) -> None:
    add_table_slide(
        prs,
        "Resumo da escrita — 5 capitulos consolidados",
        ["Cap", "Titulo", "O que entrega"],
        [
            ["1",
             "Introducao",
             "Contexto (NIST 2019, FaceScanPaliGemma 75.7%), problema (gap Latinx 30 pp), objetivo (FiLM+MST sobre ConvNeXt-T), 3 eixos de contribuicao, European AI Act 2024"],
            ["2",
             "Revisao bibliografica",
             "12 secoes: 6 frentes cronologicas (dados > algoritmos > VLMs > MST > heterogeneidade > refutacao) + backbones + metricas + auditoria + alternativas conditioning + 5 lacunas cientificas. 104 fichas, 5 figuras integradas"],
            ["3",
             "Objetivos, hipoteses, contribuicoes",
             "Objetivo geral + 6 objetivos especificos + 6 hipoteses testaveis (com criterios de confirmacao/refutacao) + 7 contribuicoes esperadas em 3 eixos"],
            ["4",
             "Metodologia",
             "Pipeline em 6 etapas + 4 configs ablation (A/B/C/D) + 6 baselines + escolha do ConvNeXt-T + mecanismo FiLM + triangulacao de metricas em 2 cenarios + ISO 19795-10 + rigor experimental (3 sementes)"],
            ["5",
             "Cronograma e riscos",
             "Cronograma alinhado as 6 etapas (Nov/2026 - Jun/2027) + 4 riscos com estrategias de mitigacao documentadas"],
        ],
        col_widths=[0.6, 3.0, 8.9],
        footer="Resumo PT: 488 palavras, 5 keywords. Abstract EN: 474 palavras. Ambos dentro dos limites ABNT NBR 6028.",
        font_size=11,
    )


# ============================================================
# SLIDE 3 — Pipeline (6 etapas)
# ============================================================
def slide_pipeline(prs: Presentation) -> None:
    add_table_slide(
        prs,
        "Pipeline experimental em 6 etapas",
        ["#", "Etapa", "Insumo / Modelo", "Saida esperada"],
        [
            ["1",
             "Classificador MST",
             "SkinToneNet pre-treinado (Matias 2026)",
             "Vetor MST 10-dim por imagem do FairFace"],
            ["2",
             "Auditoria fenotipica FairFace",
             "Saida da Etapa 1 sobre FairFace val",
             "Matriz publica MST x classes raciais (Contribuicao 2)"],
            ["3",
             "Classificador racial condicionado",
             "ConvNeXt-T + camadas FiLM (vetor MST)",
             "F1 macro + F1 por classe (config B principal + A/C/D em ablation)"],
            ["4",
             "Comparacao com baselines",
             "vs 6 baselines: ResNet-34, ConvNeXt-T puro, FSCL+, Group DRO, FineFACE, adv. debiasing",
             "Triangulacao de metricas (DR, worst-class F1, EO) + visualizacao Pareto"],
            ["5",
             "Transferencia fair para FR",
             "Backbone fair da Etapa 3 aplicado a RFW / BFW",
             "Ganho de fairness em face recognition + controle pixel information"],
            ["6",
             "Sintese decompositiva",
             "Resultados Etapas 2 + 5",
             "Decomposicao erro Latinx: fenotipico irredutivel vs algoritmico mitigavel"],
        ],
        col_widths=[0.5, 2.8, 4.8, 4.4],
        highlight_rows=[2, 5],
        footer="Etapas 3 e 6 concentram as contribuicoes metodologica e diagnostica principais.",
        font_size=10,
    )


# ============================================================
# SLIDE 4 — Contribuicoes esperadas (3 eixos, 7 contribuicoes)
# ============================================================
def slide_contribuicoes(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Contribuicoes esperadas — 3 eixos, 7 contribuicoes")

    # Coluna esquerda: 3 eixos
    tx1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.3))
    tf1 = tx1.text_frame
    tf1.word_wrap = True

    eixos = [
        ("EIXO FENOTIPICO-EMPIRICO",
         "Contribuicoes 1 e 2",
         "Documenta a distribuicao real do fenotipo (MST) dentro das classes raciais"),
        ("EIXO METODOLOGICO-ARQUITETURAL",
         "Contribuicoes 3, 4 e 7",
         "Introduz o tom de pele como sinal condicionante via mecanismo arquitetural"),
        ("EIXO DIAGNOSTICO-ESTRUTURAL",
         "Contribuicoes 5 e 6",
         "Decompoe o erro em componentes fenotipico e algoritmico, com transferencia fair"),
    ]

    for i, (titulo, mapa, resumo) in enumerate(eixos):
        p1 = tf1.paragraphs[0] if i == 0 else tf1.add_paragraph()
        p1.text = titulo
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.space_before = Pt(0 if i == 0 else 12)

        p2 = tf1.add_paragraph()
        p2.text = mapa
        p2.font.size = Pt(11)
        p2.font.italic = True
        p2.font.color.rgb = ACCENT

        p3 = tf1.add_paragraph()
        p3.text = resumo
        p3.font.size = Pt(11)
        p3.font.color.rgb = GRAY_DK
        p3.space_after = Pt(4)

    # Coluna direita: 7 contribuicoes numeradas
    tx2 = slide.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(8.0), Inches(5.3))
    tf2 = tx2.text_frame
    tf2.word_wrap = True

    contribs = [
        ("Contribuicao 1.", "Avaliacao comparativa de modelos MST pre-treinados com sensitivity analysis"),
        ("Contribuicao 2.", "Primeira matriz publica MST x classes raciais sobre o FairFace val set"),
        ("Contribuicao 3.", "Primeira aplicacao documentada de FiLM-conditioning a race classification multi-classe"),
        ("Contribuicao 4.", "Triangulacao de metricas multi-classe (DR + worst-class F1 + EO/EqOdds), aderente a ISO 19795-10"),
        ("Contribuicao 5.", "Demonstracao empirica de transferencia fair de classification para face recognition, com controle de pixel information"),
        ("Contribuicao 6.", "Decomposicao quantitativa do erro Latinx: fenotipico irredutivel vs algoritmico mitigavel"),
        ("Contribuicao 7.", "Estudo comparativo entre 4 configuracoes de conditioning (baseline, FiLM linear, FiLM com porta, FiLM CLIP)"),
    ]

    for i, (head, body) in enumerate(contribs):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        run1 = p.add_run()
        run1.text = head + "  "
        run1.font.size = Pt(11)
        run1.font.bold = True
        run1.font.color.rgb = NAVY
        run2 = p.add_run()
        run2.text = body
        run2.font.size = Pt(11)
        run2.font.color.rgb = GRAY_DK
        p.space_after = Pt(5)

    add_footer(slide, "Cada contribuicao esta vinculada a um objetivo especifico e a uma hipotese testavel do Cap 3.")


# ============================================================
# SLIDE 5 — Estrategia Backbone + FiLM
# ============================================================
def slide_backbone_film(prs: Presentation) -> None:
    add_bullets(
        prs,
        "Estrategia: ConvNeXt-T como backbone + FiLM como conditioning",
        [
            ("BACKBONE — ConvNeXt-T (Liu et al. 2022, CVPR)", ""),
            ("Por que ConvNeXt-T:",
             "28M parametros, CNN moderna com LayerNorm e depthwise 7x7 (inspirada em ViT)"),
            ("Nao e ViT:",
             "Permanece dentro do paradigma convolucional — 82.0% top-1 ImageNet-1K, supera Swin-T (81.3%)"),
            ("Comparabilidade:",
             "Config A e o ConvNeXt-T puro (isola contribuicao do backbone vs ResNet-34 canonico)"),
            ("Compatibilidade:",
             "4 estagios hierarquicos = 4 pontos naturais de insercao de camadas FiLM"),
            ("", ""),
            ("MECANISMO — FiLM (Perez et al. 2018, AAAI)", ""),
            ("Formulacao:",
             "FiLM(x | gamma, beta) = gamma * x + beta, modulacao afim canal-a-canal"),
            ("Sinal condicionante:",
             "Vetor MST 10-dim produzido pelo SkinToneNet (Etapa 1)"),
            ("Overhead:",
             "~380k parametros adicionais (~1.3% sobre o backbone) — parameter-efficient"),
            ("Lacuna coberta:",
             "Primeira aplicacao documentada de FiLM em race classification multi-classe (Contribuicao 3)"),
            ("Ablation:",
             "Config B (linear, principal) vs C (porta multiplicativa) vs D (CLIP-text embedding)"),
        ],
        footer="Justificativa detalhada em Cap 4, Sec. 5 (Escolha do backbone) e Sec. 6 (Mecanismo FiLM).",
    )


# ============================================================
# Main
# ============================================================
def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    slide_cronograma(prs)
    slide_resumo_escrita(prs)
    slide_pipeline(prs)
    slide_contribuicoes(prs)
    slide_backbone_film(prs)

    return prs


def main() -> None:
    prs = build_presentation()
    out_dir = Path(__file__).parent
    out = out_dir / "material_reuniao_orientador_2026-07.pptx"
    prs.save(out)
    print(f"OK: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
