"""Gera apresentação PowerPoint do material da reunião.

Conteúdo derivado de: docs/ativo/material_reuniao_orientador_2026-06.md

Uso:
    python _gerar_apresentacao.py
    → produz: docs/ativo/material_reuniao_orientador_2026-06.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- Paleta acadêmica sóbria ---
NAVY = RGBColor(0x1F, 0x2A, 0x4E)
GRAY_DK = RGBColor(0x3D, 0x42, 0x4E)
GRAY_MD = RGBColor(0x70, 0x76, 0x82)
GRAY_LT = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)  # vermelho-tinto para destaques (frentes 🔬)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Faixa lateral navy
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Título principal
    tx = slide.shapes.add_textbox(Inches(2.5), Inches(1.5), Inches(10.5), Inches(2.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Equidade Racial em Classificação Facial"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Auditoria de viés em sistemas de visão computacional sobre o dataset FairFace"
    p2.font.size = Pt(20)
    p2.font.color.rgb = GRAY_DK

    # Bloco de metadados
    meta = slide.shapes.add_textbox(Inches(2.5), Inches(4.2), Inches(10.5), Inches(2.5))
    mf = meta.text_frame
    mf.word_wrap = True

    rows = [
        ("Mestrando:", "Marcello Ozzetti"),
        ("Orientador:", "Prof. Marcos Quiles"),
        ("Programa:", "Mestrado em Ciência da Computação — Unifesp / ICT"),
        ("Reunião:", "Junho de 2026 (segunda)"),
        ("Versão da tese:", "v3.2 — prescritiva, pós-feedback orientador"),
    ]
    for i, (k, v) in enumerate(rows):
        p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        run_k = p.add_run() if i > 0 else None
        if i == 0:
            p.text = f"{k}  {v}"
        else:
            p.text = f"{k}  {v}"
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY_DK


def add_section_divider(prs: Presentation, num: str, title: str, subtitle: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Fundo navy
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Número grande
    tx_num = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3), Inches(3))
    pn = tx_num.text_frame.paragraphs[0]
    pn.text = num
    pn.font.size = Pt(140)
    pn.font.bold = True
    pn.font.color.rgb = WHITE

    # Título
    tx_t = slide.shapes.add_textbox(Inches(4.3), Inches(2.5), Inches(8.5), Inches(2))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = GRAY_LT


def add_content_slide(prs: Presentation, title: str, bullets: list, footer: str = "") -> None:
    """Slide padrão com título + bullet list."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Título
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Linha separadora
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    # Body
    tx_b = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.3), Inches(5.4))
    bf = tx_b.text_frame
    bf.word_wrap = True

    for i, b in enumerate(bullets):
        if isinstance(b, tuple):
            text, level = b
        else:
            text, level = b, 0

        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = text
        p.level = level
        if level == 0:
            p.font.size = Pt(20)
            p.font.color.rgb = GRAY_DK
            p.font.bold = False
        else:
            p.font.size = Pt(17)
            p.font.color.rgb = GRAY_MD
        p.space_after = Pt(8)

    # Footer
    if footer:
        tx_f = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.4))
        pf = tx_f.text_frame.paragraphs[0]
        pf.text = footer
        pf.font.size = Pt(15)
        pf.font.color.rgb = GRAY_MD
        pf.font.italic = True


def add_table_slide(prs: Presentation, title: str, headers: list, rows: list, footer: str = "") -> None:
    """Slide com título + tabela."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Título
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    p = tx_t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Linha separadora
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    # Tabela
    nrows = len(rows) + 1
    ncols = len(headers)
    table_h = Inches(5.4)
    table = slide.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.5), Inches(12.5), table_h).table

    # Headers
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.color.rgb = WHITE
                r.font.size = Pt(16)

    # Rows
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(v)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else GRAY_LT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(15)
                    r.font.color.rgb = GRAY_DK

    if footer:
        tx_f = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.4))
        pf = tx_f.text_frame.paragraphs[0]
        pf.text = footer
        pf.font.size = Pt(15)
        pf.font.color.rgb = GRAY_MD
        pf.font.italic = True


def add_thesis_slide(prs: Presentation, title: str, statement: str) -> None:
    """Slide especial para a tese — formato destacado."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Título
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    p = tx_t.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = NAVY

    # Quote box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.4))
    box.fill.solid()
    box.fill.fore_color.rgb = GRAY_LT
    box.line.color.rgb = NAVY
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.4)
    tf.margin_bottom = Inches(0.4)

    p = tf.paragraphs[0]
    p.text = statement
    p.font.size = Pt(20)
    p.font.color.rgb = NAVY
    p.font.italic = True
    p.alignment = PP_ALIGN.LEFT


def build_presentation(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- SLIDE 1: TÍTULO ---
    add_title_slide(prs)

    # --- SLIDE 2: SUMÁRIO ---
    add_content_slide(
        prs,
        "Sumário",
        [
            "1. Metodologia de pesquisa",
            "2. Racional — historyline da evolução da pesquisa",
            "3. Visão dos principais artigos",
            "4. Gaps, considerações e achados",
            "5. Próximos passos",
        ],
        footer="23 fichas catalogadas | 14 perguntas respondidas | 5 frentes de pesquisa identificadas",
    )

    # ==================== CAP 1 ====================
    add_section_divider(prs, "1", "Metodologia de pesquisa", "5 critérios + 4 fases + rigor de citação")

    add_content_slide(
        prs,
        "1.1  Gatilho da reestruturação (25/05/2026)",
        [
            "Três problemas estruturais identificados pelo orientador:",
            ("(i) Síntese baseada em abstracts, não leitura integral", 1),
            ("(ii) Autoria incorreta em parte das citações principais", 1),
            ("(iii) Framing 'evolução do MBA' recusado", 1),
            "",
            "Resposta: protocolo formal de Pesquisa Bibliográfica em 4 fases, com:",
            ("• Critérios de seleção explícitos", 1),
            ("• Verificação de autoria em fonte primária (arXiv DOI)", 1),
            ("• Leitura integral obrigatória dos PDFs", 1),
            ("• Catalogação normativa em fichas estruturadas", 1),
        ],
    )

    add_content_slide(
        prs,
        "1.2  Os 5 critérios de seleção (ordem hierárquica)",
        [
            "1. Escopo temático — fairness / biometria / raça / auditoria em sistemas faciais",
            "2. Recência — 2019–2026 (10 anos pós-Buolamwini & Gebru) + exceções seminais",
            "3. Verificabilidade — autoria/ano/venue confirmáveis em fonte primária",
            "4. Relevância editorial — CVPR, NeurIPS, ICCV, ECCV, ICML, ICLR, FAccT, TBIOM, TPAMI, ACM CSur, Information Fusion",
            "5. Impacto (citações) — ≥50 (≥3 anos), ≥20 (2024), aberto para 2025–26 com cobertura única",
            "",
            "Preprints arXiv aceitos apenas com cobertura única OU impacto demonstrado.",
        ],
    )

    add_content_slide(
        prs,
        "1.3  Fluxo de trabalho em 4 fases",
        [
            "Fase A — Coleta sistemática",
            ("Google Scholar, Semantic Scholar, arXiv, DBLP, anais de venues-alvo", 1),
            "Fase B — Triagem editorial",
            ("Critérios 1–5 aplicados; decisões registradas em _triagem.md", 1),
            ("Decisões: APROVADO | STANDBY | DESCARTADO com justificativa explícita", 1),
            "Fase C — Leitura integral",
            ("Download PDF; extração de texto via Python (pypdf)", 1),
            ("Ficha normativa em 11 seções por paper", 1),
            "Fase D — Síntese transversal",
            ("Cross-reference matricial; identificação de gaps; thesis statement v3", 1),
        ],
    )

    add_content_slide(
        prs,
        "1.4  Template de ficha — 11 seções normativas",
        [
            "Toda ficha segue estrutura idêntica:",
            ("1. Frontmatter YAML (metadata: autores, venue, citações, status)", 1),
            ("2. Resumo do problema atacado", 1),
            ("3. Método", 1),
            ("4. Datasets e setup experimental", 1),
            ("5. Métricas reportadas", 1),
            ("6. Resultados principais (valores numéricos)", 1),
            ("7. Limitações declaradas pelos autores", 1),
            ("8. Limitações que identifiquei (leitura crítica)", 1),
            ("9. Relação com nossa pesquisa", 1),
            ("10. Pontos para citar / posicionar", 1),
            ("11. Trabalhos sugeridos pelos autores (Future Work)", 1),
        ],
        footer="A Seção 11 foi incorporada na Rodada 4 para mapear 'gaps endossados pelo campo'",
    )

    add_content_slide(
        prs,
        "1.5  Metodologia Q&A e protocolo de rigor bibliográfico",
        [
            "Metodologia Q&A — interrogação ativa do corpus:",
            ("• 14 perguntas de pesquisa formuladas e respondidas em _perguntas.md", 1),
            ("• Status: ANSWERED | PARTIAL | OPEN | NEW RESEARCH FRONT", 1),
            ("• 5 perguntas geraram frentes de pesquisa (alimentam 06_gap.md)", 1),
            "",
            "Verificação de autoria — toda citação passa por:",
            ("1. Download do PDF original (arXiv ou DOI editorial)", 1),
            ("2. Inspeção da primeira página para confirmação visual", 1),
            ("3. Cross-check em Semantic Scholar (citações + venue)", 1),
            ("4. Registro em 00_referencias.md com status VERIFIED", 1),
            "",
            "Princípio operacional: sem verificação, sem citação.",
        ],
    )

    # ==================== CAP 2 ====================
    add_section_divider(prs, "2", "Racional / Historyline", "Evolução da pesquisa em 4 rodadas + 2.5")

    add_table_slide(
        prs,
        "2.1  Mapa temporal das 4 rodadas",
        ["Rodada", "Foco da rodada", "Papers", "Achado principal"],
        [
            ["R1 — Seeds iniciais", "Cobertura temática direta", "9", "FairFace + Gender Shades + AlDahoul formam a base"],
            ["R2 — Snowballing", "Técnicas metodológicas", "+5", "Mitigações existem (FSCL+, Group DRO) mas nunca em race 7-class"],
            ["R2.5 — Verificação", "Triangulação de SOTA", "—", "FaceScanPaliGemma 75.7% confirmado como SOTA único"],
            ["R3 — Broadening", "Tracks paralelos (não-FairFace)", "+5", "BFW, Casual Conv, MST-E ampliam o panorama"],
            ["R4 — Fundamentação", "Substrato teórico", "+4", "AAPA + Lewontin: race não é biologia (fundamenta a tese)"],
        ],
        footer="Total: 23 fichas catalogadas | todas com autoria verificada em fonte primária",
    )

    add_content_slide(
        prs,
        "2.2  Aprendizados metodológicos do processo",
        [
            "1. A literatura confirma o problema mas raramente executa a solução",
            ("7 papers sugerem mitigação em race classification multi-classe; ZERO executam", 1),
            "",
            "2. Categorias raciais são instrumentais, não biológicas",
            ("A escolha do FairFace 7-class é pragmática e alinhada com a literatura dominante", 1),
            ("Mesma escolha que AlDahoul (SOTA), FairGRAPE, FineFACE etc.", 1),
            "",
            "3. A confusão Fitzpatrick ↔ race é endêmica",
            ("1/3 dos próprios dermatologistas confunde. Adotar MST resolve no nível instrumental.", 1),
            "",
            "4. O método Q&A é o ativo principal",
            ("Transformou leitura passiva em 5 hipóteses falsificáveis (H1–H5)", 1),
        ],
    )

    # ==================== CAP 3 ====================
    add_section_divider(prs, "3", "Visão dos principais artigos", "10 papers centrais detalhados")

    add_table_slide(
        prs,
        "3.0  Achados por track temático (7 tracks após R5)",
        ["Track", "Foco", "Achado central"],
        [
            ["A — Race classification", "Treino e teste sobre raça", "SOTA = 75.7% F1; Latinx persistentemente em 60% F1"],
            ["B — Face recognition fairness", "Verification balanceada", "Neto 2025 questiona discretização per se"],
            ["C — Skin tone (Fitzpatrick / MST)", "Auditoria com escala física", "MST > Fitzpatrick para fairness research"],
            ["D — Mitigação algorítmica", "Técnicas anti-viés", "7 papers sugerem; ZERO executam em race 7-class FairFace"],
            ["E — Auditoria e metodologia", "Surveys e ferramentas", "Kotwal 2025 (TBIOM) é referência mais atualizada"],
            ["F — Fundamentação científica", "Substrato teórico", "Race ≠ biologia: AAPA 2019 + Lewontin 1972 (85% intra-pop)"],
            ["G — Mecanismos ML / RN (NOVO R5)", "FiLM, FRL, LAFTR, EO/EOD, impossibility", "Toolkit formal para pipeline v3.2 (FiLM + LAFTR centrais)"],
        ],
        footer="A seguir: 10 papers detalhados (1 por slide) | Total corpus: 29 fichas, 7 tracks",
    )

    # 10 main papers
    papers = [
        ("Fuentes et al. 2019 — AAPA Statement on Race and Racism",
         "Am J Phys Anthropol, vol 169, no 3 | Statement institucional AABA",
         [
             "Papel na tese: FUNDAMENTO TEÓRICO — 'race is not biology'",
             "Citação central: 'Race does not provide an accurate representation of human biological variation'",
             "Posição consolidada da comunidade de antropologia biológica desde 1996, atualizada 2019",
             "Sustenta a limitação reconhecida em §6.1 da tese v3.1 com citação institucional formal",
         ]),
        ("Lewontin 1972 — The Apportionment of Human Diversity",
         "Evolutionary Biology vol 6, Springer | 5000+ citações, 50+ anos de literatura derivada",
         [
             "Papel na tese: FUNDAMENTO GENÉTICO",
             "Partição clássica: 85.4% intra-pop / 8.3% between-pop within-race / 6.3% between-race",
             "Confirmado por estudos genômicos modernos (Rosenberg 2002, Bergstrom 2020)",
             "Sobreposição fenotípica entre 7 classes FairFace, especialmente Latinx, é previsão direta",
             "Caveat: Edwards 2003 ('Lewontin's Fallacy') — classificação possível com loci suficientes, mas isso não justifica taxonomia binária discreta",
         ]),
        ("Buolamwini & Gebru 2018 — Gender Shades",
         "PMLR vol 81 / ACM FAccT | 4 933 citações Semantic Scholar (VERIFICADO)",
         [
             "Papel na tese: MARCO FUNDADOR do campo de auditoria de fairness em CV comercial",
             "Achado-bandeira: gap de 34.7 pp na taxa de erro (lighter male 0% vs darker female 34.7% em IBM)",
             "Decisão metodológica: escolhem Fitzpatrick (não race) por argumento de instabilidade",
             "Precedente direto da nossa Q11 (race não tem fundamento biológico)",
         ]),
        ("Kärkkäinen & Joo 2021 — FairFace",
         "WACV 2021 IEEE/CVF | 263+ citações Google Scholar",
         [
             "Papel na tese: DATASET CENTRAL — toda execução experimental ocorre sobre FairFace",
             "Estrutura: 108 501 imagens; 7 categorias raciais; origem YFCC-100M Flickr; anotação MTurk 3-anotador",
             "Caveat crítico: cross-dataset evaluation usa só 4 raças — não 7",
             "7-class evaluation só aparece em external sets (Twitter/Media/Protest), com Latino F1 = 0.247",
         ]),
        ("AlDahoul et al. 2024/2026 — FaceScanPaliGemma",
         "Nature Scientific Reports 2026 (DOI 10.1038/s41598-026-39584-3) + arXiv 2410.24148",
         [
             "Papel na tese: SOTA ATUAL para FairFace race 7-class",
             "Achado: 75.7% accuracy / 75% F1 macro — número que precisamos alcançar/superar",
             "Baseline FairFace ResNet-34 confirmado em 72% (single seed, sem CI)",
             "Per-class: Black 90% | White 80% | Indian 79% | East Asian 78% | ME 73% | SE Asian 67% | LATINX 60% (pior)",
         ]),
        ("Lin, Kim & Joo 2022 — FairGRAPE",
         "ECCV 2022 | 49 citações Semantic Scholar",
         [
             "Papel na tese: VALIDAÇÃO CRUZADA do baseline 72% + documentação independente do gap Hispanic",
             "Confirmação: ResNet-34 baseline = 72.0% sobre FairFace race 7-class (Tabela 2)",
             "Per-class: Hispanic 59.6% | Black 83.2% | White 73.9% — MESMO padrão hierárquico que AlDahoul",
             "Implicação metodológica: 72% é número robusto através de splits e seeds independentes",
         ]),
        ("Schumann et al. 2023 — MST Annotation Consensus",
         "NeurIPS 2023 Datasets and Benchmarks Track",
         [
             "Papel na tese: PROTOCOLO DE ANOTAÇÃO MST para a frente Q10",
             "Contribuições: (i) Monk Skin Tone Scale 10 pontos; (ii) MST-E dataset (1 515 imgs)",
             "Achado: anotadores de diferentes regiões geográficas têm variação SISTEMÁTICA em MST",
             "Recomendação operacional: pool diverso + alta replicação — importável para Q10",
         ]),
        ("Hazirbas et al. 2021 — Casual Conversations",
         "CVPRW 2021 (Meta / Facebook AI)",
         [
             "Papel na tese: PARADIGMA ALTERNATIVO de anotação — self-reported demographics",
             "Estrutura: 3 011 sujeitos pagos, 45 000+ vídeos; age + gender SELF-REPORTED; Fitzpatrick anotado",
             "Argumento explícito contra race labels: 'Raters may have unconscious biases'",
             "Justifica adotar MST + self-identification como gold standard para Q01",
         ]),
        ("Park et al. 2022 — Fair Supervised Contrastive Learning (FSCL)",
         "CVPR 2022 | 116 citações Semantic Scholar",
         [
             "Papel na tese: TÉCNICA CENTRAL do Capítulo 2 (mitigação)",
             "Teorema 1: SupCon (Khosla 2020) incentiva encoder a aprender atributos sensíveis em datasets enviesados",
             "FSCL+: restringir negativos a mesmo grupo sensível + group-wise normalization",
             "Resultado em CelebA: reduz EO de 30.5 → 6.5 (4.7× menor), perdendo só 1.4 pp de accuracy",
             "Lacuna crítica: NÃO testado em FairFace race 7-class — exatamente nosso gap",
         ]),
        ("Sagawa et al. 2020 — Group DRO",
         "ICLR 2020",
         [
             "Papel na tese: TÉCNICA ALTERNATIVA do Capítulo 2",
             "Achado paradoxal: DRO naive ≈ ERM em redes overparameterized — ambos memorizam treino",
             "Solução: strong ℓ2 regularization (λ ∈ [0.1, 1.0], 4 ordens de magnitude acima do default) OU early stopping",
             "Resultado em CelebA hair color: worst-group accuracy 41.1 → 86.7 (+45 pp); só 1.3 pp perdidos em média",
             "Implicação: se aplicarmos a race 7-class, precisamos sweep de λ — não plug-and-play",
         ]),
    ]

    for title, subtitle, bullets in papers:
        add_content_slide(prs, title, [(subtitle, 0), ("", 0)] + [(b, 0) for b in bullets])

    # ==================== CAP 4 ====================
    add_section_divider(prs, "4", "Gaps, considerações e achados", "5 frentes de pesquisa identificadas")

    add_table_slide(
        prs,
        "4.1  Achados convergentes da literatura",
        ["Achado", "Papers endossando", "Implicação"],
        [
            ["Balanceamento de dataset não basta", "8+ (Karkkainen, Wang, Grother, Lin, AlDahoul, Kotwal, Sagawa, Klare)", "Mitigação algorítmica é necessária"],
            ["Latinx/Hispanic é estruturalmente mais difícil", "4 (Karkkainen .247, AlDahoul .60, Lin 59.6%, Buolamwini)", "Sugere causa estrutural (Q10)"],
            ["Race é construto social", "AAPA 2019, Lewontin 1972, Karkkainen explícito", "Tese reconhece e cita formalmente"],
            ["Ensembles naive não funcionam", "Bhaskaruni 2019, Sagawa 2020, Park 2022", "Mitigação exige ponderação demográfica explícita"],
            ["Métricas multi-classe são fragmentadas", "7 papers com métricas diferentes", "Q05 — triangulação DR + worst-class + CV"],
        ],
    )

    add_content_slide(
        prs,
        "4.2  Achados não-óbvios da Rodada 4",
        [
            "1. As 7 categorias do FairFace NÃO têm fundamento biológico",
            ("Construto socio-político derivado do US Census + ajustes Karkkainen", 1),
            ("AAPA 2019 (statement institucional) + Lewontin 1972 (85% intra-pop) sustentam", 1),
            "",
            "2. Fitzpatrick scale foi criada em 1975 para dosimetria PUVA, NÃO para classificação racial",
            ("1/3 dos dermatologistas ainda confunde — erro categorial endêmico", 1),
            "",
            "3. Não há 'número correto' de tons de pele",
            ("5 escalas em uso, cada uma para propósito específico", 1),
            ("MST 10-pt é a correta para fairness research (Schumann 2023, padrão Google)", 1),
            "",
            "4. A 'dificuldade Latinx' tem fundamento estatístico previsível",
            ("85% da variação genética é intra-populacional (Lewontin)", 1),
            ("Latinx abrange múltiplos clusters geográficos — sobreposição é previsão", 1),
        ],
    )

    add_table_slide(
        prs,
        "4.3  As 5 frentes de pesquisa — ranqueamento",
        ["#", "Frente", "Pergunta", "Score O+V+I", "Decisão"],
        [
            ["Q04", "Mitigação algorítmica em race 7-class", "Qual técnica de Track D funciona melhor em FairFace 7-class?", "13/15", "CAPÍTULO PRINCIPAL"],
            ["Q10", "Matriz MST × FairFace 7-race", "Como o tom de pele se distribui entre as 7 categorias raciais?", "13/15", "CAPÍTULO PRINCIPAL"],
            ["Q05", "Métrica fairness multi-classe", "Triangulação DR + worst-class + CV é robusta?", "11/15", "Adotar transversalmente"],
            ["Q06", "Decomposição do ceiling 72%", "Quanto é arquitetura vs metodologia vs dados?", "10/15", "Ablação dentro de Q04"],
            ["Q01", "Confiabilidade anotação Latinx", "Inter-annotator agreement é uniforme entre classes?", "9/15", "Fundir com Q10"],
        ],
        footer="O = Originalidade | V = Viabilidade | I = Impacto (cada de 0–5)",
    )

    add_content_slide(
        prs,
        "4.4  Decisão de escopo: Q04 + Q10 são complementares",
        [
            "Q04 ataca o 'como mitigar' (técnicas algorítmicas)",
            ("Forte endosso da literatura: 7 papers sugerem essa direção", 1),
            ("Execução zero: ninguém testou em FairFace race 7-class", 1),
            "",
            "Q10 ataca o 'por que existe erro' (diagnóstico fenotípico)",
            ("Zero precedentes: matriz MST × race nunca foi construída em contexto de fairness", 1),
            ("Originalidade total — Draelos 2025 é dermatologia, dados não-públicos", 1),
            "",
            "Combinadas — decomposição inédita na literatura de race classification:",
            ("Q04 mostra QUANTO se pode reduzir do erro", 1),
            ("Q10 mostra QUANTO sobra como irredutível", 1),
            ("Síntese: erro_total = erro_fenotípico_irredutível + erro_redutível_modelo", 1),
        ],
    )

    # ==================== CAP 5 ====================
    add_section_divider(prs, "5", "Próximos passos", "Tese v3.2 prescritiva + pipeline integrado")

    add_thesis_slide(
        prs,
        "5.1  Tese v3.2 — prescritiva (pós-reunião)",
        '"A incorporação explícita de informação de tom de pele (Monk Skin Tone) como sinal '
        "auxiliar condicionante no treinamento de classificadores faciais profundos melhora "
        "métricas de fairness demográfica em classificação racial multi-classe e estende essa "
        "melhoria a tarefas downstream de reconhecimento facial em grupos sub-representados.\n\n"
        "Contribuições originais:\n"
        "(C1) Pipeline MST → FiLM → race classifier — primeira instância documentada\n"
        "(C2) Matriz pública P(MST | race) sobre FairFace — zero precedentes\n"
        "(C3) Demonstração empírica de fair transferência race → face recognition (LAFTR)\n"
        "(C4) Triangulação DR + worst-class F1 + CV como métrica padrão multi-classe\n"
        "(C5) Quantificação fenotípico vs algorítmico do erro Latinx",
    )

    add_content_slide(
        prs,
        "5.2  Pipeline operacional (6 fases)",
        [
            "Fase 1 — Classifier MST auxiliar (Schumann 2023)",
            ("Treinar sobre MST-E + Casual Conversations", 1),
            "Fase 2 — Auditoria da matriz P(MST | race) sobre FairFace val",
            ("10 954 imgs automatizadas + subset 500-700 manual (3 anotadores)", 1),
            "Fase 3 — Race classifier condicionado",
            ("ConvNeXt-T (28M) + camadas FiLM (Perez et al. 2018) condicionadas por f_MST", 1),
            ("Treino 3-seed casado sobre FairFace train", 1),
            "Fase 4 — Comparação fairness com vs sem MST",
            ("Baselines: ConvNeXt-T vanilla, FSCL+ (Park), Group DRO (Sagawa), Adversarial (Zhang)", 1),
            "Fase 5 — Extensão a face recognition",
            ("Pipeline análogo em RFW (4 etnias) ou BFW (4×2 raça × gênero)", 1),
            "Fase 6 — Validação da fair transferência (LAFTR-style, Madras 2018)",
            ("Métrica primária: TAR @ FAR fixo por raça, foco em Black/African", 1),
        ],
    )

    add_table_slide(
        prs,
        "5.3  Hipóteses falsificáveis v3.2 (H1–H5)",
        ["ID", "Hipótese", "Critério de confirmação"],
        [
            ["H1 (CENTRAL)", "Pipeline MST → FiLM → ConvNeXt-T supera baseline ResNet-34 em F1 ≥ +2 pp E reduz DR ≥ 20%", "Ambos simultaneamente"],
            ["H2", "ConvNeXt-T vanilla (sem MST) ganha +2 a +5 pp F1 sobre ResNet-34; Latinx F1 ≈ 60% (±3 pp)", "Ganho no range E Latinx invariante"],
            ["H3", "Matriz P(MST | race) revela Latinx spread ≥ 5 categorias MST com pico distribuído", "Spread ≥ 5, não-concentrado"],
            ["H4 (CENTRAL)", "≥ 50% das misclassificações Latinx→outras estão em zonas MST de sobreposição", "%_overlap ≥ 50%"],
            ["H5 (CENTRAL)", "Pipeline em face recognition (RFW/BFW) melhora accuracy Black/African ≥ +3 pp", "Ganho ≥ 3 pp sobre baseline sem MST"],
        ],
        footer="H1, H4 e H5 são CENTRAIS. Plano B documentado em 07_thesis_statement.md §6.4 — observação negativa é cientificamente válida.",
    )

    add_content_slide(
        prs,
        "5.4  Três capítulos experimentais (v3.2)",
        [
            "Capítulo 1 — Classifier MST + matriz P(MST | race) — ~4 semanas",
            ("Fase 1: Treinar classifier MST sobre MST-E + Casual Conversations", 1),
            ("Fase 2: Aplicar sobre FairFace val + validação manual subset", 1),
            ("Testa H3 (spread Latinx ≥ 5 categorias MST)", 1),
            "",
            "Capítulo 2 — Race classifier condicionado por MST — ~10-12 semanas",
            ("Baselines: ResNet-34 vanilla, ConvNeXt-T vanilla, FSCL+, Group DRO, Adversarial", 1),
            ("Pipeline: ConvNeXt-T + FiLM condicionado em f_MST (Perez 2018)", 1),
            ("3 seeds casados, comparação pareada com paired t-test", 1),
            ("Testa H1 (pipeline supera baseline), H2 (Latinx invariante), H4 (overlap)", 1),
            "",
            "Capítulo 3 — Extensão a face recognition — ~6 semanas (NOVO)",
            ("Pipeline análogo em RFW ou BFW (verification 1:1)", 1),
            ("Métrica: TAR @ FAR por raça, foco em Black/African", 1),
            ("Testa H5 (fair transferência LAFTR-style — Madras 2018)", 1),
        ],
    )

    add_table_slide(
        prs,
        "5.5  Cronograma estimado v3.2 (28-32 semanas)",
        ["Bloco", "Duração", "Marco"],
        [
            ["Aprovação tese v3.2 + plano", "Esta reunião", "Ajustes do orientador"],
            ["Setup metodológico (02, 03, 08)", "2 semanas", "Especificações executáveis"],
            ["Cap 1 — Classifier MST + matriz", "4 semanas", "H3 testada"],
            ["Cap 2 — Race com MST conditioning", "10–12 semanas", "H1, H2, H4 testadas"],
            ["Cap 3 — Face recognition extension", "6 semanas", "H5 testada (LAFTR-style)"],
            ["Síntese — decomposição final", "4 semanas", "5 hipóteses joint"],
            ["Escrita capítulos (paralelo)", "8–12 semanas", "Defesa"],
            ["Total estimado", "~28–32 semanas (6–8 meses ativos)", "Defesa ~Jan–Mar 2027"],
        ],
    )

    add_table_slide(
        prs,
        "5.6  Riscos identificados e mitigação (v3.2)",
        ["Risco", "Probabilidade", "Mitigação"],
        [
            ["H1, H4 ou H5 refutada", "Médio", "Plano B em 07_thesis §6.4 — observação negativa é válida cientificamente"],
            ["Recrutamento de anotadores manuais (Cap 1 Fase 2)", "Alto", "Prolific Academic com filtro regional (~$1400 por 700 imgs × 3)"],
            ["Adaptação multi-classe FSCL+ / Group DRO", "Médio", "Github code disponível para todos os baselines"],
            ["FiLM com saída MST (10 logits) como contexto", "Baixo", "Mecanismo bem documentado; ablation com/sem FiLM trivial"],
            ["Compute para 3 seeds × 5 baselines em Cap 2", "Médio", "ConvNeXt-T 28M é leve; estimativa ~200-400 GPU-hours total"],
            ["Transferência MST→face recognition não generaliza (H5)", "Médio", "Achado interessante per se; reformular como limitação"],
        ],
    )

    add_content_slide(
        prs,
        "5.7  Outputs esperados (v3.2)",
        [
            "Acadêmicos:",
            ("3 capítulos de dissertação completos", 1),
            ("1 paper standalone sobre matriz MST × race (Cap 1 + Cap 2)", 1),
            ("Possível paper sobre fair transferência race → face recognition (Cap 3)", 1),
            ("Submissão alvo: CVPRW Fair AI, FAccT, ou IEEE TBIOM", 1),
            "",
            "Recursos públicos:",
            ("Código open-source: classifier MST + pipeline FiLM + adaptações multi-classe", 1),
            ("Matriz MST × FairFace 7-race pública como CSV (primeira do tipo)", 1),
            ("Dataset auxiliar: MST labels para subset do FairFace val", 1),
            "",
            "Defesa prevista: Jan-Mar 2027",
        ],
    )

    # ==================== SLIDE FINAL ====================
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Discussão e feedback"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Aprovação da tese v3.1 e plano experimental"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GRAY_LT
    p2.alignment = PP_ALIGN.CENTER

    # --- save ---
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    out = Path(__file__).resolve().parent / "material_reuniao_orientador_v3.2.pptx"
    build_presentation(out)
    print(f"Gerado: {out}")
    print(f"Tamanho: {out.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
