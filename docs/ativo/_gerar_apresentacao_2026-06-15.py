"""Gera apresentacao PowerPoint da reuniao 2026-06-15.

Uso:
    python _gerar_apresentacao_2026-06-15.py
    -> produz: docs/ativo/material_reuniao_orientador_2026-06-15.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Paleta academica sobria
NAVY = RGBColor(0x1F, 0x2A, 0x4E)
GRAY_DK = RGBColor(0x3D, 0x42, 0x4E)
GRAY_MD = RGBColor(0x70, 0x76, 0x82)
GRAY_LT = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
YELLOW = RGBColor(0xF5, 0xB7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(2.5), Inches(1.3), Inches(10.5), Inches(2.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Evolucao da Semana"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Corpus consolidado (101 fichas) + analise CLIP vs FiLM"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GRAY_DK

    p3 = tf.add_paragraph()
    p3.text = "Pos pente fino critico do corpus e resposta a sugestao do orientador"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRAY_MD
    p3.font.italic = True

    meta = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(10.5), Inches(2.5))
    mf = meta.text_frame
    mf.word_wrap = True
    rows = [
        ("Mestrando:", "Marcello Ozzetti"),
        ("Orientador:", "Prof. Marcos Quiles"),
        ("Programa:", "Mestrado em Ciencia da Computacao - Unifesp / ICT"),
        ("Reuniao:", "15 de junho de 2026"),
    ]
    for i, (k, v) in enumerate(rows):
        p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        p.text = f"{k}  {v}"
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY_DK


def add_section_divider(prs: Presentation, num: str, title: str, subtitle: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tx_num = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3), Inches(3))
    pn = tx_num.text_frame.paragraphs[0]
    pn.text = num
    pn.font.size = Pt(140)
    pn.font.bold = True
    pn.font.color.rgb = WHITE

    tx_t = slide.shapes.add_textbox(Inches(4.3), Inches(2.5), Inches(8.5), Inches(2))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = GRAY_LT


def add_title(slide, text: str) -> None:
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, prs, text: str) -> None:
    if not text:
        return
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
    pf = ft.text_frame.paragraphs[0]
    pf.text = text
    pf.font.size = Pt(10)
    pf.font.color.rgb = GRAY_MD
    pf.font.italic = True


def add_bullets(prs: Presentation, title: str, bullets: list, footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(5.4))
    tf = tx.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            head, body = item
            run1 = p.add_run()
            run1.text = head + "  "
            run1.font.size = Pt(16)
            run1.font.bold = True
            run1.font.color.rgb = NAVY
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(16)
            run2.font.color.rgb = GRAY_DK
        else:
            p.text = "- " + item
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_DK
        p.space_after = Pt(8)

    add_footer(slide, prs, footer)


def add_table_slide(prs: Presentation, title: str, headers: list, rows: list, footer: str = "",
                    col_widths: list | None = None, highlight_rows: list | None = None) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.5)
    height = Inches(5.2)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = WHITE

    highlight_rows = highlight_rows or []
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = GRAY_DK

    add_footer(slide, prs, footer)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Capa
    add_title_slide(prs)

    # 2. Agenda
    add_bullets(
        prs,
        "Agenda da reuniao",
        [
            ("1.", "Status das 8 decisoes da reuniao anterior (08/jun)"),
            ("2.", "Corpus consolidado: 101 fichas em 11 tracks tematicos"),
            ("3.", "Pente fino critico do corpus (classificacao por impacto)"),
            ("4.", "Resposta a sugestao do orientador: CLIP (Contrastive Language-Image Pre-training)"),
            ("5.", "FiLM em fairness: gap na literatura como contribuicao original"),
            ("6.", "Proposta: Contribuicao C7 - ablation arquitetural"),
            ("7.", "Roadmap LaTeX 30 dias ate qualificacao"),
            ("8.", "Discussao e proximos passos"),
        ],
        footer="Tempo estimado: 12-15 min de apresentacao + Q&A",
    )

    # SECAO 1 - Status decisoes
    add_section_divider(prs, "1", "Status das 8 decisoes", "Reuniao anterior - 08/jun/2026")

    add_table_slide(
        prs,
        "Decisoes da reuniao anterior - status",
        ["#", "Decisao", "Status", "Observacoes"],
        [
            ["1", "Pipeline 6 etapas aprovado", "OK", "Validado e detalhado por etapa"],
            ["2", "SkinToneNet pre-treinado", "OK", "Leitura integral concluida"],
            ["3", "Corpus >=100 artigos", "OK 101", "+46 fichas na ultima rodada"],
            ["4", ">=20 artigos 2025-2026", "OK 25", "Meta superada"],
            ["5", "Adicionar CLIP conditioning", "OK Track I", "14 fichas + analise tecnica"],
            ["6", "Submissao qualificacao 15/jul/2026", "Em curso", "LaTeX em andamento"],
        ],
        col_widths=[0.5, 4.8, 1.8, 5.4],
        highlight_rows=[5],
        footer="4 itens fechados + 1 com track dedicada + 1 em execucao conforme planejado",
    )

    # SECAO 2 - Corpus 101 fichas
    add_section_divider(prs, "2", "Corpus consolidado", "101 fichas / 11 tracks tematicos")

    add_bullets(
        prs,
        "Marco atingido: 101 fichas em 11 tracks",
        [
            ("Total:", "101 fichas distribuidas em 11 tracks tematicos"),
            ("Meta 2025-2026:", "25 fichas (meta era >=20) - ATINGIDA"),
            ("Crescimento:", "+46 fichas em 3 levas da Rodada 7 (55 -> 101)"),
            ("Novos tracks:", "I, J, K criados em resposta a recomendacoes"),
            ("Track I:", "VLM/CLIP fairness (14 fichas) - resposta direta ao orientador"),
            ("Track J:", "Conditioning moderno LoRA/ViT (5 fichas)"),
            ("Track K:", "FR fundadores ArcFace/FaceNet/AdaFace (6 fichas) - gap fix"),
        ],
        footer="",
    )

    add_table_slide(
        prs,
        "Justificativa de cada track (por que existe)",
        ["Track", "Tema", "Fichas", "Papel na tese"],
        [
            ["A", "Race classification 7-class", "4", "Tarefa central"],
            ["B", "FR fairness", "16", "Paradigma dominante"],
            ["C", "Skin tone (MST)", "7", "Sinal auxiliar do pipeline"],
            ["D", "Mitigacao algoritmica", "10", "Baselines a superar"],
            ["E", "Auditoria & surveys", "14", "Infraestrutura intelectual"],
            ["F", "Fundamentacao etica", "3", "Posicionamento da tese"],
            ["G", "Mecanismos ML paradigmaticos", "9", "Esqueleto teorico (FiLM, LAFTR)"],
            ["I (R7)", "VLM/CLIP fairness", "14", "Resposta ao orientador"],
            ["J (R7)", "Conditioning moderno (LoRA, ViT)", "5", "Alternativas ao FiLM"],
            ["K (R7)", "FR fundadores (ArcFace, FaceNet)", "6", "Preenche gap obvio"],
            ["L", "Auxiliar / complementar", "13", "Direcoes adjacentes"],
        ],
        col_widths=[1.0, 4.0, 0.8, 6.7],
        highlight_rows=[7, 8, 9],
        footer="Tracks I, J, K criados na Rodada 7 - destacados para discussao",
    )

    # SECAO 3 - Pente fino
    add_section_divider(prs, "3", "Pente fino", "Classificacao por impacto na tese")

    add_table_slide(
        prs,
        "Classificacao critica das 101 fichas",
        ["Categoria", "Fichas", "%", "Implicacao para a tese"],
        [
            ["Forte favoravel", "12", "11.9%", "Fundamentam etapas/contribuicoes diretamente"],
            ["Favoravel", "38", "37.6%", "Suporte / alinhadas ao argumento"],
            ["Neutra / contextual", "18", "17.8%", "Background, nao influenciam decisoes"],
            ["Caminho alternativo", "26", "25.7%", "Cobertos via Contribuicao C7 (ablation)"],
            ["Conflito moderado", "5", "5.0%", "Endereçaveis - resposta defensiva mapeada"],
            ["Conflito forte", "2", "2.0%", "Pangelinan (via H6), Neto (limitacao reconhecida)"],
        ],
        col_widths=[3.2, 1.0, 1.0, 7.3],
        highlight_rows=[0, 5],
        footer="Apenas 2 conflitos fortes - ambos com resposta defensiva pronta. Nenhuma refutacao categorica sem resposta.",
    )

    add_bullets(
        prs,
        "Conclusao do pente fino",
        [
            ("Resumo:", "68 fichas (67.3%) alinhadas ou neutras / 26 caminho alternativo / 7 conflitos"),
            ("Conflito forte 1:", "Pangelinan 2023 - 'pixel info' explica disparidades -> endereçado via Hipotese H6"),
            ("Conflito forte 2:", "Neto 2025 - balanceamento nao basta -> reconhecido como limitacao + parte do framing"),
            ("Caminho alternativo:", "26 fichas (VLM, CLIP) cobertos via C7 - ablation arquitetural"),
            ("Pontos fortes:", "12 fichas forte favoravel fundamentam etapas 1-5 do pipeline"),
            ("", ""),
            ("Decisao:", "Corpus esta solido para sustentar a qualificacao"),
        ],
        footer="",
    )

    # SECAO 4 - CLIP
    add_section_divider(prs, "4", "Resposta ao orientador",
                        "CLIP = Contrastive Language-Image Pre-training (Radford et al., ICML 2021)")

    add_bullets(
        prs,
        "O que e o CLIP - confirmando referencia",
        [
            ("Paper original:", "Radford et al. - 'Learning Transferable Visual Models From Natural Language Supervision' ICML 2021"),
            ("Mecanismo:", "Dois encoders (vision + text) treinados via contrastive loss em 400M pares (imagem, texto)"),
            ("Resultado:", "Espaco de embedding compartilhado - imagens e textos relacionados ficam proximos"),
            ("Capacidade central:", "Zero-shot classification - basta descricao textual, sem fine-tuning"),
            ("Adocao em fairness:", "CLIP virou espinha dorsal da pesquisa em vision-language fairness desde 2022"),
            ("", ""),
            ("No contexto da tese:", "CLIP produz EMBEDDINGS; FiLM e MECANISMO de injecao - sao complementares, nao substitutos"),
        ],
        footer="Radford et al. 2021 - arxiv.org/abs/2103.00020 - >10k citacoes",
    )

    add_table_slide(
        prs,
        "Estudos de CLIP em fairness no corpus (Track I)",
        ["Paper", "Venue", "Abordagem", "Aplicacao na tese"],
        [
            ["FairCLIP (Luo 2024)", "CVPR 2024", "Optimal transport / Sinkhorn", "Baseline principal C7"],
            ["FairerCLIP (Dehdashtian 2024)", "ICLR 2024", "Debias zero-shot via RKHS", "Baseline leve C7"],
            ["BendVLM (2024)", "arXiv", "Test-time debiasing", "Alternativa sem retraining"],
            ["LoRA-FAIR (Bian 2025)", "ICCV 2025", "LoRA + fairness aggregation", "Track J - C7 ablation"],
            ["FaceScanPaliGemma (AlDahoul 2024)", "Nature SR 2026", "VLM fine-tuned FairFace", "Baseline SOTA 7-class 75.7%"],
            ["FairViT (Tian 2024)", "ECCV 2024", "Adaptive attention masking", "Backbone alternativo"],
            ["Closed-form debias (2026)", "arXiv", "Analitico", "Alternativa eficiente"],
            ["Bias subspace VLM (2025)", "arXiv", "Projecao subspace", "Mecanismo alternativo"],
            ["Survey multimodal fairness (2024)", "arXiv", "Survey", "Contexto Cap 2"],
        ],
        col_widths=[3.3, 1.8, 3.2, 4.2],
        highlight_rows=[0, 1, 4],
        footer="14 fichas no total na Track I + paper original CLIP (Radford 2021).",
    )

    # SECAO 5 - FiLM
    add_section_divider(prs, "5", "FiLM em fairness",
                        "Achado critico: gap na literatura = contribuicao original")

    add_bullets(
        prs,
        "Estudos de FiLM em fairness",
        [
            ("Paper FiLM original:", "Perez et al. 2018 (AAAI) - VERIFIED no corpus"),
            ("Aplicacao em fairness pelos autores:", "ZERO - paper nao aborda fairness"),
            ("Adocao em fairness na literatura:", "RARA - busca extensiva confirma escassez"),
            ("Mecanismos analogos no corpus:", ""),
            ("  - LAFTR (Madras 2018)", "adversarial conditioning (nao FiLM)"),
            ("  - FairViT (Tian 2024)", "adaptive masking em attention (ViT)"),
            ("  - LoRA-FAIR (Bian 2025)", "low-rank weight modulation"),
            ("", ""),
            ("Implicacao:", "FiLM aplicado a fairness em race classification e DIRECAO ORIGINAL"),
        ],
        footer="Perez et al. 2018 (AAAI) - paper original nao discute fairness nem na Future Work.",
    )

    add_bullets(
        prs,
        "FiLM vs CLIP - esclarecimento conceitual",
        [
            ("Confusao categorial:", "FiLM e CLIP NAO sao objetos comparaveis no mesmo nivel"),
            ("FiLM (Perez 2018):", "MECANISMO / camada que modula features via gamma*x + beta"),
            ("CLIP (Radford 2021):", "MODELO que produz embeddings 512-dim de texto e imagem"),
            ("Relacao real:", "CLIP fornece o SINAL; FiLM (ou cross-attn) e UMA forma de injetar esse sinal"),
            ("Cronologia:", "FaceNet 2015, ArcFace 2019, FiLM 2018, CLIP 2021, LoRA 2021"),
            ("Status 2024-2026:", "TODOS continuam padrao na literatura - idade nao implica obsolescencia"),
            ("", ""),
            ("Combinacao valida:", "Stable Diffusion = CLIP embeddings + cross-attention. Analogo: CLIP + FiLM e possivel"),
        ],
        footer="",
    )

    # SECAO 6 - Proposta C7
    add_section_divider(prs, "6", "Proposta",
                        "Contribuicao C7 - ablation arquitetural no Cap 2")

    add_table_slide(
        prs,
        "C7 - 5 configuracoes comparadas",
        ["Config", "Arquitetura", "Mecanismo de conditioning", "Origem"],
        [
            ["A", "ConvNeXt-T baseline", "Sem conditioning", "Controle"],
            ["B (proposta)", "ConvNeXt-T + FiLM", "MST 10-dim -> gamma, beta", "Nossa contribuicao central"],
            ["C", "ConvNeXt-T + FiLM", "CLIP-text embedding -> gamma, beta", "Hibrido CLIP+FiLM"],
            ["D", "ConvNeXt-T + cross-attn", "CLIP-text embedding via attention", "FairCLIP-style"],
            ["E", "ConvNeXt-T + LoRA-FAIR", "Low-rank weight modulation", "LoRA-FAIR (Bian 2025)"],
        ],
        col_widths=[1.0, 3.0, 4.5, 4.5],
        highlight_rows=[1],
        footer="Se CLIP superar FiLM -> reportamos como informacao util (nao invalidacao). Honra recomendacao do orientador.",
    )

    add_bullets(
        prs,
        "Justificativa para manter FiLM como mecanismo central",
        [
            ("Adequacao a baixa-dim:", "MST e vetor 10-dim - FiLM e o sweet spot"),
            ("Custo parametrico:", "~1% sobre backbone (vs ~3% cross-attention equivalente)"),
            ("Interpretabilidade:", "gamma e beta sao por canal e visualizaveis"),
            ("Compatibilidade:", "Substitui parametros do BatchNorm - 'conditional BatchNorm'"),
            ("Compatibilidade C7:", "Ablation limpa - liga/desliga FiLM mantendo backbone fixo"),
            ("", ""),
            ("Argumento de originalidade:", "FiLM em fairness e direcao nao explorada - vira contribuicao"),
        ],
        footer="Perez et al. 2018 (AAAI) + Schumann 2023 (MST) + Pereira 2026 (SkinToneNet) ancoragem direta",
    )

    # SECAO 7 - Roadmap
    add_section_divider(prs, "7", "Roadmap LaTeX", "30 dias ate 15/jul/2026")

    add_table_slide(
        prs,
        "Cronograma de execucao - 4 semanas",
        ["Semana", "Periodo", "Entregavel", "Observacoes"],
        [
            ["1", "15-21/jun (esta)", "Setup Overleaf + Cap 1 (Introducao)", "Narrativa pre-qualificacao consolidada"],
            ["2", "22-28/jun", "Cap 2 (Revisao - 11 tracks)", "Pente fino do corpus + 101 fichas"],
            ["3", "29/jun-05/jul", "Cap 3 (Objetivos) + Cap 4 (Metodologia)", "Pipeline 6 etapas validado"],
            ["4", "06-12/jul", "Cap 5 (Cronograma) + revisao final", "Integracao + ABNT"],
            ["Buffer", "13-15/jul", "Ajustes finais + submissao", "-"],
        ],
        col_widths=[1.0, 2.4, 5.0, 4.6],
        highlight_rows=[0],
        footer="",
    )

    # SECAO 8 - Discussao
    add_section_divider(prs, "8", "Discussao", "Pontos para alinhamento hoje")

    add_bullets(
        prs,
        "Perguntas ao orientador",
        [
            ("1.", "Concorda com manter FiLM central + CLIP-conditioning como C7 (ablation)?"),
            ("2.", "11 tracks bem cobertos? Algum precisa de mais fichas?"),
            ("3.", "Esclarecimento conceitual CLIP (modelo) vs FiLM (mecanismo) faz sentido?"),
            ("4.", "Tem template Overleaf Unifesp/ICT ou usamos padrao?"),
            ("5.", "Co-orientador - ja podemos formalizar?"),
            ("6.", "Sugestoes para banca preliminar?"),
        ],
        footer="",
    )

    add_bullets(
        prs,
        "Proximos passos imediatos (proxima semana)",
        [
            ("1.", "Setup Overleaf com template institucional"),
            ("2.", "Gerar .bib consolidado a partir das 101 fichas"),
            ("3.", "Escrever Cap 1 - Introducao usando narrativa pre-qualificacao"),
            ("4.", "Importar imagens ja produzidas (FiLM pipeline, ConvNeXt-T, metricas Hardt)"),
            ("5.", "Promover fichas OVERVIEW_ONLY para VERIFIED nas 12 forte favoravel"),
            ("", ""),
            ("Entrega proxima reuniao:", "Cap 1 escrito em LaTeX + bibliografia consolidada"),
        ],
        footer="30 dias contados ate 15/jul/2026 - margem confortavel sem folga para retrabalho",
    )

    # Slide final
    add_section_divider(prs, "", "Obrigado",
                        "Discussao livre - duvidas e ajustes antes do LaTeX")

    return prs


def main() -> None:
    here = Path(__file__).parent
    out = here / "material_reuniao_orientador_2026-06-15.pptx"
    prs = build_presentation()
    prs.save(out)
    print(f"Apresentacao gerada: {out}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
