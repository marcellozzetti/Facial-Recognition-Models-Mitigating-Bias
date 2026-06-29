"""Gera apresentacao PowerPoint da reuniao 2026-06-28.

Uso:
    python _gerar_apresentacao_2026-06-28.py
    -> produz: docs/ativo/material_reuniao_orientador_2026-06-28.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Paleta academica sobria
NAVY = RGBColor(0x1F, 0x2A, 0x4E)
GRAY_DK = RGBColor(0x3D, 0x42, 0x4E)
GRAY_MD = RGBColor(0x70, 0x76, 0x82)
GRAY_LT = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(2.5), Inches(1.3), Inches(10.5), Inches(2.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Evolução das últimas 2 semanas"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Revisão completa, cross-reference e validação externa via NotebookLM"
    p2.font.size = Pt(20)
    p2.font.color.rgb = GRAY_DK

    p3 = tf.add_paragraph()
    p3.text = "Corpus 99 % verificado · Tese fundamentada · Pronta para escrita"
    p3.font.size = Pt(15)
    p3.font.color.rgb = GRAY_MD
    p3.font.italic = True

    meta = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(10.5), Inches(2.5))
    mf = meta.text_frame
    mf.word_wrap = True
    rows = [
        ("Mestrando:", "Marcello Ozzetti"),
        ("Orientador:", "Prof. Marcos Quiles"),
        ("Programa:", "Mestrado em Ciência da Computação — Unifesp / ICT"),
        ("Reunião:", "28 de junho de 2026"),
    ]
    for i, (k, v) in enumerate(rows):
        p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        p.text = f"{k}  {v}"
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY_DK


def add_section_divider(prs: Presentation, num: str, title: str, subtitle: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tx_num = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3), Inches(3))
    pn = tx_num.text_frame.paragraphs[0]
    pn.text = num
    pn.font.size = Pt(140)
    pn.font.bold = True
    pn.font.color.rgb = WHITE

    tx_t = slide.shapes.add_textbox(Inches(4.3), Inches(2.5), Inches(8.5), Inches(2))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = GRAY_LT


def add_title(slide, text: str) -> None:
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, prs, text: str) -> None:
    if not text:
        return
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
    pf = ft.text_frame.paragraphs[0]
    pf.text = text
    pf.font.size = Pt(10)
    pf.font.color.rgb = GRAY_MD
    pf.font.italic = True


def add_bullets(prs: Presentation, title: str, bullets: list, footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(5.4))
    tf = tx.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            head, body = item
            run1 = p.add_run()
            run1.text = head + "  "
            run1.font.size = Pt(16)
            run1.font.bold = True
            run1.font.color.rgb = NAVY
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(16)
            run2.font.color.rgb = GRAY_DK
        else:
            p.text = "— " + item
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_DK
        p.space_after = Pt(8)

    add_footer(slide, prs, footer)


def add_table_slide(prs: Presentation, title: str, headers: list, rows: list, footer: str = "",
                    col_widths: list | None = None, highlight_rows: list | None = None) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.5)
    height = Inches(5.2)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = WHITE

    highlight_rows = highlight_rows or []
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = GRAY_DK

    add_footer(slide, prs, footer)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Capa
    add_title_slide(prs)

    # 2. Agenda
    add_bullets(
        prs,
        "Agenda da reunião",
        [
            ("1.", "Status das decisões da reunião anterior (15/jun)"),
            ("2.", "Revisão bibliográfica concluída — 103 de 104 fichas verificadas"),
            ("3.", "Cross-reference sistemático tese × 104 fichas"),
            ("4.", "Validação externa via NotebookLM (Google AI Plus) — segunda opinião"),
            ("5.", "Rodada 8 — fundamentação Latinx (Telles, Bryc, Pew)"),
            ("6.", "Estudo comparativo Cap 2 — 4 configurações de conditioning"),
            ("7.", "Decisões a alinhar + próximos passos da escrita"),
        ],
        footer="",
    )

    # SEÇÃO 1 — Status decisões anteriores
    add_section_divider(prs, "1", "Status das decisões anteriores", "Reunião 15/jun/2026 — 4 decisões registradas")

    add_table_slide(
        prs,
        "Status das 4 decisões da reunião anterior",
        ["#", "Decisão", "Status", "Observações"],
        [
            ["1", "Embasamento aprovado — escrita liberada", "Concluído", "Corpus expandido para 104 fichas"],
            ["2", "Deadline 15/jul/2026 (primeira revisão)", "Em curso", "17 dias restantes"],
            ["3", "ResNet+FiLM tradicional + CLIP avaliação", "Formalizado", "4 configurações no Cap 2 definidas"],
            ["4", "Mais uma passada no corpus", "Concluído", "Camadas 2-5 completadas + auditoria"],
        ],
        col_widths=[0.5, 5.0, 1.8, 5.2],
        footer="Resultado: 4 de 4 decisões executadas conforme planejado.",
    )

    # SEÇÃO 2 — Revisão bibliográfica concluída
    add_section_divider(prs, "2", "Revisão bibliográfica concluída", "Corpus de 104 fichas com auditoria sistemática")

    add_table_slide(
        prs,
        "Evolução do corpus em 2 semanas",
        ["Métrica", "15/jun", "28/jun", "Variação"],
        [
            ["Fichas no corpus", "101", "104", "+3 (Rodada 8 Latinx)"],
            ["Fichas VERIFIED", "14 (Camada 1)", "103 (99 %)", "+89 fichas verificadas"],
            ["PDFs no repositório", "26", "103", "+77 PDFs"],
            ["Autoria correta verificada", "43 %", "100 %", "100 % das fichas"],
            ["Bibliografia consolidada", "—", "104 entradas", "Pronta para Overleaf"],
            ["Auditoria de qualidade", "—", "29 A / 14 B / 57 C / 4 D", "Sistemática"],
        ],
        col_widths=[3.5, 2.0, 3.0, 4.0],
        highlight_rows=[1, 4],
        footer="Detalhamento em _auditoria_fichas_relatorio.md.",
    )

    # SEÇÃO 3 — Cross-reference
    add_section_divider(prs, "3", "Cross-reference sistemático",
                        "Auditoria tese × 104 fichas — todos os elementos cruzados")

    add_bullets(
        prs,
        "Cross-reference: o que foi auditado",
        [
            ("Objetivo geral:", "fundamentado por mais de 12 fichas centrais sem conflito não endereçado"),
            ("Pipeline de 6 etapas:", "cada etapa cruzada com fichas que fundamentam"),
            ("7 contribuições (C1-C7):", "originalidade verificada para cada uma"),
            ("6 hipóteses (H1-H6):", "fichas que sustentam + fichas em conflito + status"),
            ("Storytelling em 6 partes:", "mapeamento ficha-por-ficha de cada seção"),
            ("Decisão arquitetural FiLM:", "justificada vs 8 alternativas avaliadas"),
            ("", ""),
            ("Veredito:", "tese fundamentada, sem conflitos não endereçados, pronta para escrita"),
        ],
        footer="Documento: _validacao_cross_reference_v3.md",
    )

    # SEÇÃO 4 — NotebookLM
    add_section_divider(prs, "4", "Validação externa via NotebookLM",
                        "Segunda opinião independente — Google AI Plus")

    add_bullets(
        prs,
        "Por que validar com NotebookLM",
        [
            ("Independência:", "segunda opinião externa, não enviesada pela construção interna do argumento"),
            ("Operacionalização:", "corpus organizado em 4 tiers de relevância para importação seletiva"),
            ("Auditabilidade:", "cada resposta cita as fontes específicas do corpus"),
            ("", ""),
            ("Cobertura da análise:", "7 perguntas-chave sobre originalidade, valor científico, divergências, gaps, recência, defensibilidade do pipeline"),
            ("Resultado:", "convergente com a análise interna em 6 dimensões + 2 sugestões novas valiosas incorporadas"),
        ],
        footer="Documento de tiers: _tiers_relevancia_pdfs.md",
    )

    add_table_slide(
        prs,
        "Convergência: análise interna × NotebookLM",
        ["Asserção", "Análise interna", "NotebookLM"],
        [
            ["FiLM em fairness é inédito", "Sim", "Sim"],
            ["Matriz MST × FairFace preenche gap real", "Sim", "Sim"],
            ["Pangelinan 2023 é refutação central", "Sim", "Sim"],
            ["Pipeline tem sequenciamento lógico defensável", "Sim", "Sim"],
            ["Decomposição erro Latinx é maior diferencial", "Sim", "Sim"],
            ["Corpus na fronteira (papers 2026)", "Sim", "Sim"],
        ],
        col_widths=[7.0, 2.75, 2.75],
        footer="6 de 6 asserções convergentes — alta confiabilidade da análise interna.",
    )

    add_bullets(
        prs,
        "Veredito do NotebookLM",
        [
            ("", '"O trabalho está muito bem estruturado, possui contribuições claras'),
            ("", "e está fundamentado em literatura de ponta. O ponto crítico da"),
            ("", "defesa será a decomposição do erro Latinx, que é onde reside a sua"),
            ("", 'maior contribuição intelectual e diagnóstica."'),
            ("", ""),
            ("Fonte:", "análise paralela conduzida em NotebookLM (Google AI Plus, 2026-06-16) sobre as 100 fichas do corpus"),
        ],
        footer="",
    )

    add_table_slide(
        prs,
        "2 sugestões novas — todas incorporadas",
        ["Sugestão do NotebookLM", "Ação tomada", "Onde está documentada"],
        [
            ["Heterogeneidade intra-Latinx — categoria tratada como bloco monolítico", "Rodada 8 com 3 fichas integradas (Telles, Bryc, Pew)", "Seção 5 desta apresentação"],
            ["Sensitivity analysis SkinToneNet — risco de propagação de viés", "OE-2 expandido — validação com 2 a 3 classificadores MST alternativos", "Objetivo da tese v3.5"],
        ],
        col_widths=[5.0, 4.5, 3.0],
        footer="",
    )

    add_bullets(
        prs,
        "Furos identificados pelo NotebookLM — respondidos",
        [
            ("1) Dependência do SkinToneNet:", "mitigado por sensitivity analysis no OE-2 (validar conclusões com 2 a 3 classificadores MST diferentes)"),
            ("2) Escalabilidade industrial:", "FairFace 108 mil imagens vs NIST FRVT 18 milhões — limite reconhecido formalmente"),
            ("Posicionamento:", "esta dissertação é demonstração metodológica de viabilidade (mestrado), não recomendação de deploy comercial"),
            ("Trabalho futuro:", "replicação em escala industrial declarada como direção subsequente"),
        ],
        footer="Reconhecimento formal nas Seções 9 e 10 do _objetivo_tese v3.5.",
    )

    # SEÇÃO 5 — Rodada 8 Latinx
    add_section_divider(prs, "5", "Rodada 8 — fundamentação Latinx",
                        "Tripé empírico antropologia + genética + sociologia")

    add_bullets(
        prs,
        "Motivação da Rodada 8",
        [
            ("Problema observado:", "F1 Latinx aproximadamente 60 % persistente em modelos do estado da arte (AlDahoul 2024)"),
            ("Gap no corpus original:", "categoria Latinx tratada como bloco monolítico sem fundamentação empírica explícita de heterogeneidade"),
            ("Resposta:", "Rodada 8 enxuta — 3 fichas integradas para sustentar empíricamente a hipótese H3 e a contribuição C6"),
            ("", ""),
            ("Aporte específico:", "permite defender o erro Latinx como tendo componente estrutural irredutível (overlap MST), não apenas algorítmico"),
        ],
        footer="",
    )

    add_table_slide(
        prs,
        "3 fichas integradas na Rodada 8",
        ["Ficha", "Autor / Venue", "Aporte ao argumento"],
        [
            ["telles_2014 (Pigmentocracies)", "Edward Telles + PERLA — UNC Press, 320 pp", "Antropologia — pigmentocracia em 4 países LatAm (Brasil, Colômbia, México, Peru)"],
            ["bryc_2015 (Genetic Ancestry)", "Bryc, Reich et al — AJHG (Harvard + 23andMe)", "Genética — 162 mil indivíduos, heterogeneidade Native + European + African"],
            ["pew_2017_hispanic_identity", "Lopez, Gonzalez-Barrera et al — Pew Research", "Sociologia — identidade Hispanic declina 97 % para 50 % cross 4 gerações nos EUA"],
        ],
        col_widths=[3.0, 4.5, 5.0],
        footer="Tripé empírico que sustenta H3 + C6 com 3 disciplinas independentes convergentes.",
    )

    # SEÇÃO 6 — 4 configurações Cap 2
    add_section_divider(prs, "6", "Estudo comparativo do Capítulo 2",
                        "4 configurações de conditioning — proposta + alternativa moderna")

    add_table_slide(
        prs,
        "4 configurações de conditioning comparadas",
        ["Config", "Arquitetura", "Sinal de conditioning", "Origem"],
        [
            ["A", "ConvNeXt-T baseline", "Sem conditioning", "Controle"],
            ["B", "ConvNeXt-T + FiLM", "MST 10-dim → γ, β", "Proposta principal (linha tradicional)"],
            ["C", "ConvNeXt-T + Gated FiLM", "MST 10-dim → γ, β não-lineares", "Ablação metodológica"],
            ["D", "ConvNeXt-T + FiLM", "Embedding CLIP-text 512-dim", "Avaliação alternativa (orientador)"],
        ],
        col_widths=[1.0, 3.5, 4.5, 4.0],
        highlight_rows=[1],
        footer="Atende recomendação da reunião 15/jun: linha tradicional + avaliação CLIP.",
    )

    # SEÇÃO 7 — Decisões + próximos passos
    add_section_divider(prs, "7", "Decisões a alinhar e próximos passos",
                        "Reta final para a primeira revisão em 15/jul/2026")

    add_bullets(
        prs,
        "Decisões a alinhar hoje",
        [
            ("1.", "Promover H6 para OE-6 formal — decomposição quantitativa pixel info × skin tone?"),
            ("2.", "Estudo comparativo com 4 configurações atende à recomendação anterior?"),
            ("3.", "Há template Overleaf institucional Unifesp / ICT ou seguimos com o padrão?"),
            ("4.", "Co-orientador — alguma definição?"),
            ("5.", "Sugestões para a banca preliminar?"),
        ],
        footer="",
    )

    add_table_slide(
        prs,
        "Próximas 3 semanas — plano de escrita",
        ["Semana", "Período", "Entregável"],
        [
            ["1 (esta)", "29/jun - 05/jul", "Setup Overleaf + Capítulo 1 (Introdução)"],
            ["2", "06 - 12/jul", "Capítulo 2 (Revisão) + Capítulo 3 (Objetivos)"],
            ["3", "13 - 14/jul", "Capítulo 4 (Metodologia) + Capítulo 5 (Cronograma) + revisão"],
            ["Target", "15/jul", "Entrega da primeira revisão ao orientador"],
        ],
        col_widths=[1.5, 2.5, 8.5],
        highlight_rows=[0, 3],
        footer="Plano enxuto — 17 dias para entregar 5 capítulos + revisão final.",
    )

    add_bullets(
        prs,
        "Próximos passos imediatos (esta semana)",
        [
            ("1.", "Setup Overleaf com template institucional"),
            ("2.", "Importar bibliografia consolidada — referencias.bib com 104 entradas"),
            ("3.", "Iniciar escrita do Capítulo 1 — Introdução"),
            ("4.", "Promover OE-6 se aprovado nesta reunião"),
            ("5.", "Continuar leitura aprofundada da Camada 2 em paralelo à escrita"),
            ("", ""),
            ("Entrega para próxima reunião:", "Capítulo 1 escrito em LaTeX + estrutura dos capítulos 2 a 5"),
        ],
        footer="",
    )

    # Slide final
    add_section_divider(prs, "", "Obrigado",
                        "Discussão livre — dúvidas e alinhamento")

    return prs


def main() -> None:
    here = Path(__file__).parent
    out = here / "material_reuniao_orientador_2026-06-28.pptx"
    prs = build_presentation()
    prs.save(out)
    print(f"Apresentacao gerada: {out}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
