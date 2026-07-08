"""Gera apresentacao PowerPoint da reuniao 2026-07-08.

Voz autoral, densa em conteudo tecnico, com espaco para debate.
- Plano de trabalho consolidado
- Estado da escrita capitulo a capitulo (com substancia)
- Postura autoral sobre uso de LLM
- Adequacao etica CEP + trade-offs discutiveis
- Extensao de prazo
- Perguntas de debate

Uso:
    python _gerar_apresentacao_2026-07-08.py
    -> produz: docs/ativo/material_reuniao_orientador_2026-07-08.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x2A, 0x4E)
GRAY_DK = RGBColor(0x3D, 0x42, 0x4E)
GRAY_MD = RGBColor(0x70, 0x76, 0x82)
GRAY_LT = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(2.5), Inches(1.3), Inches(10.5), Inches(2.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Reta final até 15 de julho"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Estado da escrita, decisões éticas e ajuste de cronograma"
    p2.font.size = Pt(20)
    p2.font.color.rgb = GRAY_DK

    p3 = tf.add_paragraph()
    p3.text = "Equidade racial em classificação facial via conditioning por tom de pele (MST) sobre FairFace"
    p3.font.size = Pt(14)
    p3.font.color.rgb = GRAY_MD
    p3.font.italic = True

    meta = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(10.5), Inches(2.5))
    mf = meta.text_frame
    mf.word_wrap = True
    rows = [
        ("Mestrando:", "Marcello Ozzetti"),
        ("Orientador:", "Prof. Marcos Quiles"),
        ("Programa:", "Mestrado em Ciência da Computação — Unifesp / ICT"),
        ("Reunião:", "8 de julho de 2026"),
    ]
    for i, (k, v) in enumerate(rows):
        p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        p.text = f"{k}  {v}"
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY_DK


def add_section_divider(prs: Presentation, num: str, title: str, subtitle: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tx_num = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3), Inches(3))
    pn = tx_num.text_frame.paragraphs[0]
    pn.text = num
    pn.font.size = Pt(140)
    pn.font.bold = True
    pn.font.color.rgb = WHITE

    tx_t = slide.shapes.add_textbox(Inches(4.3), Inches(2.5), Inches(8.5), Inches(2))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = GRAY_LT


def add_title(slide, text: str) -> None:
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, prs, text: str) -> None:
    if not text:
        return
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
    pf = ft.text_frame.paragraphs[0]
    pf.text = text
    pf.font.size = Pt(10)
    pf.font.color.rgb = GRAY_MD
    pf.font.italic = True


def add_bullets(prs: Presentation, title: str, bullets: list, footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(5.4))
    tf = tx.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            head, body = item
            run1 = p.add_run()
            run1.text = head + "  "
            run1.font.size = Pt(16)
            run1.font.bold = True
            run1.font.color.rgb = NAVY
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(16)
            run2.font.color.rgb = GRAY_DK
        else:
            p.text = "— " + item
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_DK
        p.space_after = Pt(8)

    add_footer(slide, prs, footer)


def add_quote_slide(prs: Presentation, title: str, quote: str, attribution: str = "", footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(0.08), Inches(3.5))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = NAVY
    quote_box.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(1.8), Inches(2.2), Inches(10.5), Inches(3.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = GRAY_DK
    p.space_after = Pt(20)

    if attribution:
        p2 = tf.add_paragraph()
        p2.text = "— " + attribution
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY_MD

    add_footer(slide, prs, footer)


def add_table_slide(prs: Presentation, title: str, headers: list, rows: list, footer: str = "",
                    col_widths: list | None = None, highlight_rows: list | None = None) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.5)
    height = Inches(5.2)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = WHITE

    highlight_rows = highlight_rows or []
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = GRAY_DK

    add_footer(slide, prs, footer)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Capa
    add_title_slide(prs)

    # 2. Objetivo da reunião
    add_bullets(
        prs,
        "Objetivo desta reunião",
        [
            ("", ""),
            ("1.", "Plano até a defesa"),
            ("", ""),
            ("2.", "Estado real da escrita"),
            ("", ""),
            ("3.", "Decisão de pesquisa ética"),
        ],
        footer="Solicitação pessoal e debate técnico fecham a reunião como pontos complementares.",
    )

    # 3. Recap do objetivo da tese
    add_bullets(
        prs,
        "Recapitulando o objetivo da tese",
        [
            ("Objetivo geral:",
             "desenvolver e avaliar um pipeline de classificação racial em imagens faciais que incorpora tom de pele (escala Monk Skin Tone) como sinal auxiliar condicionante via mecanismo arquitetural"),
            ("", ""),
            ("Problema empírico atacado:",
             "disparidade severa entre Black F1 ≈ 90 % e Latinx F1 ≈ 60 % em SOTA sobre FairFace (FaceScanPaliGemma, AlDahoul 2024)"),
            ("", ""),
            ("Contribuição principal esperada:",
             "primeira instância empírica documentada do uso de tom de pele explícito como contexto arquitetural para race classification multi-classe"),
            ("", ""),
            ("Diagnóstico central da tese:",
             "o erro Latinx tem componente fenotípico irredutível (heterogeneidade MST intra-categoria) + componente algorítmico — a decomposição desses dois é o principal diferencial"),
        ],
        footer="Base para tudo que vem a seguir — pipeline, quatro configurações, sete contribuições, seis hipóteses.",
    )

    # 4. Consolidado do corpus
    add_table_slide(
        prs,
        "Corpus consolidado — 104 fichas em 11 linhas de pesquisa",
        ["Linha de pesquisa", "N", "% corpus", "Papers-âncora"],
        [
            ["A — Race classification multi-classe", "4", "3,8 %", "Kärkkäinen 2021 (FairFace), AlDahoul 2024"],
            ["B — Face recognition fairness", "16", "15,4 %", "Wang RFW 2019, Robinson BFW 2020, Deng ArcFace"],
            ["C — Skin tone alternativo (Fitzpatrick, MST, ITA)", "7", "6,7 %", "Schumann 2023 (MST), Pereira 2026 (SkinToneNet)"],
            ["D — Mitigação algorítmica", "10", "9,6 %", "Madras LAFTR 2018, Sagawa Group DRO 2020"],
            ["E — Auditoria e Surveys", "14", "13,5 %", "Buolamwini 2018, Mehrabi 2021 survey"],
            ["F — Fundamentação científica e ética", "3", "2,9 %", "Fuentes AAPA 2019, Lewontin 1972"],
            ["G — Mecanismos ML paradigmáticos", "9", "8,7 %", "Perez FiLM 2018, Ba LayerNorm, He ResNet"],
            ["I — VLM / CLIP em fairness", "14", "13,5 %", "Radford CLIP 2021, Luo FairCLIP 2024"],
            ["J — Conditioning moderno", "5", "4,8 %", "Liu ConvNeXt 2022, Dhariwal ADM 2021"],
            ["K — Fundadores de FR (DeepFace, ArcFace, CosFace)", "6", "5,8 %", "Taigman DeepFace 2014, Wang CosFace 2018"],
            ["L — Auxiliar / complementar", "13", "12,5 %", "Karkkainen datasets, Hendrycks robustness"],
            ["R8 — Diversidade Latinx (antropo + genética + socio)", "3", "2,9 %", "Telles 2014, Bryc 2015, Pew 2017"],
        ],
        col_widths=[4.8, 0.7, 1.2, 5.8],
        footer="Corpus preparado ao longo de 8 rodadas de expansão. 99 % das fichas verified via leitura direta do PDF.",
    )

    add_table_slide(
        prs,
        "Densidade do corpus — distribuição temporal",
        ["Período", "N fichas", "% corpus", "Papel na tese"],
        [
            ["Pré-2015 (fundações históricas)", "5", "4,8 %", "Lewontin 1972, Taigman DeepFace 2014, Telles 2014 — fundamentam OG"],
            ["2015-2017 (fairness formal)", "6", "5,8 %", "Hardt 2016 (Equal Opportunity), Kleinberg 2017 (impossibility)"],
            ["2018-2020 (marcos canônicos)", "14", "13,5 %", "Perez FiLM 2018, Buolamwini 2018, Madras LAFTR 2018"],
            ["2021-2022 (state of the art anterior)", "17", "16,3 %", "Radford CLIP 2021, Kärkkäinen FairFace 2021, Liu ConvNeXt 2022"],
            ["2023 (contexto imediato)", "5", "4,8 %", "Schumann MST-E, Pangelinan (refutação central)"],
            ["2024-2026 (fronteira)", "55", "52,9 %", "AlDahoul 2024, Luo FairCLIP 2024, Pereira SkinToneNet 2026"],
        ],
        col_widths=[4.0, 1.2, 1.3, 6.0],
        highlight_rows=[5],
        footer="Mais de metade do corpus é 2024+. Cinco papers 2026 sustentam o argumento de trabalho na fronteira.",
    )

    # SEÇÃO 1 — Plano até a defesa
    add_section_divider(prs, "1", "Plano até a defesa",
                        "Três marcos formais, novo horizonte de outubro")

    # SEÇÃO 1 — Plano até a defesa
    add_section_divider(prs, "1", "Plano até a defesa",
                        "Três marcos formais, novo horizonte de outubro")

    add_table_slide(
        prs,
        "Cronograma consolidado",
        ["#", "Marco", "Data", "Comentário"],
        [
            ["1", "Primeira revisão da qualificação ao orientador", "15/jul/2026", "Sete dias a partir de hoje — em curso"],
            ["2", "Pedido formal de qualificação ao PPG-CC / ICT", "30/jul/2026", "Prazo regimental do Programa — mantido"],
            ["3", "Defesa da qualificação", "outubro/2026", "Ajuste de agosto → outubro (extensão de 2 meses solicitada)"],
        ],
        col_widths=[0.5, 5.5, 2.0, 5.0],
        highlight_rows=[2],
        footer="Os dois primeiros marcos ficam intactos. A extensão só empurra o marco final.",
    )

    # SEÇÃO 2 — Estado da escrita
    add_section_divider(prs, "2", "Estado da escrita",
                        "O que já está escrito e como vai virar Overleaf")

    add_bullets(
        prs,
        "O que cada capítulo já tem",
        [
            ("Capítulo 1 — Introdução",
             "argumento central: F1 macro alto mascara disparidade Black 90 % × Latinx 60 % em SOTA sobre FairFace"),
            ("Capítulo 2 — Revisão",
             "fairness (Buolamwini, Hardt, Kleinberg), MST (Schumann, Pereira), FiLM (Perez), 104 papers verificados"),
            ("Capítulo 3 — Objetivos",
             "objetivo geral, seis específicos, seis hipóteses testáveis, sete contribuições esperadas"),
            ("Capítulo 4 — Metodologia",
             "pipeline em seis etapas, quatro configurações comparativas, adequação ética já resolvida"),
            ("Capítulo 5 — Cronograma",
             "único capítulo que só fecha depois desta reunião (marcos ajustados hoje)"),
            ("", ""),
            ("Formato:",
             "corpo escrito em Markdown estruturado; próximos sete dias faço a transposição para LaTeX/Overleaf"),
        ],
        footer="Se quiser ver, abro qualquer capítulo agora — está tudo no GitHub.",
    )

    add_bullets(
        prs,
        "A pergunta central da tese, em uma frase",
        [
            ("", ""),
            ("Pergunta de pesquisa:",
             "condicionar arquiteturalmente um classificador racial pelo tom de pele (MST) reduz a disparidade Black × Latinx sem sacrificar acurácia agregada?"),
            ("", ""),
            ("Por que MST em vez de Fitzpatrick:",
             "Fitzpatrick foi projetada para dermatologia (queimadura solar), tem apenas 6 classes e viés fototípico documentado. MST tem 10 classes e foi desenhada para representatividade em CV — Google 2022, Schumann 2023."),
            ("", ""),
            ("Por que FiLM em vez de concatenação simples:",
             "FiLM (Perez, AAAI 2018) modula features por escala e deslocamento condicionados — permite que o mesmo backbone reaja diferente ao mesmo pixel dependendo do contexto MST. Concatenação exige que o modelo re-aprenda a semântica do sinal."),
        ],
        footer="Slide-âncora do Cap 1. Toda a argumentação da tese sai daqui.",
    )

    add_table_slide(
        prs,
        "Pipeline proposto — seis etapas encadeadas",
        ["Etapa", "O que faz", "Entrega / contribuição"],
        [
            ["1. Classificador MST", "SkinToneNet pré-treinado + validação interna em subset FairFace estratificado", "Insumo do pipeline (não é contribuição)"],
            ["2. Auditoria FairFace", "Aplicar SkinToneNet sobre todo FairFace validation set", "Matriz pública MST × race classes — Contribuição C2"],
            ["3. Race classifier com conditioning", "ConvNeXt-T fine-tuned em FairFace, camadas FiLM por estágio recebendo vetor MST", "Configurações A/B/C/D — Contribuição C3"],
            ["4. Comparação contra baselines", "6 baselines: ResNet-34, ConvNeXt-T puro, FSCL+, Group DRO, FineFACE, Adversarial", "Triangulação DR + worst-class F1 + EO_h — Contribuição C4"],
            ["5. Fair transferência", "Aplicar backbone fair em face recognition downstream (RFW ou BFW)", "Demonstração de fair transfer — Contribuição C5"],
            ["6. Síntese decompositiva", "Combinar C2 + C5 para separar componente fenotípico do algorítmico do erro Latinx", "Decomposição pixel info × skin tone — Contribuição C6/OE-6"],
        ],
        col_widths=[2.8, 5.5, 4.2],
        highlight_rows=[2, 5],
        footer="Etapa 3 é onde o mecanismo FiLM entra. Etapa 6 é o diagnóstico central da tese.",
    )

    add_table_slide(
        prs,
        "Quatro configurações do estudo comparativo (Etapa 3)",
        ["Config", "Backbone + conditioning", "Sinal", "Papel no estudo"],
        [
            ["A", "ConvNeXt-T sem conditioning", "—", "Baseline de controle"],
            ["B", "ConvNeXt-T + FiLM linear", "MST 10-dim", "Proposta principal — hipótese central"],
            ["C", "ConvNeXt-T + Gated FiLM", "MST 10-dim (não-linear)", "Ablação: linear vs não-linear importa?"],
            ["D", "ConvNeXt-T + FiLM", "Embedding CLIP-text 512-dim", "Recomendação da reunião 15/jun — CLIP como alternativa moderna"],
        ],
        col_widths=[1.0, 3.8, 3.5, 5.0],
        highlight_rows=[1],
        footer="B é a proposta central. C isola a não-linearidade. D compara com sinal semântico rico via CLIP.",
    )

    add_bullets(
        prs,
        "O achado que ancora a tese (Rodada 8, junho)",
        [
            ("Problema conhecido:",
             "modelos SOTA reportam F1 Latinx próximo de 60 % há três anos — persiste mesmo em modelos multimodais recentes (AlDahoul 2024)"),
            ("", ""),
            ("Diagnóstico novo:",
             "não é só bias algorítmico. É heterogeneidade fenotípica intra-categoria que a rotulagem monolítica de FairFace apaga."),
            ("", ""),
            ("Fundamentação empírica agora tripla:"),
            ("Antropologia:",
             "Telles 2014 (PERLA, 4 países latino-americanos) documenta pigmentocracia — tom de pele varia mais dentro da categoria do que entre elas"),
            ("Genética:",
             "Bryc 2015 (AJHG, 162 mil indivíduos) — ancestralidade Native + European + African altamente variável em Latinos"),
            ("Sociologia:",
             "Pew 2017 — identidade Hispanic cai de 97 % para 50 % em quatro gerações nos EUA. Categoria é sociopolítica, não fenotípica."),
        ],
        footer="Esse é o argumento pelo qual C6 (decomposição fenotípico × algorítmico) vira contribuição central.",
    )

    add_quote_slide(
        prs,
        "Postura sobre uso de LLM — como quero deixar explícito na defesa",
        "A LLM me poupou tempo em bibliotecas de referência, formatação de fichas e "
        "revisão de coerência entre capítulos. Mas cada afirmação central da tese — "
        "a pergunta de pesquisa, a escolha do backbone, o mecanismo FiLM, o ajuste "
        "ético do OE-1, o argumento de heterogeneidade intra-Latinx — nasceu de "
        "leitura minha e de conversas nesta sala. A transposição para o Overleaf é "
        "manual, linha a linha, exatamente para forçar essa apropriação.",
        attribution="posição que vou defender se questionado na banca",
        footer="Preferi trazer isso já, para o senhor calibrar comigo o quanto explicitar em texto vs quanto só demonstrar em resposta.",
    )

    # SEÇÃO 3 — Decisão ética
    add_section_divider(prs, "3", "Decisão ética que tomei sozinho",
                        "Preciso validar o enquadramento no Art. 8º")

    add_bullets(
        prs,
        "O contexto",
        [
            ("O que investiguei:",
             "Resolução 200/2021 do Conselho Universitário — dispõe sobre projetos que devem passar pelo CEP"),
            ("", ""),
            ("Por que fui atrás disso:",
             "no OE-1 original eu previa validar a classificação MST em ~700 imagens usando três anotadores externos via Prolific. Isso é crowdsourcing pago — margem cinzenta."),
            ("", ""),
            ("O que a resolução diz:",
             "Art. 1º — todo projeto que envolve seres humanos direta ou indiretamente precisa passar pelo CEP antes de começar. Não pode retroagir."),
            ("Art. 8º:",
             "pesquisas que não envolvem seres humanos ficam dispensadas do cadastro, mas exigem Declaração de Responsabilidade assinada por estudante, orientador e chefe de departamento."),
        ],
        footer="",
    )

    add_bullets(
        prs,
        "O trade-off que enfrentei",
        [
            ("Opção A — manter Prolific:",
             "amostra maior (700 × 3 = 2.100 anotações), diversidade de anotadores, mas exige submissão ao CEP com prazo estimado de 2-3 meses"),
            ("Opção B — abandonar validação humana:",
             "confiaria só no SkinToneNet (Pereira 2026), mas isso vira ponto único de falha — foi exatamente a crítica do NotebookLM"),
            ("Opção C — validação interna reduzida:",
             "eu + o senhor anotamos ~200-300 imagens estratificadas por raça e MST. Menor escala, mas suficiente estatisticamente para inter-annotator agreement + concordância com SkinToneNet"),
            ("", ""),
            ("Decisão que tomei:",
             "Opção C — preserva rigor metodológico via estratificação, elimina risco regulatório, não bloqueia o cronograma"),
            ("", ""),
            ("O que preciso do senhor hoje:",
             "1) validar a decisão; 2) me dizer quem é o chefe do departamento que assina a Declaração"),
        ],
        footer="Se preferir voltar para Prolific, eu topo — mas aí a defesa em outubro fica em risco por causa do prazo CEP.",
    )

    add_table_slide(
        prs,
        "O que muda tecnicamente entre v3.5 e v3.6 do OE-1",
        ["Aspecto", "v3.5 (com Prolific)", "v3.6 (interno)"],
        [
            ["Escala da validação", "~700 imagens × 3 anotadores externos", "~200-300 imagens × 2-3 anotadores internos"],
            ["Diversidade dos anotadores", "Crowdworkers com perfil variado", "Mestrando + orientador (potencial viés declarado)"],
            ["Métrica de concordância", "Fleiss' kappa (múltiplos anotadores)", "Cohen's kappa (par a par)"],
            ["Poder estatístico", "IC 95 % com margem menor", "IC 95 % com margem maior — declarado como limitação"],
            ["Risco regulatório", "Requer CEP (2-3 meses)", "Art. 8º — apenas Declaração"],
            ["Cronograma", "Bloqueia defesa em outubro", "Compatível com outubro"],
        ],
        col_widths=[3.0, 4.5, 5.0],
        highlight_rows=[4, 5],
        footer="A perda de escala é real. Vou declarar explicitamente como limitação metodológica no Cap 4.",
    )

    # SEÇÃO 4 — Solicitação pessoal
    add_section_divider(prs, "4", "Solicitação pessoal — extensão de dois meses",
                        "Nascimento da minha filha em 28/março e afastamento subsequente")

    add_bullets(
        prs,
        "O contexto e o que estou pedindo",
        [
            ("O que aconteceu:",
             "minha filha nasceu em 28 de março. Fiquei afastado das atividades acadêmicas nas semanas seguintes para dar o cuidado que ela e a mãe precisaram."),
            ("", ""),
            ("O que preservei:",
             "os marcos de 15/jul e 30/jul continuam de pé. Não estou pedindo prazo para a primeira revisão nem para o pedido formal."),
            ("", ""),
            ("O que estou pedindo:",
             "extensão de dois meses no prazo da defesa da qualificação — de agosto para outubro. Isso me dá margem para escrever a versão final com o rigor que o trabalho merece."),
            ("", ""),
            ("Documento:",
             "carta em dois parágrafos redigida, pronta para envio. Preciso só do senhor me confirmar o trâmite — se é carta direta, requerimento via SEI ou processo formal com anexos."),
        ],
        footer="",
    )

    # SEÇÃO 5 — Debate
    add_section_divider(prs, "5", "Debate aberto",
                        "Seis perguntas onde preciso da opinião do orientador")

    add_bullets(
        prs,
        "Perguntas técnicas",
        [
            ("1.",
             "A quantidade de configurações comparativas (A, B, C, D) é suficiente ou o senhor recomenda incluir uma quinta variante — por exemplo, FiLM com skin tone contínuo (ITA/L*) em vez de MST discreto?"),
            ("2.",
             "No Cap 2, faz sentido eu incluir uma subseção de meia página comparando FiLM com CBN, cross-attention e AdaIN — ou deixo isso para a defesa oral?"),
            ("3.",
             "Sobre a extensão para face recognition (OE-4) — o senhor concorda que a métrica principal deva ser TPR@FAR=1e-4 estratificada por raça, ou prefere DR agregado?"),
        ],
        footer="",
    )

    add_bullets(
        prs,
        "Perguntas administrativas",
        [
            ("4.",
             "Quem é o chefe do departamento que deve assinar a Declaração de Responsabilidade do CEP?"),
            ("5.",
             "Qual o trâmite correto para a solicitação de extensão de dois meses — carta direta, requerimento SEI ao PPG ou processo formal com anexos?"),
            ("6.",
             "O senhor tem preferência sobre template Overleaf (padrão institucional Unifesp / ICT vs template ABNT genérico do abnTeX2)?"),
            ("", ""),
            ("Bônus (se der tempo):",
             "alguma sugestão sobre composição da banca preliminar? Estou pensando em alguém de fairness (pode ser externo) + alguém de visão computacional aplicada."),
        ],
        footer="",
    )

    add_bullets(
        prs,
        "Próximos sete dias — como vou organizar a semana",
        [
            ("Hoje (08/jul):",
             "aplicar o que sair desta reunião; setar Overleaf; iniciar transposição do Cap 1"),
            ("09-10/jul:",
             "finalizar Cap 1 no Overleaf; transpor Cap 2 (Revisão bibliográfica — o mais denso)"),
            ("11-12/jul:",
             "transpor Cap 3 (Objetivos) e Cap 4 (Metodologia)"),
            ("13/jul:",
             "fechar Cap 5 (Cronograma) com marcos ajustados hoje; revisão final integrada"),
            ("14/jul:",
             "leitura completa de ponta a ponta; ajustes de coerência e ABNT"),
            ("15/jul:",
             "entrega da primeira revisão ao senhor via Overleaf compartilhado"),
            ("", ""),
            ("Em paralelo:",
             "iniciar tramitação da Declaração de Responsabilidade + protocolar carta de extensão"),
        ],
        footer="",
    )

    add_section_divider(prs, "", "Obrigado",
                        "Aberto para o debate")

    return prs


def main() -> None:
    here = Path(__file__).parent
    out = here / "material_reuniao_orientador_2026-07-08.pptx"
    prs = build_presentation()
    prs.save(out)
    print(f"Apresentacao gerada: {out}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
