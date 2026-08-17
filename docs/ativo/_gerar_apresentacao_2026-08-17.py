"""Gera apresentacao PowerPoint para reuniao Ago/2026 com Prof. Marcos Quiles.

5 slides tematicos:
  1. Capa (meta + contagem T-menos qualificacao)
  2. Cronograma atual vs planejado
  3. Resumo dos pipelines (6 etapas com status)
  4. Contribuicoes esperadas (3 eixos, 7 contribuicoes)
  5. Pontos de decisao (SkinToneNet + itens abertos)

Uso:
    python docs/ativo/_gerar_apresentacao_2026-08-17.py
    -> produz: docs/ativo/material_reuniao_orientador_2026-08-17.pptx
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Paleta academica sobria (consistente com scripts anteriores)
NAVY = RGBColor(0x1F, 0x2A, 0x4E)
BLUE_MID = RGBColor(0x3D, 0x68, 0xBB)
BLUE_LIGHT = RGBColor(0xB4, 0xC4, 0xE8)
GRAY_DK = RGBColor(0x3D, 0x42, 0x4E)
GRAY_MD = RGBColor(0x70, 0x76, 0x82)
GRAY_LT = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
YELLOW = RGBColor(0xF5, 0xB7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

REUNIAO = date(2026, 8, 17)
QUALIFICACAO = date(2026, 9, 30)
DIAS_ATE_QUALI = (QUALIFICACAO - REUNIAO).days


# ============================================================
# Helpers de layout (mesmo padrao dos scripts anteriores)
# ============================================================
def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Contador T-menos na barra lateral
    tcount = slide.shapes.add_textbox(Inches(0.15), Inches(2.7), Inches(2.0), Inches(2.0))
    tcf = tcount.text_frame
    tcf.word_wrap = True
    p_t = tcf.paragraphs[0]
    p_t.text = f"T-{DIAS_ATE_QUALI}"
    p_t.font.size = Pt(56)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE
    p_sub = tcf.add_paragraph()
    p_sub.text = "dias ate\nqualificacao"
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = BLUE_LIGHT

    tx = slide.shapes.add_textbox(Inches(2.5), Inches(1.3), Inches(10.5), Inches(2.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Andamento — Ago/2026"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Etapa 1 iniciada + pipeline organizado + achado bloqueante"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GRAY_DK

    p3 = tf.add_paragraph()
    p3.text = "3 meses de adiantamento sobre o cronograma proposto na qualificacao"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRAY_MD
    p3.font.italic = True

    meta = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(10.5), Inches(2.5))
    mf = meta.text_frame
    mf.word_wrap = True
    rows = [
        ("Mestrando:", "Marcello Ozzetti"),
        ("Orientador:", "Prof. Marcos Quiles"),
        ("Programa:", "Mestrado em Ciencia da Computacao — Unifesp / ICT"),
        ("Reuniao:", f"{REUNIAO.strftime('%d/%m/%Y')} (semana 17-Ago)"),
        ("Qualificacao:", f"{QUALIFICACAO.strftime('%d/%m/%Y')}"),
    ]
    for i, (k, v) in enumerate(rows):
        p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        p.text = f"{k}  {v}"
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY_DK


def add_title(slide, text: str) -> None:
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.5), Inches(0.9))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, text: str) -> None:
    if not text:
        return
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
    pf = ft.text_frame.paragraphs[0]
    pf.text = text
    pf.font.size = Pt(10)
    pf.font.color.rgb = GRAY_MD
    pf.font.italic = True


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list,
    rows: list,
    footer: str = "",
    col_widths: list | None = None,
    highlight_rows: list | None = None,
    status_col: int | None = None,
    font_size: int = 11,
) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.5)
    height = Inches(5.2)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = WHITE

    highlight_rows = highlight_rows or []
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LT
            # coloracao especial da coluna de status
            if status_col is not None and ci == status_col:
                colored = _status_color(str(val))
                if colored is not None:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(font_size)
                            run.font.bold = True
                            run.font.color.rgb = colored
                    continue
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = GRAY_DK

    add_footer(slide, footer)


def _status_color(val: str) -> RGBColor | None:
    v = val.strip().lower()
    if v.startswith("adiantado") or v.startswith("pronto") or v.startswith("entregue"):
        return GREEN
    if v.startswith("bloqueado") or v.startswith("bloqueio"):
        return ACCENT
    if v.startswith("atencao") or v.startswith("em curso") or v.startswith("em risco"):
        return RGBColor(0xC6, 0x8A, 0x00)  # amber para leitura em branco
    if v.startswith("planejado") or v.startswith("skeleton"):
        return GRAY_MD
    return None


# ============================================================
# SLIDE 2 — Cronograma atual vs planejado
# ============================================================
def slide_cronograma_vs(prs: Presentation) -> None:
    add_table_slide(
        prs,
        "Cronograma — realizado vs planejado (proposta da qualificacao)",
        ["Periodo", "Etapa", "Planejado na qualificacao", "Status atual", "Delta"],
        [
            ["Jul/2026", "Escrita", "Fechamento da monografia", "Entregue", "no prazo"],
            ["Ago/2026 (hoje)", "Etapa 1 — prep.", "(nao iniciado)", "Adiantado", "+3 meses"],
            ["Set/2026", "Qualificacao", "Exame de qualificacao", "Em curso", f"T-{DIAS_ATE_QUALI}"],
            ["Out/2026", "Correcoes", "Aplicar apontamentos da banca", "Planejado", "—"],
            ["Nov/2026", "Etapa 1 — formal", "SkinToneNet + FairFace + validacao humana", "Bloqueado (weights)", "risco 2"],
            ["Dez/2026", "Etapa 2", "Auditoria fenotipica FairFace", "Skeleton pronto", "no prazo"],
            ["Jan-Mar/2027", "Etapa 3", "ConvNeXt-T + FiLM (A/B/C/D)", "Skeleton pronto", "no prazo"],
            ["Abr/2027", "Etapa 4", "6 baselines + triangulacao Pareto", "Skeleton pronto", "no prazo"],
            ["Mai/2027", "Etapa 5", "Transferencia fair RFW/BFW", "Skeleton pronto", "no prazo"],
            ["Jun/2027", "Etapa 6", "Decomposicao Latinx", "Skeleton pronto", "no prazo"],
            ["Nov/2026 - Jul/2027", "Redacao", "Elaboracao paralela", "Planejado", "—"],
            ["2o sem 2027", "Defesa", "Defesa da dissertacao", "Planejado", "—"],
        ],
        col_widths=[1.9, 1.5, 4.3, 2.6, 2.2],
        highlight_rows=[1, 2, 4],
        status_col=3,
        footer="Status alinhado ao Cap 5 (cronograma+riscos). Bloqueio da Etapa 1 formal tem 3 mitigacoes ativas (Slide 5).",
        font_size=10,
    )


# ============================================================
# SLIDE 3 — Resumo dos pipelines (6 etapas com status)
# ============================================================
def slide_pipelines(prs: Presentation) -> None:
    add_table_slide(
        prs,
        "Pipelines das 6 etapas — estado do codigo em Ago/2026",
        ["#", "Etapa", "Modulo (src/face_bias/)", "Codigo entregue", "Testes"],
        [
            [
                "1",
                "Classificador MST",
                "mst/ (skintonenet + validation + sensitivity)",
                "Wrapper + cache SQLite + CLI + plano B (reproducao)",
                "16 unit tests OK",
            ],
            [
                "2",
                "Auditoria FairFace",
                "audit/ (fairface_mst + cross_matrix)",
                "Skeleton com interfaces",
                "skip (TODO Dez)",
            ],
            [
                "3",
                "Classificador condicionado",
                "conditioning/ (film + clip_prompts + injector)",
                "Skeleton com interfaces",
                "skip (TODO Jan)",
            ],
            [
                "4",
                "Baselines + Pareto",
                "baselines/ + fairness/ (9 modulos)",
                "Skeleton com interfaces",
                "skip (TODO Abr)",
            ],
            [
                "5",
                "Transferencia fair",
                "transfer/ (rfw + bfw + bisenet + qualidade)",
                "Skeleton com interfaces",
                "skip (TODO Mai)",
            ],
            [
                "6",
                "Sintese decompositiva",
                "decomposition/ (error_decomp)",
                "Skeleton com interfaces",
                "skip (TODO Jun)",
            ],
        ],
        col_widths=[0.5, 2.8, 3.6, 4.4, 1.7],
        highlight_rows=[0],
        footer="Herdado do MBA: backbones (ConvNeXt-T ja registrado), losses (ArcFace/AdaFace/MagFace/SupCon), MTCNN, DVC, MLflow.",
        font_size=10,
    )


# ============================================================
# SLIDE 4 — Contribuicoes esperadas (3 eixos, 7 contribuicoes)
# ============================================================
def slide_contribuicoes(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Contribuicoes esperadas — 3 eixos, 7 contribuicoes")

    tx1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.3))
    tf1 = tx1.text_frame
    tf1.word_wrap = True

    eixos = [
        (
            "EIXO FENOTIPICO-EMPIRICO",
            "Contribuicoes 1 e 2",
            "Documenta a distribuicao real do fenotipo (MST) dentro das classes raciais",
        ),
        (
            "EIXO METODOLOGICO-ARQUITETURAL",
            "Contribuicoes 3, 4 e 7",
            "Introduz o tom de pele como sinal condicionante via mecanismo arquitetural",
        ),
        (
            "EIXO DIAGNOSTICO-ESTRUTURAL",
            "Contribuicoes 5 e 6",
            "Decompoe o erro em componentes fenotipico e algoritmico, com transferencia fair",
        ),
    ]

    for i, (titulo, mapa, resumo) in enumerate(eixos):
        p1 = tf1.paragraphs[0] if i == 0 else tf1.add_paragraph()
        p1.text = titulo
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.space_before = Pt(0 if i == 0 else 12)

        p2 = tf1.add_paragraph()
        p2.text = mapa
        p2.font.size = Pt(11)
        p2.font.italic = True
        p2.font.color.rgb = ACCENT

        p3 = tf1.add_paragraph()
        p3.text = resumo
        p3.font.size = Pt(11)
        p3.font.color.rgb = GRAY_DK
        p3.space_after = Pt(4)

    tx2 = slide.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(8.0), Inches(5.3))
    tf2 = tx2.text_frame
    tf2.word_wrap = True

    contribs = [
        ("Contribuicao 1.", "Avaliacao comparativa de modelos MST pre-treinados com sensitivity analysis"),
        ("Contribuicao 2.", "Primeira matriz publica MST x classes raciais sobre o FairFace val set"),
        ("Contribuicao 3.", "Primeira aplicacao documentada de FiLM-conditioning a race classification multi-classe"),
        ("Contribuicao 4.", "Triangulacao de metricas multi-classe (DR + worst-class F1 + EO/EqOdds), aderente a ISO 19795-10"),
        ("Contribuicao 5.", "Demonstracao empirica de transferencia fair de classification para face recognition, com controle de pixel information"),
        ("Contribuicao 6.", "Decomposicao quantitativa do erro Latinx: fenotipico irredutivel vs algoritmico mitigavel"),
        ("Contribuicao 7.", "Estudo comparativo entre 4 configuracoes de conditioning (baseline, FiLM linear, FiLM com porta, FiLM CLIP)"),
    ]

    for i, (head, body) in enumerate(contribs):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        run1 = p.add_run()
        run1.text = head + "  "
        run1.font.size = Pt(11)
        run1.font.bold = True
        run1.font.color.rgb = NAVY
        run2 = p.add_run()
        run2.text = body
        run2.font.size = Pt(11)
        run2.font.color.rgb = GRAY_DK
        p.space_after = Pt(5)

    add_footer(
        slide,
        "Sem mudancas de escopo desde Jul/2026. Cada contribuicao continua vinculada a um objetivo e uma hipotese testavel (Cap 3).",
    )


# ============================================================
# SLIDE 5 — Pontos de decisao para o orientador
# ============================================================
def slide_decisoes(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Pontos de decisao — pauta da reuniao")

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(5.3))
    tf = tx.text_frame
    tf.word_wrap = True

    itens = [
        (
            "1. Achado bloqueante — SkinToneNet (Matias 2026, arXiv 2603.02475)",
            "Weights e dataset STW ainda nao publicados ('code and data available soon'). Autores no ICMC/USP + IMPA.",
            "Recomendacao: enviar email institucional ao Prof. Joao Batista Neto (ICMC/USP) pedindo acesso antecipado. Rascunho pronto em docs/ativo/email_skintonenet_authors.md — precisa da sua aprovacao/ajustes.",
            "Plano B ativo: pipeline de reproducao pronto (pipelines/03b_train_mst_reproduction.py); backend classico 'stone_monk' como oraculo temporario para nao bloquear Etapas 2 e 3.",
        ),
        (
            "2. Foco Ago-Set/2026 — qualificacao vs adiantamento tecnico",
            "Foco proposto: apresentacao (prioridade maxima) + Etapa 2 skeleton avancando em paralelo.",
            "Decisao aberta: quer que eu comece a implementar a Etapa 2 (auditoria FairFace x MST) ja nesta semana, ou aguardar o pos-qualificacao?",
            "",
        ),
        (
            "3. Ordem de execucao dos baselines (Etapa 4, Abr/2027)",
            "Sao 6 baselines: ResNet-34, ConvNeXt-T puro, FSCL+, Group DRO, FineFACE, Adversarial Debias.",
            "Decisao aberta: rodar todos os 6 desde o comeco ou selecionar 3 criticos (ResNet-34, FSCL+, FineFACE) para reduzir tempo de GPU?",
            "",
        ),
        (
            "4. Sensitivity backend — licenca GPL do 'stone_monk'",
            "Biblioteca ChenglongMa/SkinToneClassifier tem paleta Monk 10-classes disponivel hoje, mas licenca GPL-3.0.",
            "Uso previsto: dep opcional, invocacao out-of-process (import lazy ja implementado). Sem linkagem ao codigo proprio.",
            "Decisao aberta: aprovado esse uso ou trocar por outro backend (ex: modelo HuggingFace)?",
        ),
    ]

    for i, item in enumerate(itens):
        titulo, contexto, decisao, extra = item
        p1 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p1.text = titulo
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.space_before = Pt(0 if i == 0 else 10)

        p2 = tf.add_paragraph()
        p2.text = "  Contexto:  " + contexto
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY_DK

        p3 = tf.add_paragraph()
        run_lbl = p3.add_run()
        run_lbl.text = "  Decisao:  "
        run_lbl.font.size = Pt(11)
        run_lbl.font.bold = True
        run_lbl.font.color.rgb = ACCENT
        run_body = p3.add_run()
        run_body.text = decisao
        run_body.font.size = Pt(11)
        run_body.font.color.rgb = GRAY_DK

        if extra:
            p4 = tf.add_paragraph()
            p4.text = "  " + extra
            p4.font.size = Pt(10)
            p4.font.italic = True
            p4.font.color.rgb = GRAY_MD

    add_footer(
        slide,
        "Todo o codigo mencionado ja esta implementado e testado; nada depende dessa reuniao para continuar tecnicamente — depende para decisoes de escopo.",
    )


# ============================================================
# Main
# ============================================================
def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    slide_cronograma_vs(prs)
    slide_pipelines(prs)
    slide_contribuicoes(prs)
    slide_decisoes(prs)

    return prs


def main() -> None:
    prs = build_presentation()
    out_dir = Path(__file__).parent
    out = out_dir / "material_reuniao_orientador_2026-08-17.pptx"
    prs.save(out)
    print(f"OK: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
