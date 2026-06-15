"""Gera apresentacao PowerPoint da reuniao 2026-06-15.

Uso:
    python _gerar_apresentacao_2026-06-15.py
    -> produz: docs/ativo/material_reuniao_orientador_2026-06-15.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Paleta academica sobria
NAVY = RGBColor(0x1F, 0x2A, 0x4E)
GRAY_DK = RGBColor(0x3D, 0x42, 0x4E)
GRAY_MD = RGBColor(0x70, 0x76, 0x82)
GRAY_LT = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
YELLOW = RGBColor(0xF5, 0xB7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(2.5), Inches(1.3), Inches(10.5), Inches(2.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Evolução da semana"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Corpus consolidado (101 fichas) e análise CLIP vs FiLM"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GRAY_DK

    p3 = tf.add_paragraph()
    p3.text = "Pente fino crítico do corpus e resposta à sugestão do orientador"
    p3.font.size = Pt(16)
    p3.font.color.rgb = GRAY_MD
    p3.font.italic = True

    meta = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(10.5), Inches(2.5))
    mf = meta.text_frame
    mf.word_wrap = True
    rows = [
        ("Mestrando:", "Marcello Ozzetti"),
        ("Orientador:", "Prof. Marcos Quiles"),
        ("Programa:", "Mestrado em Ciência da Computação — Unifesp / ICT"),
        ("Reunião:", "15 de junho de 2026"),
    ]
    for i, (k, v) in enumerate(rows):
        p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        p.text = f"{k}  {v}"
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY_DK


def add_section_divider(prs: Presentation, num: str, title: str, subtitle: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tx_num = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3), Inches(3))
    pn = tx_num.text_frame.paragraphs[0]
    pn.text = num
    pn.font.size = Pt(140)
    pn.font.bold = True
    pn.font.color.rgb = WHITE

    tx_t = slide.shapes.add_textbox(Inches(4.3), Inches(2.5), Inches(8.5), Inches(2))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = GRAY_LT


def add_title(slide, text: str) -> None:
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, prs, text: str) -> None:
    if not text:
        return
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
    pf = ft.text_frame.paragraphs[0]
    pf.text = text
    pf.font.size = Pt(10)
    pf.font.color.rgb = GRAY_MD
    pf.font.italic = True


def add_bullets(prs: Presentation, title: str, bullets: list, footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(5.4))
    tf = tx.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            head, body = item
            run1 = p.add_run()
            run1.text = head + "  "
            run1.font.size = Pt(16)
            run1.font.bold = True
            run1.font.color.rgb = NAVY
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(16)
            run2.font.color.rgb = GRAY_DK
        else:
            p.text = "— " + item
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_DK
        p.space_after = Pt(8)

    add_footer(slide, prs, footer)


def add_table_slide(prs: Presentation, title: str, headers: list, rows: list, footer: str = "",
                    col_widths: list | None = None, highlight_rows: list | None = None) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.5)
    height = Inches(5.2)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = WHITE

    highlight_rows = highlight_rows or []
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = GRAY_DK

    add_footer(slide, prs, footer)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Capa
    add_title_slide(prs)

    # 2. Agenda
    add_bullets(
        prs,
        "Agenda da reunião",
        [
            ("1.", "Status das decisões da reunião anterior (08/jun)"),
            ("2.", "Corpus consolidado: 101 fichas em 11 tracks temáticos"),
            ("3.", "Pente fino crítico do corpus (classificação por impacto)"),
            ("4.", "Resposta à sugestão do orientador: CLIP (Contrastive Language-Image Pre-training)"),
            ("5.", "FiLM em fairness: lacuna na literatura como contribuição original"),
            ("6.", "Proposta: estudo comparativo de mecanismos de conditioning"),
            ("7.", "Plano de escrita e decisões a alinhar"),
        ],
        footer="",
    )

    # SEÇÃO 1 — Status decisões
    add_section_divider(prs, "1", "Status das decisões", "Reunião anterior — 08/jun/2026")

    add_table_slide(
        prs,
        "Decisões da reunião anterior — status",
        ["#", "Decisão", "Status", "Observações"],
        [
            ["1", "Pipeline 6 etapas validado", "Concluído", "Detalhado etapa por etapa"],
            ["2", "SkinToneNet pré-treinado", "Concluído", "Leitura integral concluída"],
            ["3", "Corpus ≥100 artigos", "Concluído", "101 fichas (+46 na última rodada)"],
            ["4", "≥20 artigos 2025-2026", "Concluído", "25 fichas em 2025-2026"],
            ["5", "Adicionar CLIP conditioning", "Concluído", "Track I criada — 14 fichas + análise técnica"],
            ["6", "Primeira revisão ao orientador em 15/jul/2026", "Em andamento", "Estrutura sendo organizada — LaTeX ainda não iniciado"],
        ],
        col_widths=[0.5, 4.6, 2.0, 5.4],
        highlight_rows=[5],
        footer="5 itens concluídos e 1 em andamento conforme planejado",
    )

    # SEÇÃO 2 — Corpus 101 fichas
    add_section_divider(prs, "2", "Corpus consolidado", "101 fichas em 11 tracks temáticos")

    add_bullets(
        prs,
        "Como o corpus foi ampliado: 55 → 101 fichas em 11 tracks",
        [
            ("Critérios da busca:", "Top venues (CVPR, ICCV, ICLR, ECCV, Nature SR) com filtros menos restritivos"),
            ("Recorte temporal:", "Foco em 2025-2026 para refletir o estado da arte mais recente"),
            ("Fontes consultadas:", "arXiv, OpenAccess (CVPR/ICCV), Semantic Scholar e busca direta por DOI"),
            ("Estratégia em 3 levas:", "Top venues recentes → fundadores de FR → temas complementares"),
            ("Reorganização temática:", "3 novos tracks criados a partir dos achados (I, J, K)"),
            ("", ""),
            ("Resultado:", "+46 fichas / 25 fichas em 2025-2026 / 3 tracks novos"),
        ],
        footer="",
    )

    add_table_slide(
        prs,
        "Densidade das fichas por Tracking",
        ["Track", "Tema", "O que é o track", "Principais artigos", "Fichas"],
        [
            ["A", "Race classification", "Classificação racial em 7 categorias sobre o FairFace", "Kärkkäinen & Joo 2021, AlDahoul 2024, Lin (FairGRAPE) 2022", "4"],
            ["B", "FR fairness", "Datasets, métricas e auditorias de viés em reconhecimento facial", "Buolamwini 2018, Grother (NIST) 2019, Wang (RFW) 2019", "16"],
            ["C", "Skin tone (MST)", "Escalas e classificadores de tom de pele como sinal auxiliar", "Fitzpatrick 1988, Schumann 2023, Pereira (SkinToneNet) 2026", "7"],
            ["D", "Mitigação algorítmica", "Algoritmos para reduzir viés (contrastive, DRO, arquitetural)", "Park (FSCL+) 2022, Sagawa (Group DRO) 2020, Manzoor 2024", "10"],
            ["E", "Auditoria e surveys", "Auditorias sistêmicas e revisões da literatura", "Mehrabi 2021, Dominguez (DSAP) 2024, Buolamwini 2018", "14"],
            ["F", "Fundamentação ética", "Validade conceitual de 'raça' como variável científica", "Lewontin 1972, Fuentes 2019, Neto 2025", "3"],
            ["G", "Mecanismos ML paradigmáticos", "Mecanismos teóricos clássicos (métricas, conditioning)", "Hardt 2016, Perez (FiLM) 2018, Kleinberg 2017", "9"],
            ["I", "VLM / CLIP em fairness", "Fairness em modelos vision-language (CLIP, BLIP)", "Luo (FairCLIP) 2024, Dehdashtian (FairerCLIP) 2024, BendVLM 2024", "14"],
            ["J", "Conditioning moderno", "Adaptadores alternativos ao FiLM (LoRA, ViT-mask)", "Bian (LoRA-FAIR) 2025, Tian (FairViT) 2024, Zhao (AIM-Fair) 2025", "5"],
            ["K", "Fundadores de FR", "Losses fundadoras de reconhecimento facial", "Schroff (FaceNet) 2015, Deng (ArcFace) 2019, Kim (AdaFace) 2022", "6"],
            ["L", "Auxiliar / complementar", "Direções adjacentes (federated, synthetic, post-hoc)", "Salvador (FairCal) 2021, FairImagen 2025, VoIDFace 2025", "13"],
        ],
        col_widths=[0.7, 2.1, 3.4, 5.6, 0.7],
        footer="Tracks I, J e K são os novos tracks da ampliação recente do corpus",
    )

    # SEÇÃO 3 — Pente fino
    add_section_divider(prs, "3", "Pente fino", "Classificação por impacto na tese")

    add_table_slide(
        prs,
        "Classificação crítica das 101 fichas",
        ["Categoria", "Fichas", "%", "Implicação para a tese"],
        [
            ["Forte favorável", "12", "11.9%", "Fundamentam etapas e contribuições diretamente"],
            ["Favorável", "38", "37.6%", "Suporte / alinhadas ao argumento"],
            ["Neutra / contextual", "18", "17.8%", "Background, não influenciam decisões"],
            ["Caminho alternativo", "26", "25.7%", "Cobertos pelo estudo comparativo de mecanismos de conditioning"],
            ["Conflito moderado", "5", "5.0%", "Endereçáveis — resposta defensiva mapeada"],
            ["Conflito forte", "2", "2.0%", "Pangelinan (via H6); Neto (limitação reconhecida)"],
        ],
        col_widths=[3.2, 1.0, 1.0, 7.3],
        highlight_rows=[0, 5],
        footer="Apenas 2 conflitos fortes — ambos com resposta defensiva pronta. Nenhuma refutação categórica sem resposta.",
    )

    add_bullets(
        prs,
        "Conclusão do pente fino",
        [
            ("Resumo:", "68 fichas (67,3%) alinhadas ou neutras / 26 caminho alternativo / 7 conflitos"),
            ("Conflito forte 1:", "Pangelinan 2023 — leitura integral confirmou: skin tone isolado não causa o gap; pixel info explica parcialmente o FNMR → endossa nossa Hipótese H6"),
            ("Conflito forte 2:", "Neto 2025 — balanceamento não basta, reconhecido como limitação e parte do framing"),
            ("Caminho alternativo:", "26 fichas (VLM, CLIP) cobertas pelo estudo comparativo de mecanismos de conditioning"),
            ("Pontos fortes:", "12 fichas forte favorável fundamentam as etapas 1-5 do pipeline"),
            ("", ""),
            ("Decisão:", "Corpus está sólido para sustentar a primeira versão da escrita"),
        ],
        footer="",
    )

    add_bullets(
        prs,
        "Status das leituras críticas",
        [
            ("Cobertura de PDFs:", "92 de 101 fichas (91 %) com PDF no repositório — 9 restantes em paywall institucional ou fontes históricas"),
            ("Camada 1 (críticas):", "14 fichas — todas com leitura integral concluída e fichas em estado VERIFIED"),
            ("Inclui:", "Pereira (SkinToneNet), Luo (FairCLIP), Pangelinan (refutação central), Hardt, Perez (FiLM), Madras (LAFTR), Schumann (MST), Karkkainen (FairFace), AlDahoul (SOTA)"),
            ("Próximas leituras:", "Camada 2 — 38 fichas favoráveis para aprofundamento durante a escrita do Capítulo 2"),
            ("", ""),
            ("Achado novo (Pangelinan):", "imagens de pessoas do mesmo gênero são intrinsecamente mais similares — possível análogo intra-racial a investigar"),
        ],
        footer="",
    )

    # SEÇÃO 4 — CLIP
    add_section_divider(prs, "4", "Resposta ao orientador",
                        "CLIP = Contrastive Language-Image Pre-training (Radford et al., ICML 2021)")

    add_bullets(
        prs,
        "O que é o CLIP — confirmando a referência",
        [
            ("Paper original:", "Radford et al. — 'Learning Transferable Visual Models From Natural Language Supervision', ICML 2021"),
            ("Mecanismo:", "Dois encoders (visão e texto) treinados via contrastive loss em 400M pares (imagem, texto)"),
            ("Resultado:", "Espaço de embedding compartilhado — imagens e textos relacionados ficam próximos"),
            ("Capacidade central:", "Classificação zero-shot — basta uma descrição textual, sem fine-tuning"),
            ("Adoção em fairness:", "CLIP virou espinha dorsal da pesquisa em vision-language fairness desde 2022"),
            ("", ""),
            ("No contexto da tese:", "CLIP produz EMBEDDINGS; FiLM é MECANISMO de injeção — são complementares, não substitutos"),
        ],
        footer="Radford et al., ICML 2021 — arxiv.org/abs/2103.00020 — mais de 10 mil citações",
    )

    add_table_slide(
        prs,
        "Estudos de CLIP em fairness no corpus (Track I)",
        ["Paper", "Venue", "Abordagem", "Aplicação na tese"],
        [
            ["FairCLIP (Luo, 2024)", "CVPR 2024", "Optimal transport / Sinkhorn", "Baseline principal do estudo comparativo"],
            ["FairerCLIP (Dehdashtian, 2024)", "ICLR 2024", "Debias zero-shot via RKHS", "Baseline leve do estudo comparativo"],
            ["BendVLM (2024)", "arXiv", "Debiasing em test-time", "Alternativa sem retraining"],
            ["LoRA-FAIR (Bian, 2025)", "ICCV 2025", "LoRA + agregação fair", "Entra no estudo comparativo (Track J)"],
            ["FaceScanPaliGemma (AlDahoul, 2024)", "Nature SR 2026", "VLM fine-tuned no FairFace", "Baseline SOTA 7-class (75,7%)"],
            ["FairViT (Tian, 2024)", "ECCV 2024", "Adaptive attention masking", "Backbone alternativo"],
            ["Closed-form debias (2026)", "arXiv", "Solução analítica", "Alternativa eficiente"],
            ["Bias subspace VLM (2025)", "arXiv", "Projeção em subespaço", "Mecanismo alternativo"],
            ["Survey multimodal fairness (2024)", "arXiv", "Survey", "Contexto para o Capítulo 2"],
        ],
        col_widths=[3.3, 1.8, 3.2, 4.2],
        highlight_rows=[0, 1, 4],
        footer="14 fichas no total na Track I, além do paper original de CLIP (Radford 2021).",
    )

    # SEÇÃO 5 — FiLM
    add_section_divider(prs, "5", "FiLM em fairness",
                        "Achado crítico: lacuna na literatura = contribuição original")

    add_bullets(
        prs,
        "Estudos de FiLM em fairness",
        [
            ("Paper FiLM original:", "Perez et al., AAAI 2018 — verificado no corpus"),
            ("Aplicação em fairness pelos autores:", "Nenhuma — paper não aborda fairness"),
            ("Adoção em fairness na literatura:", "Rara — busca extensiva confirma a escassez"),
            ("Mecanismos análogos no corpus:", ""),
            ("  — LAFTR (Madras, 2018)", "adversarial conditioning (não é FiLM)"),
            ("  — FairViT (Tian, 2024)", "adaptive masking em attention (ViT)"),
            ("  — LoRA-FAIR (Bian, 2025)", "low-rank weight modulation"),
            ("", ""),
            ("Implicação:", "FiLM aplicado a fairness em classificação racial é uma direção original"),
        ],
        footer="Perez et al., AAAI 2018 — paper original não discute fairness nem na seção de Future Work.",
    )

    add_bullets(
        prs,
        "FiLM vs CLIP — esclarecimento conceitual",
        [
            ("Confusão categorial:", "FiLM e CLIP não são objetos comparáveis no mesmo nível"),
            ("FiLM (Perez, 2018):", "MECANISMO / camada que modula features via γ·x + β"),
            ("CLIP (Radford, 2021):", "MODELO que produz embeddings 512-dim de texto e imagem"),
            ("Relação real:", "CLIP fornece o SINAL; FiLM (ou cross-attention) é uma forma de injetar esse sinal"),
            ("Cronologia:", "FaceNet 2015, ArcFace 2019, FiLM 2018, CLIP 2021, LoRA 2021"),
            ("Status em 2024-2026:", "Todos continuam padrão na literatura — idade não implica obsolescência"),
            ("", ""),
            ("Combinação válida:", "Stable Diffusion = CLIP embeddings + cross-attention. Análogo: CLIP + FiLM é possível"),
        ],
        footer="",
    )

    # SEÇÃO 6 — Estudo comparativo de mecanismos
    add_section_divider(prs, "6", "Estudo comparativo de mecanismos de conditioning",
                        "FiLM, CLIP-conditioning e LoRA aplicados a race classification")

    add_table_slide(
        prs,
        "5 configurações de conditioning comparadas",
        ["Config", "Arquitetura", "Mecanismo de conditioning", "Origem"],
        [
            ["A", "ConvNeXt-T baseline", "Sem conditioning", "Controle"],
            ["B (proposta)", "ConvNeXt-T + FiLM", "MST 10-dim → γ, β", "Nossa contribuição central"],
            ["C", "ConvNeXt-T + FiLM", "Embedding CLIP-text → γ, β", "Híbrido CLIP + FiLM"],
            ["D", "ConvNeXt-T + cross-attention", "Embedding CLIP-text via attention", "Estilo FairCLIP"],
            ["E", "ConvNeXt-T + LoRA-FAIR", "Low-rank weight modulation", "LoRA-FAIR (Bian, 2025)"],
        ],
        col_widths=[1.0, 3.0, 4.5, 4.5],
        highlight_rows=[1],
        footer="Se CLIP superar FiLM, reportamos como informação útil — não como invalidação. Honra a recomendação do orientador.",
    )

    add_bullets(
        prs,
        "Justificativa para manter FiLM como mecanismo central",
        [
            ("Adequação a baixa dimensão:", "MST é vetor 10-dim — FiLM é o ponto ótimo"),
            ("Custo paramétrico:", "~1% sobre o backbone (vs. ~3% de cross-attention equivalente)"),
            ("Interpretabilidade:", "γ e β são por canal e podem ser visualizados"),
            ("Compatibilidade:", "Substitui parâmetros do BatchNorm — 'conditional BatchNorm'"),
            ("Compatibilidade com o estudo comparativo:", "Ablation limpa — liga ou desliga o FiLM mantendo o backbone fixo"),
            ("", ""),
            ("Argumento de originalidade:", "FiLM em fairness é direção pouco explorada — vira contribuição"),
        ],
        footer="Perez et al. (AAAI 2018) + Schumann (2023, MST) + Pereira (2026, SkinToneNet) — ancoragem direta",
    )

    # SEÇÃO 7 — Plano de escrita e decisões
    add_section_divider(prs, "7", "Plano de escrita e decisões",
                        "Cronograma das 4 semanas e pontos a alinhar")

    add_table_slide(
        prs,
        "Plano de escrita — 4 semanas até a primeira revisão (15/jul/2026)",
        ["Semana", "Período", "Entregável", "Observações"],
        [
            ["1", "15-21/jun (esta)", "Estrutura no repositório + Capítulo 1 (Introdução)", "Narrativa de abertura consolidada"],
            ["2", "22-28/jun", "Capítulo 2 (Revisão — 11 tracks)", "Pente fino do corpus e 101 fichas"],
            ["3", "29/jun-05/jul", "Capítulo 3 (Objetivos) e Capítulo 4 (Ampliação e busca de novas técnicas)", "Pipeline de 6 etapas validado"],
            ["4", "06-12/jul", "Capítulo 5 (Cronograma) e revisão final", "Integração e ABNT"],
            ["Target", "13-15/jul", "Ajustes finais e entrega ao orientador", "Primeira revisão"],
        ],
        col_widths=[1.0, 2.4, 5.0, 4.6],
        highlight_rows=[0],
        footer="LaTeX ainda não iniciado: estrutura sendo organizada no projeto base do GitHub. Migração para o Overleaf após consolidação.",
    )

    add_bullets(
        prs,
        "Decisões a alinhar hoje",
        [
            ("1.", "Concorda em manter FiLM como mecanismo central e incluir CLIP-conditioning no estudo comparativo?"),
            ("2.", "Os 11 tracks estão bem cobertos? Algum precisa de mais fichas?"),
            ("3.", "O esclarecimento conceitual CLIP (modelo) vs FiLM (mecanismo) faz sentido?"),
            ("4.", "Pixel information (achado de Pangelinan) pode ser formalizado como objetivo específico (OE-6) com hipótese H6 quantitativa?"),
            ("5.", "Co-orientador — já podemos formalizar?"),
            ("6.", "Sugestões para a banca preliminar?"),
        ],
        footer="",
    )

    add_bullets(
        prs,
        "Próximos passos imediatos",
        [
            ("1.", "Consolidar a estrutura no repositório do GitHub (ainda não iniciei o LaTeX)"),
            ("2.", "Migrar para o Overleaf assim que a estrutura estiver estável"),
            ("3.", "Consolidar a bibliografia a partir das 101 fichas"),
            ("4.", "Escrever Capítulo 1 (Introdução) usando a narrativa de abertura"),
            ("5.", "Incorporar as imagens já produzidas (pipeline FiLM, ConvNeXt-T, métricas de Hardt)"),
            ("6.", "Aprofundar a leitura das 12 fichas forte favorável"),
            ("", ""),
            ("Entrega para a próxima reunião:", "Capítulo 1 escrito e bibliografia consolidada"),
        ],
        footer="",
    )

    # Slide final
    add_section_divider(prs, "", "Obrigado",
                        "Discussão livre — dúvidas e ajustes antes da escrita")

    return prs


def main() -> None:
    here = Path(__file__).parent
    out = here / "material_reuniao_orientador_2026-06-15.pptx"
    prs = build_presentation()
    prs.save(out)
    print(f"Apresentacao gerada: {out}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
