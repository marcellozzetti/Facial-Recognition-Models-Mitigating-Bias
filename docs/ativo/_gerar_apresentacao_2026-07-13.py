"""Gera apresentacao PowerPoint da reuniao 2026-07-13.

Reuniao apos migracao da qualificacao para Overleaf. Marco 15/jul
antecipado — primeira revisao pronta para entrega.

- Plano de trabalho consolidado
- Estado da escrita: LaTeX/Overleaf migracao concluida
- Postura autoral sobre uso de LLM
- Adequacao etica CEP + tramite Declaracao de Responsabilidade
- Extensao de prazo
- Perguntas de debate

Uso:
    python _gerar_apresentacao_2026-07-13.py
    -> produz: docs/ativo/material_reuniao_orientador_2026-07-13.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x2A, 0x4E)
GRAY_DK = RGBColor(0x3D, 0x42, 0x4E)
GRAY_MD = RGBColor(0x70, 0x76, 0x82)
GRAY_LT = RGBColor(0xE8, 0xEA, 0xED)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(2.5), Inches(1.3), Inches(10.5), Inches(2.4))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Acompanhamento Semana"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    p2.text = "Escrita LaTeX em andamento, decisões éticas e ajuste de cronograma"
    p2.font.size = Pt(20)
    p2.font.color.rgb = GRAY_DK

    p3 = tf.add_paragraph()
    p3.text = "Mitigação de Viés Racial em Classificação Facial via Condicionamento por Tom de Pele em Arquiteturas Profundas"
    p3.font.size = Pt(13)
    p3.font.color.rgb = GRAY_MD
    p3.font.italic = True

    meta = slide.shapes.add_textbox(Inches(2.5), Inches(4.5), Inches(10.5), Inches(2.5))
    mf = meta.text_frame
    mf.word_wrap = True
    rows = [
        ("Mestrando:", "Marcello Ozzetti"),
        ("Orientador:", "Prof. Marcos Quiles"),
        ("Programa:", "Mestrado em Ciência da Computação — Unifesp / ICT"),
        ("Reunião:", "13 de julho de 2026"),
    ]
    for i, (k, v) in enumerate(rows):
        p = mf.paragraphs[0] if i == 0 else mf.add_paragraph()
        p.text = f"{k}  {v}"
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY_DK


def add_section_divider(prs: Presentation, num: str, title: str, subtitle: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tx_num = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(3), Inches(3))
    pn = tx_num.text_frame.paragraphs[0]
    pn.text = num
    pn.font.size = Pt(140)
    pn.font.bold = True
    pn.font.color.rgb = WHITE

    tx_t = slide.shapes.add_textbox(Inches(4.3), Inches(2.5), Inches(8.5), Inches(2))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = GRAY_LT


def add_title(slide, text: str) -> None:
    tx_t = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9))
    tf = tx_t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(12.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


def add_footer(slide, prs, text: str) -> None:
    if not text:
        return
    ft = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.4))
    pf = ft.text_frame.paragraphs[0]
    pf.text = text
    pf.font.size = Pt(10)
    pf.font.color.rgb = GRAY_MD
    pf.font.italic = True


def add_bullets(prs: Presentation, title: str, bullets: list, footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(5.4))
    tf = tx.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            head, body = item
            run1 = p.add_run()
            run1.text = head + "  "
            run1.font.size = Pt(16)
            run1.font.bold = True
            run1.font.color.rgb = NAVY
            run2 = p.add_run()
            run2.text = body
            run2.font.size = Pt(16)
            run2.font.color.rgb = GRAY_DK
        else:
            p.text = "— " + item
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY_DK
        p.space_after = Pt(8)

    add_footer(slide, prs, footer)


def add_quote_slide(prs: Presentation, title: str, quote: str, attribution: str = "", footer: str = "") -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.2), Inches(0.08), Inches(3.5))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = NAVY
    quote_box.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(1.8), Inches(2.2), Inches(10.5), Inches(3.8))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = GRAY_DK
    p.space_after = Pt(20)

    if attribution:
        p2 = tf.add_paragraph()
        p2.text = "— " + attribution
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY_MD

    add_footer(slide, prs, footer)


def add_table_slide(prs: Presentation, title: str, headers: list, rows: list, footer: str = "",
                    col_widths: list | None = None, highlight_rows: list | None = None) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, title)

    n_cols = len(headers)
    n_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.5)
    height = Inches(5.2)

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = WHITE

    highlight_rows = highlight_rows or []
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri in highlight_rows:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = GRAY_DK

    add_footer(slide, prs, footer)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Capa
    add_title_slide(prs)

    # 2. Objetivo da reunião
    add_bullets(
        prs,
        "Objetivo desta reunião",
        [
            ("", ""),
            ("1.", "Plano até a defesa"),
            ("", ""),
            ("2.", "Estado real da escrita — qualificação no Overleaf, Cap 2 em revisão ativa"),
            ("", ""),
            ("Já resolvido:", "Declaração de Responsabilidade CEP submetida ontem (12/jul) — Art. 8º da Res. 200/2021"),
        ],
        footer="Solicitação pessoal (extensão) e debate técnico fecham a reunião como pontos complementares.",
    )

    # 3. Recap do objetivo + pergunta central
    add_bullets(
        prs,
        "A tese em uma tela",
        [
            ("Problema empírico:",
             "F1 macro alto mascara disparidade Black ≈ 90 % vs Latinx ≈ 60 % em SOTA sobre FairFace (AlDahoul 2024)"),
            ("Pergunta de pesquisa:",
             "condicionar arquiteturalmente um classificador racial pelo tom de pele (MST) reduz essa disparidade sem sacrificar acurácia agregada?"),
            ("", ""),
            ("Objetivo geral:",
             "desenvolver e avaliar pipeline de classificação racial que incorpora tom de pele (escala Monk 10 classes) como sinal auxiliar condicionante via mecanismo arquitetural (FiLM)"),
            ("", ""),
            ("Contribuição principal:",
             "primeira instância empírica de tom de pele como contexto arquitetural para race classification multi-classe"),
            ("Diagnóstico central:",
             "erro Latinx = componente fenotípico irredutível (overlap MST) + componente algorítmico — a decomposição é o diferencial da tese"),
        ],
        footer="MST vs Fitzpatrick: 10 vs 6 classes, projetado para CV vs dermatologia. FiLM vs concatenação: modula vs re-aprende.",
    )

    # 4. Consolidado do corpus
    add_table_slide(
        prs,
        "Corpus consolidado — 104 fichas em 11 linhas de pesquisa",
        ["Linha de pesquisa", "N", "% corpus", "Papers-âncora"],
        [
            ["A — Race classification multi-classe", "4", "3,8 %", "Kärkkäinen 2021 (FairFace), AlDahoul 2024"],
            ["B — Face recognition fairness", "16", "15,4 %", "Wang RFW 2019, Robinson BFW 2020, Deng ArcFace"],
            ["C — Skin tone alternativo (Fitzpatrick, MST, ITA)", "7", "6,7 %", "Schumann 2023 (MST), Pereira 2026 (SkinToneNet)"],
            ["D — Mitigação algorítmica", "10", "9,6 %", "Madras LAFTR 2018, Sagawa Group DRO 2020"],
            ["E — Auditoria e Surveys", "14", "13,5 %", "Buolamwini 2018, Mehrabi 2021 survey"],
            ["F — Fundamentação científica e ética", "3", "2,9 %", "Fuentes AAPA 2019, Lewontin 1972"],
            ["G — Mecanismos ML paradigmáticos", "9", "8,7 %", "Perez FiLM 2018, Ba LayerNorm, He ResNet"],
            ["I — VLM / CLIP em fairness", "14", "13,5 %", "Radford CLIP 2021, Luo FairCLIP 2024"],
            ["J — Conditioning moderno", "5", "4,8 %", "Liu ConvNeXt 2022, Dhariwal ADM 2021"],
            ["K — Fundadores de FR (DeepFace, ArcFace, CosFace)", "6", "5,8 %", "Taigman DeepFace 2014, Wang CosFace 2018"],
            ["L — Auxiliar / complementar", "13", "12,5 %", "Karkkainen datasets, Hendrycks robustness"],
            ["R8 — Diversidade Latinx (antropo + genética + socio)", "3", "2,9 %", "Telles 2014, Bryc 2015, Pew 2017"],
        ],
        col_widths=[4.8, 0.7, 1.2, 5.8],
        footer="Fronteira temporal: 55 fichas (53 %) são 2024+, incluindo 5 papers de 2026 — corpus na fronteira absoluta.",
    )

    # SEÇÃO 1 — Plano até a defesa
    add_section_divider(prs, "1", "Plano até a defesa",
                        "Três marcos formais, novo horizonte de outubro")

    add_table_slide(
        prs,
        "Cronograma consolidado",
        ["#", "Marco", "Data", "Comentário"],
        [
            ["1", "Primeira revisão da qualificação ao orientador", "15/jul/2026", "Overleaf compartilhado hoje; Cap 1, 3, 4, 5 em versão inicial; Cap 2 em revisão ativa"],
            ["2", "Pedido formal de qualificação ao PPG-CC / ICT", "30/jul/2026", "Prazo regimental do Programa — mantido"],
            ["3", "Defesa da qualificação", "outubro/2026", "Ajuste de agosto → outubro (extensão de 2 meses solicitada)"],
        ],
        col_widths=[0.5, 5.5, 2.0, 5.0],
        highlight_rows=[0, 2],
        footer="Marco 1 antecipado em dois dias. Marcos 2 e 3 conforme planejado.",
    )

    # SEÇÃO 2 — Estado da escrita
    add_section_divider(prs, "2", "Estado da escrita",
                        "O que já está escrito e como vai virar Overleaf")

    add_bullets(
        prs,
        "Estado real da escrita — qualificação no Overleaf",
        [
            ("Cap 1 — Introdução:", "no Overleaf, versão inicial completa; ancora a disparidade Black 90 % × Latinx 60 %"),
            ("Cap 2 — Revisão:", "em escrita e revisão ativa; corpo bibliográfico está sendo densificado com 104 referências mapeadas"),
            ("Cap 3 — Objetivos:", "no Overleaf; OG + 6 OE + 6 hipóteses + 7 contribuições"),
            ("Cap 4 — Metodologia:", "no Overleaf; pipeline 6 etapas + 4 configurações + adequação ética CEP formalizada"),
            ("Cap 5 — Cronograma:", "no Overleaf; marcos ajustados conforme plano consolidado"),
            ("", ""),
            ("Transposição:", "esqueleto no Overleaf; bibliografia .bib com 104 entradas importada; referências cruzadas testadas"),
            ("Compartilhamento:", "convite do Overleaf enviado ao senhor a partir desta reunião — foco maior de revisão pedido no Cap 2"),
        ],
        footer="Cap 2 é o capítulo mais denso e onde ainda concentro esforço técnico esta semana.",
    )

    add_table_slide(
        prs,
        "Pipeline proposto — seis etapas encadeadas",
        ["Etapa", "O que faz", "Entrega / contribuição"],
        [
            ["1. Classificador MST", "SkinToneNet pré-treinado + validação interna em subset FairFace estratificado", "Insumo do pipeline (não é contribuição)"],
            ["2. Auditoria FairFace", "Aplicar SkinToneNet sobre todo FairFace validation set", "Matriz pública MST × race classes — Contribuição C2"],
            ["3. Race classifier com conditioning", "ConvNeXt-T fine-tuned em FairFace, camadas FiLM por estágio recebendo vetor MST", "Configurações A/B/C/D — Contribuição C3"],
            ["4. Comparação contra baselines", "6 baselines: ResNet-34, ConvNeXt-T puro, FSCL+, Group DRO, FineFACE, Adversarial", "Triangulação DR + worst-class F1 + EO_h — Contribuição C4"],
            ["5. Fair transferência", "Aplicar backbone fair em face recognition downstream (RFW ou BFW)", "Demonstração de fair transfer — Contribuição C5"],
            ["6. Síntese decompositiva", "Combinar C2 + C5 para separar componente fenotípico do algorítmico do erro Latinx", "Decomposição pixel info × skin tone — Contribuição C6/OE-6"],
        ],
        col_widths=[2.8, 5.5, 4.2],
        highlight_rows=[2, 5],
        footer="Etapa 3 é onde o mecanismo FiLM entra. Etapa 6 é o diagnóstico central da tese.",
    )

    add_table_slide(
        prs,
        "Quatro configurações do estudo comparativo (Etapa 3)",
        ["Config", "Backbone + conditioning", "Sinal", "Papel no estudo"],
        [
            ["A", "ConvNeXt-T sem conditioning", "—", "Baseline de controle"],
            ["B", "ConvNeXt-T + FiLM linear", "MST 10-dim", "Proposta principal — hipótese central"],
            ["C", "ConvNeXt-T + Gated FiLM", "MST 10-dim (não-linear)", "Ablação: linear vs não-linear importa?"],
            ["D", "ConvNeXt-T + FiLM", "Embedding CLIP-text 512-dim", "Recomendação da reunião 15/jun — CLIP como alternativa moderna"],
        ],
        col_widths=[1.0, 3.8, 3.5, 5.0],
        highlight_rows=[1],
        footer="B é a proposta central. C isola a não-linearidade. D compara com sinal semântico rico via CLIP.",
    )

    add_bullets(
        prs,
        "O achado que ancora a tese (Rodada 8, junho)",
        [
            ("Problema conhecido:",
             "modelos SOTA reportam F1 Latinx próximo de 60 % há três anos — persiste mesmo em modelos multimodais recentes (AlDahoul 2024)"),
            ("", ""),
            ("Diagnóstico novo:",
             "não é só bias algorítmico. É heterogeneidade fenotípica intra-categoria que a rotulagem monolítica de FairFace apaga."),
            ("", ""),
            ("Fundamentação empírica agora tripla:"),
            ("Antropologia:",
             "Telles 2014 (PERLA, 4 países latino-americanos) documenta pigmentocracia — tom de pele varia mais dentro da categoria do que entre elas"),
            ("Genética:",
             "Bryc 2015 (AJHG, 162 mil indivíduos) — ancestralidade Native + European + African altamente variável em Latinos"),
            ("Sociologia:",
             "Pew 2017 — identidade Hispanic cai de 97 % para 50 % em quatro gerações nos EUA. Categoria é sociopolítica, não fenotípica."),
        ],
        footer="Esse é o argumento pelo qual C6 (decomposição fenotípico × algorítmico) vira contribuição central.",
    )

    add_table_slide(
        prs,
        "Validação externa via NotebookLM — 7 perguntas de defesa antecipadas",
        ["Pergunta submetida", "Veredito", "Ponto-chave"],
        [
            ["Tese é nova?", "Sim", "FiLM em fairness facial é inédito; matriz MST × FairFace preenche gap de Pereira 2026"],
            ["Valor científico?", "Elevado", "Diagnóstico estrutural fenotípico × algorítmico + European AI Act 2024"],
            ["Divergências / SOTA?", "Sim", "FaceScanPaliGemma 75,7 % acc; Pangelinan 2023 refuta via 'pixel info' (endereçado por H6)"],
            ["Faltam artigos?", "Latinx", "Diversidade intra-Latinx monolítica — respondido pela Rodada 8 (Telles/Bryc/Pew)"],
            ["Furos e gaps?", "2 pontos", "Dependência SkinToneNet (mitigado por sensitivity analysis); escalabilidade além FairFace"],
            ["Corpus atualizado?", "Fronteira", "Papers 2026 já cobertos (Pereira, AlDahoul, Lian)"],
            ["Pipeline defensável?", "Sim", "Sequenciamento + controle de confounders + ablação + rigor estatístico"],
        ],
        col_widths=[3.5, 1.8, 7.2],
        footer='Veredito NotebookLM: "trabalho bem estruturado, contribuições claras, fundamentado em literatura de ponta. Ponto crítico da defesa: decomposição do erro Latinx".',
    )

    add_table_slide(
        prs,
        "Cinco caminhos possíveis para mitigar viés — por que escolhi FiLM + MST + Decomposição",
        ["Caminho", "Ideia central", "Posição relativa à minha proposta"],
        [
            ["1. Condicionamento arquitetural (FiLM)", "Modular features do backbone com tom de pele como sinal auxiliar", "Caminho que adotei — inédito em fairness facial"],
            ["2. Decomposição fenotípico × algorítmico", "Separar erro irredutível de erro mitigável", "Adotei — OE-6 na v3.6 formaliza isso"],
            ["3. Destilação de conhecimento (MST-KD)", "Múltiplos professores enviesados por etnia → estudante", "Alternativa avaliada — abandonada por custo computacional"],
            ["4. Intervenção geométrica em embeddings (SPD/SFID)", "Projeção ortogonal ao subespaço do viés em VLMs — training-free", "Complementar — cobrimos parcialmente via configuração D (FiLM-CLIP)"],
            ["5. Rótulos demográficos contínuos (Ethnicity Shift)", "Etnia como espectro contínuo, não categoria discreta", "Interseciona com MST — MST é já uma discretização de espectro"],
        ],
        col_widths=[3.5, 4.5, 5.0],
        highlight_rows=[0, 1],
        footer="Adotei os caminhos 1 + 2 combinados. Deixo 3-5 declarados em Cap 4 como alternativas ponderadas e não adotadas.",
    )

    # Complemento — solicitação pessoal
    add_bullets(
        prs,
        "Complemento — solicitação pessoal de extensão de 2 meses",
        [
            ("O que aconteceu:",
             "minha filha nasceu em 28/março. Fiquei afastado nas semanas seguintes para o cuidado que ela e a mãe precisaram"),
            ("", ""),
            ("O que preservei:",
             "marcos de 15/jul (1ª revisão) e 30/jul (pedido formal) continuam de pé"),
            ("O que estou pedindo:",
             "extensão de 2 meses no prazo da defesa — de agosto para outubro"),
            ("", ""),
            ("Documento:",
             "carta em 2 parágrafos redigida, pronta para envio via SEI ou trâmite indicado pelo Programa"),
        ],
        footer="Complemento à agenda, mas precisa de encaminhamento hoje.",
    )

    # Complemento — perguntas
    add_bullets(
        prs,
        "Perguntas em aberto — técnicas + administrativas",
        [
            ("1. Estrutura entregue no Overleaf:",
             "há capítulo que o senhor gostaria de ver com tratamento diferente antes da minha próxima rodada de revisão?"),
            ("2. Configurações comparativas:",
             "as 4 configs (A/B/C/D) bastam ou incluo uma 5ª variante (skin tone contínuo ITA/L* em vez de MST discreto)?"),
            ("3. Cap 2 — ablation arquitetural:",
             "faz sentido meia-página comparando FiLM × CBN × cross-attention × AdaIN, ou deixo para a defesa oral?"),
            ("4. OE-4 (fair transfer):",
             "métrica principal TPR@FAR=1e-4 estratificada por raça, ou DR agregado?"),
            ("5. Extensão de 2 meses:",
             "trâmite correto — carta direta, requerimento SEI ao PPG, ou processo formal?"),
            ("Bônus:",
             "sugestões para banca preliminar (fairness externo + CV aplicada)?"),
        ],
        footer="Nesta semana entro em modo revisão conforme feedback vindo do Overleaf.",
    )

    add_section_divider(prs, "", "Obrigado",
                        "Aberto para o debate")

    return prs


def main() -> None:
    here = Path(__file__).parent
    out = here / "material_reuniao_orientador_2026-07-13.pptx"
    prs = build_presentation()
    prs.save(out)
    print(f"Apresentacao gerada: {out}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
